import pytest
import hashlib
import json
import os
import re
import time
from datetime import datetime
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock
from pyclaw.core.agent import Agent
from pyclaw.core.message import Message, MessageRole, MessageType
from pyclaw.core.session import Session
from pyclaw.tools.base import ToolResult
from pyclaw.tools.registry import ToolRegistry
from pyclaw.tools.skill_activation import ActivateSkillTool, resolve_markdown_skill
from pyclaw.tools.terminal import TerminalTool
from pyclaw.core.skill_context import SkillContextService

@pytest.mark.asyncio
async def test_agent_self_healing_loop():
    # 1. Mock Components
    model = AsyncMock()
    # Use MagicMock for tools registry properties, but AsyncMock for its async methods
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    
    sessions = AsyncMock()
    
    # 2. Setup LLM response sequence
    # First call: Returns a tool call
    first_resp = {
        "content": "<thought>I need to run a command</thought>",
        "__tool_calls__": True,
        "tool_calls": [{
            "id": "call1",
            "function": {"name": "terminal", "arguments": '{"command": "false"}'}
        }]
    }
    # Second call (after failure): Returns a successful response
    second_resp = {
        "content": "It failed but I understand why.",
        "__tool_calls__": False
    }
    model.chat.side_effect = [first_resp, second_resp]
    
    # 3. Setup Tool results (Simulation: tool fails)
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "call1",
            "name": "terminal",
            "content": "Error: Command exited with code 1",
            "success": False,
            "metadata": {}
        }
    ]
    tools.get_all_specs.return_value = []
    
    # 4. Mock session manager get_or_create
    session = MagicMock()
    session.session_id = "s1"
    session.channel = "t"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    
    # Track messages added to session
    def add_msg_side_effect(msg):
        session.messages.append(msg)
    session.add_message.side_effect = add_msg_side_effect
    
    # Return LLM formatted history
    def get_history_side_effect(limit=10):
        return [m.to_llm_format() for m in session.messages]
    session.get_history.side_effect = get_history_side_effect
    
    # Track messages added to session via save_message
    async def save_msg_side_effect(sess, msg):
        # Only add if not already there to avoid duplicates if Agent also calls add_message
        if msg not in sess.messages:
            sess.messages.append(msg)
    sessions.save_message.side_effect = save_msg_side_effect
    
    sessions.get_or_create.return_value = session
    
    # Mock model embed for memory
    model.embed.return_value = [0.1] * 1536 # standard dim
    
    agent = Agent(model, tools, sessions)
    
    # 5. Process a message
    user_msg = Message(
        id="m1", channel="t", channel_user_id="u1", session_id="s1",
        type=MessageType.TEXT, role=MessageRole.USER, content="Run false"
    )
    response = await agent.process_message(user_msg)
    
    # 6. Assertions
    # Model should have been called twice (Initial -> After Failure)
    assert model.chat.call_count == 2
    
    # Check if <error_context> was in the conversation for the second call
    history_call_args = model.chat.call_args_list[1][1]["messages"]
    tool_msg = next(m for m in history_call_args if m["role"] == "tool")
    assert "<error_context>" in tool_msg["content"]
    assert "NOTICE: The tool call failed" in tool_msg["content"]
    
    # Final response content
    assert "It failed but I understand why." in response.content

@pytest.mark.asyncio
async def test_agent_max_iterations():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    
    sessions = AsyncMock()
    
    # Always returns a read-only tool call to simulate infinite loop
    infinite_resp = {
        "content": "looping...",
        "__tool_calls__": True,
        "tool_calls": [{
            "id": "loop",
            "function": {"name": "web_read", "arguments": '{"url": "https://example.com"}'}
        }]
    }
    model.chat.return_value = infinite_resp
    
    tools.execute_tool_calls.return_value = [
        {"role": "tool", "tool_call_id": "loop", "name": "web_read", "content": "ok", "success": True, "metadata": {}}
    ]
    
    session = MagicMock()
    session.session_id = "s2"
    session.channel = "t"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    
    def add_msg_side_effect(msg):
        session.messages.append(msg)
    session.add_message.side_effect = add_msg_side_effect
    
    def get_history_side_effect(limit=10):
        return [m.to_llm_format() for m in session.messages]
    session.get_history.side_effect = get_history_side_effect
    
    # Track messages added to session via save_message
    async def save_msg_side_effect(sess, msg):
        # Only add if not already there to avoid duplicates if Agent also calls add_message
        if msg not in sess.messages:
            sess.messages.append(msg)
    sessions.save_message.side_effect = save_msg_side_effect
    
    sessions.get_or_create.return_value = session
    
    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m2", channel="t", channel_user_id="u1", session_id="s2",
        type=MessageType.TEXT, role=MessageRole.USER, content="Loop me"
    )
    
    response = await agent.process_message(user_msg)
    
    # Should stop after max_iterations = 5
    assert model.chat.call_count == 5
    assert "达到最大思考深度" in response.content or "⚠️  思考超时" in response.content
    assert "OBSERVATION from" not in response.content


@pytest.mark.asyncio
async def test_agent_repeated_read_only_tools_force_final_answer_without_raw_observations():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    repeated_tool_response = {
        "content": "继续读取网页...",
        "__tool_calls__": True,
        "tool_calls": [{
            "id": "read-loop",
            "function": {"name": "web_read", "arguments": '{"url": "https://example.com"}'}
        }]
    }
    model.chat.side_effect = [repeated_tool_response] * 9 + [
        {"content": "基于已有信息，这是最终答复。", "__tool_calls__": False}
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "read-loop",
            "name": "web_read",
            "content": "A very long raw web page observation",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-repeated-tool"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-repeated-tool", channel="telegram", channel_user_id="u1", session_id="s-repeated-tool",
        type=MessageType.TEXT, role=MessageRole.USER, content="查一下"
    )

    response = await agent.process_message(user_msg)

    assert response.content == "基于已有信息，这是最终答复。"
    assert tools.execute_tool_calls.call_count == 4
    assert "OBSERVATION from" not in response.content
    assert "A very long raw web page observation" not in response.content


@pytest.mark.asyncio
async def test_activate_skill_resolves_nested_frontmatter_name(tmp_path, monkeypatch):
    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    skill_dir.mkdir(parents=True)
    (skill_dir / "SKILL.md").write_text(
        "---\n"
        "name: baoyu-design\n"
        "description: Make polished slide decks.\n"
        "---\n\n"
        "# Design\n\nFollow the deck workflow.\n",
        encoding="utf-8",
    )
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    resolved = resolve_markdown_skill("baoyu-design", [str(skills_root)])
    assert resolved is not None
    assert resolved.rel_path == "baoyu-design/skills/baoyu-design"

    result = await ActivateSkillTool().execute(name="baoyu-design")

    assert result.success is True
    assert result.metadata["activated_skill"] == "baoyu-design"
    assert result.metadata["activated_skill_path"] == "baoyu-design/skills/baoyu-design"
    assert result.metadata["activated_skill_md_path"] == str(skill_dir / "SKILL.md")
    assert result.metadata["activated_skill_content_sha256"]
    assert "path=\"baoyu-design/skills/baoyu-design\"" in result.content


@pytest.mark.asyncio
async def test_active_markdown_skill_context_is_persisted_and_rendered(tmp_path, monkeypatch):
    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    skill_dir.mkdir(parents=True)
    skill_md = skill_dir / "SKILL.md"
    skill_md.write_text(
        "---\nname: baoyu-design\ndescription: Deck workflow.\n---\n\n"
        "# Baoyu Design\n\nRead system-prompt.md before creating slides.\n",
        encoding="utf-8",
    )
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    activation = await ActivateSkillTool().execute(name="baoyu-design")
    session = Session(session_id="s-skill-context", user_id="u1", channel="feishu")

    record = SkillContextService().persist_activation(session, activation.metadata)

    assert record is not None
    assert session.metadata["active_skills"] == ["baoyu-design"]
    stored = session.metadata["active_skill_contexts"][0]
    assert stored["canonical_rel_path"] == "baoyu-design/skills/baoyu-design"
    assert stored["skill_md_path"] == str(skill_md)

    rendered = SkillContextService().render_prompt_context(session)
    assert "<active_skills>" in rendered
    assert "skill_md_path" in rendered
    assert "Read system-prompt.md" in rendered


@pytest.mark.asyncio
async def test_dynamic_prompt_includes_active_skill_context(tmp_path):
    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    skill_dir.mkdir(parents=True)
    skill_md = skill_dir / "SKILL.md"
    skill_md.write_text("---\nname: baoyu-design\n---\n\n# Workflow\nUse deck-stage HTML workflow.\n", encoding="utf-8")

    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [tmp_path / "skills"]
    tools._refresh_skills = MagicMock()
    agent = Agent(AsyncMock(), tools, AsyncMock())
    session = Session(
        session_id="s-dynamic-skill",
        user_id="u1",
        channel="feishu",
        metadata={
            "active_skills": ["baoyu-design"],
            "active_skill_contexts": [{
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(skill_md),
                "root_dir": str(tmp_path / "skills"),
            }],
        },
    )

    prompt = await agent._get_dynamic_system_prompt(session)

    assert "<active_skills>" in prompt
    assert "Use deck-stage HTML workflow" in prompt
    assert str(skill_md) in prompt


@pytest.mark.asyncio
async def test_duplicate_activate_skill_reinjects_canonical_context(tmp_path):
    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    skill_dir.mkdir(parents=True)
    skill_md = skill_dir / "SKILL.md"
    skill_md.write_text("---\nname: baoyu-design\n---\n\n# Workflow\nLoad references/codex.md.\n", encoding="utf-8")
    session = Session(
        session_id="s-dup-skill-context",
        user_id="u1",
        channel="feishu",
        metadata={
            "active_skills": ["baoyu-design"],
            "active_skill_contexts": [{
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(skill_md),
                "root_dir": str(tmp_path / "skills"),
            }],
        },
    )
    sessions = AsyncMock()

    async def save_msg_side_effect(sess, msg):
        sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    agent = Agent(AsyncMock(), tools, sessions)

    await agent._request_active_skill_continue(session, ["baoyu-design"])

    assert session.messages
    notice = session.messages[-1].content
    assert "already active" in notice
    assert str(skill_md) in notice
    assert "Load references/codex.md" in notice


def test_skills_index_lists_frontmatter_name_and_nested_path(tmp_path):
    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    skill_dir.mkdir(parents=True)
    (skill_dir / "SKILL.md").write_text(
        "---\nname: baoyu-design\ndescription: Polished deck workflow.\n---\n\n# Skill\n",
        encoding="utf-8",
    )
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    tools._refresh_skills = MagicMock()
    agent = Agent(AsyncMock(), tools, AsyncMock())

    index = agent._get_skills_index()

    assert "- baoyu-design: [Markdown Skill] Polished deck workflow." in index
    assert "(path: baoyu-design/skills/baoyu-design)" in index


def test_explicit_skill_deliverable_disables_generic_synthesis():
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.artifact_synthesis import SynthesisQuality

    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    agent = Agent(AsyncMock(), tools, AsyncMock())
    session = Session(session_id="s-skill-no-generic", user_id="u1", channel="feishu")
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir="/tmp/pyclaw-artifacts",
    )

    quality = agent._synthesis_quality_for_contract(
        session,
        contract,
        "## RAG 企业知识库\n\n- 检索增强生成\n- 向量数据库\n- 权限治理\n- 评测闭环\n- 运维监控\n",
    )

    assert quality == SynthesisQuality.DISABLED




def test_explicit_skill_contract_records_required_skill(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.completion_contract import CompletionContract

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path))
    session = Session(session_id="s-skill-contract", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-skill-contract",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
    ))

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.required_skills == ("baoyu-design",)
    assert CompletionContract.from_metadata(contract.to_metadata()).required_skills == ("baoyu-design",)


@pytest.mark.asyncio
async def test_explicit_skill_deliverable_prompt_injects_workspace_adapter(tmp_path, monkeypatch):
    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    (skill_dir / "built-in-skills").mkdir(parents=True)
    (skill_dir / "references").mkdir()
    (skill_dir / "SKILL.md").write_text(
        "---\nname: baoyu-design\ndescription: Deck workflow.\n---\n\n"
        "Read [system-prompt.md](system-prompt.md).",
        encoding="utf-8",
    )
    (skill_dir / "system-prompt.md").write_text("Read [references/codex.md](references/codex.md).", encoding="utf-8")
    (skill_dir / "references" / "codex.md").write_text("Codex harness", encoding="utf-8")
    (skill_dir / "built-in-skills" / "make-a-deck.md").write_text(
        "Use `<deck-stage width=\"1920\" height=\"1080\">` and `section data-label`.",
        encoding="utf-8",
    )
    (skill_dir / "built-in-skills" / "export-as-pptx-editable.md").write_text("Call gen_pptx.", encoding="utf-8")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    agent = Agent(AsyncMock(), tools, AsyncMock())
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path / "artifacts" / "rag_deck")
    agent.artifacts.root_path.return_value = str(tmp_path / "artifacts")
    session = Session(session_id="s-skill-prompt", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-skill-prompt",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
    ))

    prompt = await agent._get_dynamic_system_prompt(session)

    assert "<skill_workspace_adapter>" in prompt
    assert "bounded_artifact_dir:" in prompt
    assert "First load the required_skill_docs with the read_file tool" in prompt
    assert "built-in-skills/make-a-deck.md" in prompt
    assert "built-in-skills/export-as-pptx-editable.md" in prompt
    assert "references/*.md" in prompt


def test_explicit_skill_deliverable_requires_skill_evidence_not_just_pptx(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    artifact = tmp_path / "RAG.pptx"
    _write_test_pptx(artifact, 10)
    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    session = Session(session_id="s-skill-evidence-missing", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-skill-evidence-missing",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
    ))
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(tmp_path),
        required_skills=("baoyu-design",),
    )
    session.metadata["current_completion_contract"] = contract.to_metadata()
    pending_files = [{"file_path": str(artifact), "description": "RAG 企业知识库幻灯片"}]

    content = agent._prepare_completion_contract_final_content(
        session=session,
        content="已生成并发送文件：RAG.pptx",
        pending_files=pending_files,
    )

    assert "已生成并发送文件" not in content
    assert "skill 工作流验收" in content
    assert "未观察到已激活" in content
    assert pending_files == []
    assert session.metadata.get("current_completion_contract") is not None


def test_explicit_skill_deliverable_holds_generic_pptx_even_when_active(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    (skill_dir / "built-in-skills").mkdir(parents=True)
    (skill_dir / "SKILL.md").write_text(
        "---\nname: baoyu-design\n---\n\nRead [system-prompt.md](system-prompt.md).",
        encoding="utf-8",
    )
    (skill_dir / "system-prompt.md").write_text("Design instructions", encoding="utf-8")
    (skill_dir / "built-in-skills" / "make-a-deck.md").write_text(
        "Use `<deck-stage width=\"1920\" height=\"1080\">` with `section data-label`.",
        encoding="utf-8",
    )
    (skill_dir / "built-in-skills" / "export-as-pptx-editable.md").write_text("Call gen_pptx.", encoding="utf-8")
    artifact = tmp_path / "generic-rag.pptx"
    _write_test_pptx(artifact, 10)

    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(tmp_path),
        required_skills=("baoyu-design",),
    )
    session = Session(
        session_id="s-skill-generic-held",
        user_id="u1",
        channel="feishu",
        metadata={
            "current_completion_contract": contract.to_metadata(),
            "active_skill_contexts": [{
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(skill_dir / "SKILL.md"),
                "root_dir": str(tmp_path / "skills"),
            }],
        },
    )
    session.messages.append(Message(
        id="m-skill-generic-held",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=contract.task_text,
    ))
    pending_files = [{"file_path": str(artifact), "description": "generic PPT"}]
    agent = Agent(AsyncMock(), MagicMock(), MagicMock())

    content = agent._prepare_completion_contract_final_content(
        session=session,
        content="已生成并发送文件：generic-rag.pptx",
        pending_files=pending_files,
    )

    assert "skill 工作流验收" in content
    assert "缺少关键说明文件读取证据" in content
    assert pending_files == []
    assert session.metadata.get("current_completion_contract") is not None


def test_skill_evidence_detects_baoyu_deck_workflow_from_observations_and_artifact(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.skill_evidence import SkillEvidenceService

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    (skill_dir / "built-in-skills").mkdir(parents=True)
    (skill_dir / "references").mkdir()
    (skill_dir / "SKILL.md").write_text(
        "---\nname: baoyu-design\n---\n\nRead [system-prompt.md](system-prompt.md).",
        encoding="utf-8",
    )
    (skill_dir / "system-prompt.md").write_text("Read [references/codex.md](references/codex.md).", encoding="utf-8")
    (skill_dir / "references" / "codex.md").write_text("Codex harness", encoding="utf-8")
    (skill_dir / "built-in-skills" / "make-a-deck.md").write_text(
        "Use `<deck-stage width=\"1920\" height=\"1080\">` and `section data-label`.",
        encoding="utf-8",
    )
    (skill_dir / "built-in-skills" / "export-as-pptx-editable.md").write_text("Call gen_pptx.", encoding="utf-8")
    artifact_dir = tmp_path / "artifacts"
    artifact_dir.mkdir()
    (artifact_dir / "deck.html").write_text(
        '<deck-stage width="1920" height="1080"><section data-label="Title">RAG</section></deck-stage>',
        encoding="utf-8",
    )
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(artifact_dir),
        source_message_id="m-skill-evidence-ok",
        created_at=time.time() - 1,
        task_fingerprint=Agent(AsyncMock(), MagicMock(), MagicMock())._task_fingerprint("走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"),
        required_skills=("baoyu-design",),
    )
    session = Session(
        session_id="s-skill-evidence-ok",
        user_id="u1",
        channel="feishu",
        metadata={"active_skill_contexts": [{
            "name": "baoyu-design",
            "canonical_rel_path": "baoyu-design/skills/baoyu-design",
            "skill_md_path": str(skill_dir / "SKILL.md"),
            "root_dir": str(tmp_path / "skills"),
        }]},
    )
    session.messages.append(Message(
        id="m-skill-evidence-ok",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=contract.task_text,
    ))
    for rel in ("SKILL.md", "system-prompt.md", "references/codex.md", "built-in-skills/make-a-deck.md", "built-in-skills/export-as-pptx-editable.md"):
        session.messages.append(Message(
            id=f"tool-{rel}",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content=f"OBSERVATION from read_file:\nFile: {skill_dir / rel}\n...",
            metadata={"tool_name": "read_file"},
        ))

    result = SkillEvidenceService().evaluate(session=session, contract=contract, pending_files=[])

    assert result is not None
    assert result.satisfied is True
    assert "built-in-skills/make-a-deck.md" in result.observed_paths
    assert "<deck-stage" in result.observed_markers


def test_skill_evidence_ignores_stale_deck_html_for_new_contract(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.skill_evidence import SkillEvidenceService

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    (skill_dir / "built-in-skills").mkdir(parents=True)
    (skill_dir / "references").mkdir()
    (skill_dir / "SKILL.md").write_text(
        "---\nname: baoyu-design\n---\n\nRead [system-prompt.md](system-prompt.md).",
        encoding="utf-8",
    )
    (skill_dir / "system-prompt.md").write_text("Read [references/codex.md](references/codex.md).", encoding="utf-8")
    (skill_dir / "references" / "codex.md").write_text("Codex harness", encoding="utf-8")
    (skill_dir / "built-in-skills" / "make-a-deck.md").write_text(
        "Use `<deck-stage width=\"1920\" height=\"1080\">` and `section data-label`.",
        encoding="utf-8",
    )
    (skill_dir / "built-in-skills" / "export-as-pptx-editable.md").write_text("Call gen_pptx.", encoding="utf-8")
    artifact_dir = tmp_path / "artifacts"
    artifact_dir.mkdir()
    stale_deck = artifact_dir / "deck.html"
    stale_deck.write_text(
        '<deck-stage width="1920" height="1080"><section data-label="Old">RAG</section></deck-stage>',
        encoding="utf-8",
    )
    old_time = time.time() - 3600
    os.utime(stale_deck, (old_time, old_time))
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(artifact_dir),
        required_skills=("baoyu-design",),
        created_at=time.time(),
    )
    session = Session(
        session_id="s-stale-skill-evidence",
        user_id="u1",
        channel="feishu",
        metadata={"active_skill_contexts": [{
            "name": "baoyu-design",
            "canonical_rel_path": "baoyu-design/skills/baoyu-design",
            "skill_md_path": str(skill_dir / "SKILL.md"),
            "root_dir": str(tmp_path / "skills"),
        }]},
    )
    session.messages.append(Message(
        id="m-stale-skill-evidence",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=contract.task_text,
    ))
    for rel in ("SKILL.md", "system-prompt.md", "references/codex.md", "built-in-skills/make-a-deck.md", "built-in-skills/export-as-pptx-editable.md"):
        session.messages.append(Message(
            id=f"tool-stale-{rel}",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content=f"OBSERVATION from read_file:\nFile: {skill_dir / rel}\n...",
            metadata={"tool_name": "read_file"},
        ))

    result = SkillEvidenceService().evaluate(session=session, contract=contract, pending_files=[])

    assert result is not None
    assert result.satisfied is False
    assert "<deck-stage" not in result.observed_markers
    assert any("缺少工作流产物证据" in reason for reason in result.reasons)


def test_explicit_skill_deliverable_rejects_stale_pptx_from_previous_turn(tmp_path):
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.completion_contract import CompletionContract

    artifact = tmp_path / "RAG.pptx"
    _write_test_pptx(artifact, 10)
    old_time = time.time() - 3600
    os.utime(artifact, (old_time, old_time))
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(tmp_path),
        required_skills=("baoyu-design",),
        created_at=time.time(),
    )

    result = ArtifactAcceptanceService().evaluate(contract, [{"file_path": str(artifact)}])

    assert result.accepted is False
    assert any("早于当前任务" in reason for reason in result.reasons)


def test_explicit_skill_deck_adapter_generates_evidence_and_deliverable(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    (skill_dir / "built-in-skills").mkdir(parents=True)
    (skill_dir / "references").mkdir()
    (skill_dir / "SKILL.md").write_text(
        "---\nname: baoyu-design\n---\n\nRead [system-prompt.md](system-prompt.md).",
        encoding="utf-8",
    )
    (skill_dir / "system-prompt.md").write_text("Read [references/codex.md](references/codex.md).", encoding="utf-8")
    (skill_dir / "references" / "codex.md").write_text("Codex harness", encoding="utf-8")
    (skill_dir / "built-in-skills" / "make-a-deck.md").write_text(
        "Use `<deck-stage width=\"1920\" height=\"1080\">` and `section data-label`.",
        encoding="utf-8",
    )
    (skill_dir / "built-in-skills" / "export-as-pptx-editable.md").write_text("Call gen_pptx.", encoding="utf-8")
    artifact_dir = tmp_path / "artifacts"
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(artifact_dir),
        source_message_id="m-skill-adapter-ok",
        task_fingerprint=Agent(AsyncMock(), MagicMock(), MagicMock())._task_fingerprint("走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"),
        required_skills=("baoyu-design",),
    )
    session = Session(
        session_id="s-skill-adapter-ok",
        user_id="u1",
        channel="feishu",
        metadata={
            "current_completion_contract": contract.to_metadata(),
            "active_skill_contexts": [{
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(skill_dir / "SKILL.md"),
                "root_dir": str(tmp_path / "skills"),
            }],
        },
    )
    session.messages.append(Message(
        id="m-skill-adapter-ok",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=contract.task_text,
    ))
    for rel in ("SKILL.md", "system-prompt.md", "references/codex.md", "built-in-skills/make-a-deck.md", "built-in-skills/export-as-pptx-editable.md"):
        session.messages.append(Message(
            id=f"tool-adapter-{rel}",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content=f"OBSERVATION from read_file:\nFile: {skill_dir / rel}\n...",
            metadata={"tool_name": "read_file"},
        ))
    artifact_dir.mkdir(parents=True)
    (artifact_dir / "deck.html").write_text(
        """
<deck-stage width="1920" height="1080">
  <section data-label="RAG 企业知识库"><h1>RAG 企业知识库</h1><ul><li>检索增强生成</li><li>企业知识治理</li></ul></section>
  <section data-label="架构"><h1>架构</h1><ul><li>采集</li><li>索引</li><li>检索</li><li>生成</li></ul></section>
  <section data-label="落地路径"><h1>落地路径</h1><ul><li>权限治理</li><li>效果评测</li><li>持续迭代</li></ul></section>
</deck-stage>
""".strip(),
        encoding="utf-8",
    )
    pending_files: list[dict[str, str]] = []
    agent = Agent(AsyncMock(), MagicMock(), MagicMock())

    content = agent._prepare_completion_contract_final_content(
        session=session,
        content="没能顺利落地，稍后重试。",
        pending_files=pending_files,
    )

    assert "任务未完成" not in content
    assert "已生成并发送文件" in content
    assert len(pending_files) == 1
    pptx_path = pending_files[0]["file_path"]
    assert pptx_path.endswith(".pptx")
    deck_html = artifact_dir / "deck.html"
    assert deck_html.exists()
    html = deck_html.read_text(encoding="utf-8")
    assert "<deck-stage" in html
    assert "section data-label" in html
    text = _pptx_text(pptx_path)
    assert "RAG" in text
    assert "企业知识库" in text
    assert session.metadata.get("current_completion_contract") is None


def test_explicit_skill_deck_adapter_creates_missing_fresh_deck_html(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    artifact_dir = tmp_path / "artifacts"
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(artifact_dir),
        source_message_id="m-no-deck",
        task_fingerprint=Agent(AsyncMock(), MagicMock(), MagicMock())._task_fingerprint("走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"),
        required_skills=("baoyu-design",),
        created_at=time.time(),
    )
    session = Session(
        session_id="s-skill-adapter-no-deck",
        user_id="u1",
        channel="feishu",
        metadata={
            "current_completion_contract": contract.to_metadata(),
            "active_skill_contexts": [{
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(skill_dir / "SKILL.md"),
                "root_dir": str(skill_dir),
            }],
        },
    )
    session.messages.append(Message(
        id="m-no-deck",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=contract.task_text,
    ))
    for rel in ("SKILL.md", "system-prompt.md", "references/codex.md", "built-in-skills/make-a-deck.md", "built-in-skills/export-as-pptx-editable.md"):
        session.messages.append(Message(
            id=f"tool-no-deck-{rel}",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content=f"OBSERVATION from read_file:\nFile: {skill_dir / rel}\n...",
            metadata={"tool_name": "read_file"},
        ))

    pending_files: list[dict[str, str]] = []
    content = Agent(AsyncMock(), MagicMock(), MagicMock())._prepare_completion_contract_final_content(
        session=session,
        content="已生成的文件未通过 skill 工作流验收，任务未完成。",
        pending_files=pending_files,
    )

    assert len(pending_files) == 1
    assert pending_files[0]["file_path"].endswith(".pptx")
    assert "已生成并发送文件" in content
    assert "任务未完成" not in content
    assert "缺少工作流产物证据" not in content
    assert (artifact_dir / "deck.html").exists()
    assert (artifact_dir / "skill-workflow-evidence.md").exists()
    html = (artifact_dir / "deck.html").read_text(encoding="utf-8")
    assert "<deck-stage" in html
    assert "section data-label" in html
    text = _pptx_text(pending_files[0]["file_path"])
    assert "RAG" in text
    assert "企业知识库" in text


def test_explicit_skill_deck_adapter_overwrites_fresh_bad_deck_html(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    artifact_dir = tmp_path / "artifacts"
    artifact_dir.mkdir(parents=True)
    bad_deck = artifact_dir / "deck.html"
    bad_deck.write_text(
        """
<deck-stage width="1920" height="1080">
  <section data-label="当前进展"><h1>当前进展</h1><ul><li>工作目录已创建</li><li>稍后重试</li></ul></section>
  <section data-label="AI Agent"><h1>AI Agent</h1><ul><li>补充要点</li><li>后续可替换</li></ul></section>
</deck-stage>
""".strip(),
        encoding="utf-8",
    )
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-bad-deck",
        task_fingerprint=Agent(AsyncMock(), MagicMock(), MagicMock())._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time() - 1,
    )
    session = Session(
        session_id="s-bad-deck",
        user_id="u1",
        channel="feishu",
        metadata={
            "current_completion_contract": contract.to_metadata(),
            "active_skill_contexts": [{
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(skill_dir / "SKILL.md"),
                "root_dir": str(skill_dir),
            }],
        },
    )
    session.messages.append(Message(
        id="m-bad-deck",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))
    for rel in ("SKILL.md", "system-prompt.md", "references/codex.md", "built-in-skills/make-a-deck.md", "built-in-skills/export-as-pptx-editable.md"):
        session.messages.append(Message(
            id=f"tool-bad-deck-{rel}",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content=f"OBSERVATION from read_file:\nFile: {skill_dir / rel}\n...",
            metadata={"tool_name": "read_file"},
        ))

    pending_files: list[dict[str, str]] = []
    content = Agent(AsyncMock(), MagicMock(), MagicMock())._prepare_completion_contract_final_content(
        session=session,
        content="已生成的文件未通过交付验收，任务未完成。",
        pending_files=pending_files,
    )

    assert len(pending_files) == 1
    assert pending_files[0]["file_path"].endswith(".pptx")
    assert "已生成并发送文件" in content
    assert "任务未完成" not in content
    html = bad_deck.read_text(encoding="utf-8")
    assert "当前进展" not in html
    assert "补充要点" not in html
    assert html.count("<section") >= 12
    text = _pptx_text(pending_files[0]["file_path"])
    for keyword in ("RAG", "企业知识库", "检索", "向量", "权限", "评测"):
        assert keyword in text
    for marker in ("当前进展", "任务未完成", "稍后重试", "补充要点", "后续可替换"):
        assert marker not in text


def test_explicit_skill_deck_adapter_replaces_rejected_pending_pptx(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    artifact_dir = tmp_path / "artifacts"
    artifact_dir.mkdir(parents=True)
    rejected = artifact_dir / "wrong_agent_deck.pptx"
    _write_process_report_pptx(rejected, 12)
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-rejected-pending",
        task_fingerprint=Agent(AsyncMock(), MagicMock(), MagicMock())._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time() - 1,
    )
    session = Session(
        session_id="s-rejected-pending",
        user_id="u1",
        channel="feishu",
        metadata={
            "current_completion_contract": contract.to_metadata(),
            "active_skill_contexts": [{
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(skill_dir / "SKILL.md"),
                "root_dir": str(skill_dir),
            }],
        },
    )
    session.messages.append(Message(
        id="m-rejected-pending",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))
    for rel in ("SKILL.md", "system-prompt.md", "references/codex.md", "built-in-skills/make-a-deck.md", "built-in-skills/export-as-pptx-editable.md"):
        session.messages.append(Message(
            id=f"tool-rejected-pending-{rel}",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content=f"OBSERVATION from read_file:\nFile: {skill_dir / rel}\n...",
            metadata={"tool_name": "read_file"},
        ))

    pending_files = [{"file_path": str(rejected), "description": "rejected process report"}]
    content = Agent(AsyncMock(), MagicMock(), MagicMock())._prepare_completion_contract_final_content(
        session=session,
        content="已生成的文件未通过交付验收，任务未完成。",
        pending_files=pending_files,
    )

    assert len(pending_files) == 1
    assert pending_files[0]["file_path"] != str(rejected)
    assert pending_files[0]["file_path"].endswith(".pptx")
    assert "已生成并发送文件" in content
    assert "任务未完成" not in content
    text = _pptx_text(pending_files[0]["file_path"])
    assert "RAG" in text
    assert "企业知识库" in text
    assert "当前进展" not in text


def test_target_baoyu_rag_prompt_generates_polished_skill_pptx(tmp_path):
    import re
    import zipfile

    pytest.importorskip("pptx")
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService
    from pyclaw.core.completion_contract import CompletionContract

    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    artifact_dir = tmp_path / "artifacts"
    contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-polished-rag",
        task_fingerprint=Agent(AsyncMock(), MagicMock(), MagicMock())._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time(),
    )

    artifact = ArtifactSynthesisService().synthesize_deck_stage_workflow(
        contract,
        draft="已生成的文件未通过 skill 工作流验收，任务未完成。",
    )

    assert artifact is not None
    pptx_path = artifact.file_path
    assert os.path.getsize(pptx_path) > 30_000
    acceptance = ArtifactAcceptanceService().evaluate(contract, [{"file_path": pptx_path}])
    assert acceptance.accepted, acceptance.summary
    html = (artifact_dir / "deck.html").read_text(encoding="utf-8")
    assert "<deck-stage" in html
    assert "section data-label" in html
    text = _pptx_text(pptx_path)
    for keyword in ("RAG", "企业知识库", "检索", "向量", "权限", "评测"):
        assert keyword in text
    for marker in ("任务未完成", "工作流验收", "当前进展", "延伸分析"):
        assert marker not in text
    with zipfile.ZipFile(pptx_path) as zf:
        slide_names = [name for name in zf.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)]
        slide_xml = "".join(zf.read(name).decode("utf-8", errors="ignore") for name in slide_names)
    assert len(slide_names) >= 12
    assert slide_xml.count("solidFill") >= 20
    assert len(set(re.findall(r"srgbClr val=\"([0-9A-Fa-f]+)\"", slide_xml))) >= 6


def test_controller_deck_stage_adapter_requires_current_doc_evidence(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    artifact_dir = tmp_path / "artifacts"
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-no-docs",
        task_fingerprint=Agent(AsyncMock(), MagicMock(), MagicMock())._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time(),
    )
    session = Session(
        session_id="s-skill-adapter-no-docs",
        user_id="u1",
        channel="feishu",
        metadata={
            "current_completion_contract": contract.to_metadata(),
            "active_skill_contexts": [{
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(skill_dir / "SKILL.md"),
                "root_dir": str(skill_dir),
            }],
        },
    )
    session.messages.append(Message(
        id="m-no-docs",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    pending_files: list[dict[str, str]] = []
    content = Agent(AsyncMock(), MagicMock(), MagicMock())._prepare_completion_contract_final_content(
        session=session,
        content="没能顺利落地，需要稍后重试。",
        pending_files=pending_files,
    )

    assert pending_files == []
    assert "已生成并发送文件" not in content
    assert "缺少关键说明文件读取证据" in content


def test_explicit_skill_directory_probe_does_not_count_as_doc_evidence(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    artifact_dir = tmp_path / "artifacts"
    artifact_dir.mkdir()
    stale = artifact_dir / "RAG.pptx"
    _write_test_pptx(stale, 10)
    (artifact_dir / "deck.html").write_text(
        '<deck-stage width="1920" height="1080"><section data-label="RAG">RAG</section></deck-stage>',
        encoding="utf-8",
    )
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-ls-only",
        task_fingerprint=Agent(AsyncMock(), MagicMock(), MagicMock())._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time() - 1,
    )
    session = Session(
        session_id="s-ls-only",
        user_id="u1",
        channel="feishu",
        metadata={
            "current_completion_contract": contract.to_metadata(),
            "active_skill_contexts": [{
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(skill_dir / "SKILL.md"),
                "root_dir": str(skill_dir),
            }],
        },
    )
    session.messages.append(Message(
        id="m-ls-only",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))
    session.messages.append(Message(
        id="tool-ls-only",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.TOOL,
        content="OBSERVATION from terminal:\nsystem-prompt.md\nreferences/codex.md\nbuilt-in-skills/make-a-deck.md\nbuilt-in-skills/export-as-pptx-editable.md",
        metadata={"tool_name": "terminal"},
    ))

    pending_files = [{"file_path": str(stale), "description": "stale RAG deck"}]
    content = Agent(AsyncMock(), MagicMock(), MagicMock())._prepare_completion_contract_final_content(
        session=session,
        content="已生成并发送文件：RAG.pptx",
        pending_files=pending_files,
    )

    assert pending_files == []
    assert "已生成并发送文件" not in content
    assert "缺少关键说明文件读取证据" in content


@pytest.mark.asyncio
async def test_explicit_skill_deck_e2e_controller_completes_after_skill_hydration(tmp_path, monkeypatch):
    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    (skill_dir / "built-in-skills").mkdir(parents=True)
    (skill_dir / "references").mkdir()
    (skill_dir / "SKILL.md").write_text(
        "---\nname: baoyu-design\ndescription: Deck workflow.\n---\n\n"
        "Read [system-prompt.md](system-prompt.md).",
        encoding="utf-8",
    )
    (skill_dir / "system-prompt.md").write_text("Read [references/codex.md](references/codex.md).", encoding="utf-8")
    (skill_dir / "references" / "codex.md").write_text("Codex harness", encoding="utf-8")
    (skill_dir / "built-in-skills" / "make-a-deck.md").write_text(
        "Use `<deck-stage width=\"1920\" height=\"1080\">` and `section data-label`.",
        encoding="utf-8",
    )
    (skill_dir / "built-in-skills" / "export-as-pptx-editable.md").write_text("Call gen_pptx.", encoding="utf-8")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "我会使用 baoyu-design skill。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "activate-baoyu",
                "function": {
                    "name": "activate_skill",
                    "arguments": json.dumps({"name": "baoyu-design"}),
                },
            }],
        },
        {"content": "没能顺利落地，需要稍后重试。", "__tool_calls__": False},
    ]

    async def execute_tool_calls(tool_calls):
        assert len(tool_calls) == 1
        activation = await ActivateSkillTool().execute(name="baoyu-design")
        return [{
            "role": "tool",
            "tool_call_id": "activate-baoyu",
            "name": "activate_skill",
            "content": activation.content,
            "success": activation.success,
            "metadata": activation.metadata,
        }]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(side_effect=execute_tool_calls)
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    tools.get_all_specs.return_value = [{"name": "activate_skill"}, {"name": "send_file_to_user"}]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-explicit-skill-e2e"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=6)
    artifact_root = tmp_path / "artifacts"
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(artifact_root / "rag_deck")
    agent.artifacts.root = str(artifact_root)
    agent.artifacts.root_path.return_value = str(artifact_root)

    user_msg = Message(
        id="m-explicit-skill-e2e",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-explicit-skill-e2e",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的 12 页幻灯片",
    )

    response = await agent.process_message(user_msg)

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    assert pending_files[0]["file_path"].endswith(".pptx")
    assert "已生成并发送" in response.content
    assert "skill 工作流验收" not in response.content
    assert "任务未完成" not in response.content
    assert (artifact_root / "rag_deck" / "deck.html").exists()
    assert (artifact_root / "rag_deck" / "skill-workflow-evidence.md").exists()
    text = _pptx_text(pending_files[0]["file_path"])
    assert "RAG" in text
    assert "企业知识库" in text
    assert session.metadata.get("current_completion_contract") is None


def test_explicit_skill_repair_notice_contains_workspace_adapter_details(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.deliverable_workflow import DeliverableWorkflow
    from pyclaw.core.skill_evidence import SkillEvidenceResult

    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(tmp_path / "artifacts"),
        required_skills=("baoyu-design",),
    )
    notice = DeliverableWorkflow().repair_notice(
        contract,
        skill_evidence=SkillEvidenceResult(False, ("missing deck evidence",), (), ()),
    )

    assert contract.artifact_dir in notice
    assert "workflow artifacts" in notice
    assert "html2pptx" in notice
    assert "deck.html only for deck-stage skills" in notice
    assert "~/designs" in notice
    assert "send_file_to_user" in notice


def test_explicit_skill_file_deliverable_does_not_patch_first(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path))
    session = Session(session_id="s-no-patch-first-skill", user_id="u1", channel="feishu")
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    session.messages.append(Message(
        id="m-no-patch-first-skill",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    assert agent._infer_completion_contract(session) is not None
    assert not agent._should_run_patch_first_gate(
        session=session,
        task_text=task,
        changed_files=set(),
        already_repaired=False,
        is_final_iteration=False,
        force_final_answer=False,
        soft_deadline_reached=False,
    )

def test_activate_skill_repeat_bucket_is_semantic_and_strict():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.metadata = {}

    short = agent._tool_repeat_counter_name("activate_skill", json.dumps({"name": "baoyu-design"}))
    nested = agent._tool_repeat_counter_name(
        "activate_skill",
        json.dumps({"name": "baoyu-design/skills/baoyu-design"}),
    )

    assert short == "activate_skill:baoyu-design"
    assert nested == "activate_skill:baoyu-design"
    assert agent._tool_repeat_limit(short, 8, session) == 2


@pytest.mark.asyncio
async def test_agent_skips_already_active_skill_and_continues_deliverable():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "我再激活一下技能。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "skill1",
                "function": {"name": "activate_skill", "arguments": json.dumps({"name": "pptx"})},
            }],
        },
        {"content": "已继续生成，不再重复激活技能。", "__tool_calls__": False},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-active-skill-skip"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {"active_skills": ["pptx"]}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m-active-skill-skip",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-active-skill-skip",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="讲讲 pptx 技能怎么用",
    )

    response = await agent.process_message(user_msg)

    tools.execute_tool_calls.assert_not_called()
    assert any("already active" in str(m.content) for m in session.messages)
    assert "重复调用过多" not in response.content
    assert "已继续生成" in response.content


@pytest.mark.asyncio
async def test_file_deliverable_gate_forces_create_and_send_in_same_turn():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    artifact = "/Users/bytedance/.pyclaw/artifacts/ai_agent_slides/AI_Agent_Slides.pptx"
    model.chat.side_effect = [
        {"content": "# AI Agent 幻灯片方案\n\n你可以说“生成 pptx”，我下一轮直接出文件。", "__tool_calls__": False},
        {
            "content": "现在直接创建并发送文件。",
            "__tool_calls__": True,
            "tool_calls": [
                {
                    "id": "write1",
                    "function": {
                        "name": "write_file",
                        "arguments": json.dumps({"path": "~/.pyclaw/artifacts/ai_agent_slides/notes.txt", "content": "slides"}),
                    },
                },
                {
                    "id": "send1",
                    "function": {
                        "name": "send_file_to_user",
                        "arguments": json.dumps({"file_path": artifact, "description": "AI Agent 幻灯片"}),
                    },
                },
            ],
        },
        {"content": "AI Agent 幻灯片已生成并发送。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {"role": "tool", "tool_call_id": "write1", "name": "write_file", "content": "File written", "success": True, "metadata": {}},
        {
            "role": "tool",
            "tool_call_id": "send1",
            "name": "send_file_to_user",
            "content": "File sent",
            "success": True,
            "metadata": {"is_file_transfer": True, "file_path": artifact, "description": "AI Agent 幻灯片"},
        },
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-file-deliverable-gate"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=6)
    user_msg = Message(
        id="m-file-deliverable-gate",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-file-deliverable-gate",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    assert any("File deliverable gate failed" in str(m.content) for m in session.messages)
    assert tools.execute_tool_calls.call_count == 1
    assert response.metadata["pending_files"] == [{"file_path": artifact, "description": "AI Agent 幻灯片"}]
    assert "你可以说" not in response.content
    assert "已生成并发送" in response.content


def test_file_sandbox_repair_notice_points_to_artifacts_dir():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.metadata = {}
    notice = agent._tool_failure_repair_notice(
        tool_name="write_file",
        tool_content="Access denied: Path '~/gen_agent_ppt.py' is outside the allowed workspace(s)",
        session=session,
    )

    assert "~/.pyclaw/artifacts" in notice
    assert "~/Desktop" in notice
    assert "send_file_to_user" in notice


def test_terminal_stderr_redirect_does_not_create_file_mutation_key():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())

    assert agent._terminal_file_redirect_target("pip3 install python-pptx 2>&1 | tail -5") == ""
    assert agent._terminal_mutation_target_key("pip3 install python-pptx 2>&1 | tail -5") == ""
    assert "redirect:&1" not in str(
        agent._terminal_side_effect_call_key(
            json.dumps({"command": "pip3 install python-pptx 2>&1 | tail -5", "approved": True})
        )
    )
    assert agent._terminal_file_redirect_target("python build_ppt.py > out.log 2>&1") == "out.log"


def test_file_deliverable_auto_approves_artifact_scoped_terminal_command():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.session_id = "s-artifact-approval"
    session.channel = "feishu"
    session.metadata = {}
    session.messages = [
        Message(
            id="u-artifact-approval",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-artifact-approval",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="做一个关于 AI Agent 的幻灯片",
        )
    ]
    contract = agent._infer_completion_contract(session)
    assert contract is not None
    command = f'mkdir -p "{contract.artifact_dir}" && cd "{contract.artifact_dir}" && python3 -c "import pptx; print(\"ok\")"'
    tool_calls = [{
        "id": "artifact-prep",
        "function": {"name": "terminal", "arguments": json.dumps({"command": command})},
    }]

    updated = agent._auto_approve_explicit_terminal_calls(tool_calls, session=session)
    args = json.loads(updated[0]["function"]["arguments"])

    assert args["approved"] is True
    decisions = session.metadata["last_exec_approval_decisions"]
    assert decisions[-1]["reason"] == "artifact-scoped file delivery command"


def test_file_deliverable_does_not_auto_approve_outside_artifact_command():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.session_id = "s-artifact-deny"
    session.channel = "feishu"
    session.metadata = {}
    session.messages = [
        Message(
            id="u-artifact-deny",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-artifact-deny",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="做一个关于 AI Agent 的幻灯片",
        )
    ]
    tool_calls = [{
        "id": "bad-prep",
        "function": {"name": "terminal", "arguments": json.dumps({"command": "mkdir -p /tmp/ai-agent-ppt"})},
    }]

    updated = agent._auto_approve_explicit_terminal_calls(tool_calls, session=session)
    args = json.loads(updated[0]["function"]["arguments"])

    assert "approved" not in args


def test_continue_html_deliverable_recovers_prior_fenced_artifact(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    session = Session(session_id="s-html-recovery", user_id="u1", channel="feishu")
    original_task = "做一个可视化的教学网页，图解chatgpt结构和训练流程"
    html = """<!doctype html>
<html lang="zh-CN">
<head><meta charset="utf-8"><title>ChatGPT 结构和训练流程</title></head>
<body>
  <main>
    <h1>ChatGPT 结构和训练流程</h1>
    <section><h2>模型结构</h2><p>Transformer 解码器、Token、上下文窗口与注意力机制。</p></section>
    <section><h2>训练流程</h2><p>预训练、监督微调、偏好学习与安全对齐。</p></section>
    <svg role="img" aria-label="ChatGPT training flow"><text x="10" y="20">Pretrain -> SFT -> RLHF</text></svg>
  </main>
</body>
</html>"""
    session.messages.extend([
        Message(
            id="u-html-original",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content=original_task,
        ),
        Message(
            id="a-html-failed",
            channel="feishu",
            channel_user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.ASSISTANT,
            content=(
                "网页文件没有成功写入磁盘（`~/.pyclaw/artifacts/chatgpt-teach/index.html` 未创建）。\n"
                "```html\n" + html + "\n```\n"
                "下轮你说一句「继续生成文件」，我会重新落盘并发送。"
            ),
        ),
        Message(
            id="u-html-continue",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="继续生成文件",
        ),
    ])

    contract = agent._infer_completion_contract(session)
    assert contract is not None
    assert "教学网页" in contract.task_text
    pending_files: list[dict[str, str]] = []
    content = agent._prepare_completion_contract_final_content(
        session=session,
        content="这一轮仍然没能落盘，请手动保存。",
        pending_files=pending_files,
    )

    assert pending_files
    output_path = pending_files[0]["file_path"]
    assert output_path.endswith("index.html")
    assert Path(output_path).is_file()
    written = Path(output_path).read_text(encoding="utf-8")
    assert "<!doctype html>" in written.lower()
    assert "ChatGPT" in written
    assert "已生成并发送文件" in content
    assert "手动保存" not in content
    assert "任务未完成" not in content
    assert session.metadata.get("current_completion_contract") is None


def test_html_deliverable_delegation_text_synthesizes_verified_page(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    session = Session(session_id="s-html-gpt3", user_id="u1", channel="feishu")
    task = "做一个可视化的教学网页，图解GPT-3结构和训练流程"
    session.messages.append(Message(
        id="u-html-gpt3",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    contract = agent._infer_completion_contract(session)
    assert contract is not None
    assert contract.kind == "file_deliverable"
    pending_files: list[dict[str, str]] = []
    content = agent._prepare_completion_contract_final_content(
        session=session,
        content=(
            "交给外部代码生成器了，正在做 GPT-3 版本的可视化教学网页。"
            "这次我特意加了 In-Context Learning 和 GPT-2 vs GPT-3 对比，做好马上发你。"
        ),
        pending_files=pending_files,
    )

    assert pending_files
    output_path = Path(pending_files[0]["file_path"])
    assert output_path.name == "index.html"
    assert output_path.is_file()
    written = output_path.read_text(encoding="utf-8")
    assert "GPT-3" in written
    assert "175B" in written
    assert "In-Context Learning" in written
    assert "Few-shot" in written
    assert "Decoder-only Transformer" in written
    assert "交给外部代码生成器" not in written
    assert "做好马上发" not in content
    assert "已生成并发送文件" in content
    assert session.metadata.get("current_completion_contract") is None


def test_rich_html_deliverable_send_file_tool_is_not_required_skill_and_keeps_path(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    artifact_root = tmp_path / ".pyclaw" / "artifacts"
    agent.artifacts = ArtifactManager(root=str(artifact_root))
    explicit_file = artifact_root / "chatgpt-teach-rich" / "index.html"
    session = Session(session_id="s-rich-html-contract", user_id="u1", channel="feishu")
    task = (
        "接着上次做 ChatGPT 图解教学网页的丰富版。\n"
        f"单文件 HTML，写到 {explicit_file}，通过 send_file_to_user 交付给我。\n"
        "【板块】10 个板块：\n"
        "1. 总览时间线：GPT-1 → GPT-2 → GPT-3 → InstructGPT → ChatGPT → GPT-4\n"
        "2. Tokenizer 可视化：BPE 切分动画\n"
        "3. Transformer 架构：Embedding / Positional Encoding / Multi-Head Attention\n"
        "4. 注意力机制：Q/K/V 矩阵动画 + 注意力热力图\n"
        "5. Pretraining：loss 曲线 + 数据来源饼图\n"
        "6. SFT：人类示范数据卡片\n"
        "7. RLHF 三步曲：SFT → Reward Model → PPO\n"
        "8. 推理与采样：Greedy / Top-k / Top-p / Temperature\n"
        "9. Scaling Law：参数量 vs 性能曲线\n"
        "10. 术语速查表：20+ 关键概念\n"
        "技术栈：TailwindCSS (CDN) + Chart.js (CDN) + vanilla JS。"
    )
    session.messages.append(Message(
        id="u-rich-html-contract",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.required_skills == ()
    assert Path(contract.artifact_dir) == explicit_file.parent


def test_rich_html_partial_progress_synthesizes_verified_file_instead_of_deferral(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    artifact_root = tmp_path / ".pyclaw" / "artifacts"
    agent.artifacts = ArtifactManager(root=str(artifact_root))
    explicit_file = artifact_root / "chatgpt-teach-rich" / "index.html"
    session = Session(session_id="s-rich-html-final", user_id="u1", channel="feishu")
    task = (
        "接着上次做 ChatGPT 图解教学网页的丰富版。要求：\n\n"
        f"【产物】单文件 HTML，写到 {explicit_file}，通过 send_file_to_user 交付给我。\n\n"
        "【板块】10 个板块：\n"
        "1. 总览时间线：GPT-1 → GPT-2 → GPT-3 → InstructGPT → ChatGPT → GPT-4 的参数量/关键突破对比\n"
        "2. Tokenizer 可视化：一句话 BPE 切分成 token id 的动画\n"
        "3. Transformer 架构（增强）：Embedding / Positional Encoding / Multi-Head Attention / FFN / LayerNorm / Residual 分层拆解\n"
        "4. 注意力机制（增强）：Q/K/V 矩阵动画 + Multi-Head 并行示意 + 可交互注意力热力图\n"
        "5. Pretraining：下一个 token 预测的 loss 曲线 + 数据来源饼图\n"
        "6. SFT：人类示范数据卡片，prompt→ideal response 配对示例\n"
        "7. RLHF 三步曲：SFT → Reward Model → PPO 流程图，含为什么需要 RM 对比\n"
        "8. 推理与采样：Greedy / Top-k / Top-p / Temperature 四种策略的输出分布对比\n"
        "9. 彩蛋 Scaling Law：参数量 vs 性能对数曲线\n"
        "10. 术语速查表：折叠面板，20+ 关键概念\n\n"
        "【技术栈】TailwindCSS (CDN) + Chart.js (CDN) + 纯 CSS 动画 + vanilla JS 交互。"
    )
    session.messages.extend([
        Message(
            id="u-rich-html-final",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content=task,
        ),
        Message(
            id="tool-invalid-json",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content="OBSERVATION from write_file (FAILED): Invalid JSON arguments for tool 'write_file'.",
        ),
    ])
    contract = agent._infer_completion_contract(session)
    assert contract is not None
    assert contract.required_skills == ()
    pending_files: list[dict[str, str]] = []

    content = agent._prepare_completion_contract_final_content(
        session=session,
        content=(
            "# ChatGPT 图解教学网页（丰富版）- 进度报告\n"
            "已完成前 5 个板块。尚未执行 build.py，尚未 send_file_to_user。"
            "请在新一轮会话继续生成剩余板块。"
        ),
        pending_files=pending_files,
    )

    assert pending_files
    output_path = Path(pending_files[0]["file_path"])
    assert output_path == explicit_file
    assert output_path.is_file()
    html = output_path.read_text(encoding="utf-8")
    for marker in ("Tokenizer", "Transformer", "Q/K/V", "Pretraining", "SFT", "RLHF", "Top-k", "Top-p", "Scaling Law", "术语"):
        assert marker in html
    assert "进度报告" not in html
    assert "下一轮" not in html
    assert "已生成并发送文件" in content
    assert "新一轮" not in content
    assert "尚未" not in content
    assert session.metadata.get("current_completion_contract") is None


def test_internal_guardrail_text_does_not_pollute_file_deliverable_contract():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.session_id = "s-contract-context"
    session.channel = "feishu"
    session.metadata = {}
    guardrail = "未观察到目标文件已生成并通过 send_file_to_user 发送，任务未完成。"
    session.messages = [
        Message(
            id="u-real",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-contract-context",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="做一个关于 AI Agent 的幻灯片",
        ),
        Message(
            id="a-guardrail",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-contract-context",
            type=MessageType.TEXT,
            role=MessageRole.ASSISTANT,
            content=guardrail,
        ),
        Message(
            id="u-confirm",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-contract-context",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="生成 pptx",
        ),
    ]

    assert agent._pending_action_context_for_short_confirmation(session, "生成 pptx") == ""
    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert "未观察到" not in contract.artifact_dir
    assert "任务未完成" not in contract.artifact_dir


def test_internal_notice_does_not_become_current_task_boundary():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.session_id = "s-current-boundary"
    session.metadata = {}
    session.messages = [
        Message(
            id="u-real-boundary",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-current-boundary",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="做一个关于 AI Agent 的幻灯片",
        ),
        Message(
            id="internal-boundary",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-current-boundary",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="NOTICE: Completion contract failed. Create/export the file now under bad/path.",
            metadata={"internal_notice": True},
        ),
    ]

    bounded = agent._add_current_task_boundary(session, [])

    assert bounded
    assert "做一个关于 AI Agent 的幻灯片" in bounded[-1]["content"]
    assert "bad/path" not in bounded[-1]["content"]


def test_persisted_completion_contract_survives_internal_notice_repairs():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.session_id = "s-persisted-contract"
    session.channel = "feishu"
    session.metadata = {}
    session.messages = [
        Message(
            id="u-contract-real",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-persisted-contract",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="做一个关于 AI Agent 的幻灯片",
        )
    ]

    original = agent._infer_completion_contract(session)
    assert original is not None
    session.messages.append(
        Message(
            id="u-contract-notice",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-persisted-contract",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="NOTICE: Completion contract failed. Create/export the file now under polluted/internal/path.",
            metadata={"internal_notice": True},
        )
    )

    recovered = agent._infer_completion_contract(session)

    assert recovered == original
    assert recovered is not None
    assert "polluted" not in recovered.artifact_dir
    assert "internal" not in recovered.artifact_dir


@pytest.mark.asyncio
async def test_file_deliverable_contract_repairs_after_duplicate_side_effect_block():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    script_path = "/Users/bytedance/.pyclaw/artifacts/ai_agent_slides/build_ppt.py"
    artifact = "/Users/bytedance/.pyclaw/artifacts/ai_agent_slides/AI_Agent_Slides.pptx"
    write_args = json.dumps({"path": script_path, "content": "print('build ppt')"})
    model.chat.side_effect = [
        {
            "content": "先写生成脚本。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "write1", "function": {"name": "write_file", "arguments": write_args}}],
        },
        {
            "content": "再写一次脚本。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "write2", "function": {"name": "write_file", "arguments": write_args}}],
        },
        {
            "content": "重复写入已跳过，现在发送生成文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send1",
                "function": {
                    "name": "send_file_to_user",
                    "arguments": json.dumps({"file_path": artifact, "description": "AI Agent 幻灯片"}),
                },
            }],
        },
        {"content": "AI Agent 幻灯片已生成并发送。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "write1", "name": "write_file", "content": "File written", "success": True, "metadata": {}}],
        [{
            "role": "tool",
            "tool_call_id": "send1",
            "name": "send_file_to_user",
            "content": "File sent",
            "success": True,
            "metadata": {"is_file_transfer": True, "file_path": artifact, "description": "AI Agent 幻灯片"},
        }],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-deliverable-duplicate-side-effect"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=8)
    user_msg = Message(
        id="m-deliverable-duplicate-side-effect",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-deliverable-duplicate-side-effect",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.metadata["pending_files"] == [{"file_path": artifact, "description": "AI Agent 幻灯片"}]
    assert "副作用工具重复调用" not in response.content
    assert "未观察到目标文件" not in response.content
    assert any("Completion contract failed" in str(m.content) for m in session.messages)


@pytest.mark.asyncio
async def test_completion_contract_rejects_false_file_delivery_claim():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {"content": "PDF 报告已生成并发送。", "__tool_calls__": False},
        {"content": "我已经把 PDF 文件发给你了。", "__tool_calls__": False},
        {"content": "PDF 报告已生成并发送。", "__tool_calls__": False},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-contract-false-claim"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m-contract-false-claim",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-contract-false-claim",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="帮我生成一份 AI Agent PDF 报告文件",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 0
    assert response.metadata == {}
    assert "未观察到目标文件已生成" in response.content
    assert "已生成并发送" not in response.content
    assert sum("Completion contract failed" in str(m.content) for m in session.messages) == 2


def test_completion_contract_infers_capture_and_artifact_dir():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.session_id = "s-capture-contract"
    session.messages = [
        Message(
            id="u-capture",
            channel="telegram",
            channel_user_id="u1",
            session_id="s-capture-contract",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="截个屏",
        )
    ]

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.kind == "capture_artifact"
    assert os.path.expanduser("~/.pyclaw/artifacts") in contract.artifact_dir
    assert os.path.isabs(contract.artifact_dir)


def test_short_ppt_confirmation_recovers_pending_action_context():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.messages = [
        Message(
            id="assistant-outline",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-ppt-confirm",
            type=MessageType.TEXT,
            role=MessageRole.ASSISTANT,
            content="AI Agent 幻灯片大纲已准备好。你回复『生成 pptx』，我就创建 PPTX 文件并发送。",
        ),
        Message(
            id="user-confirm",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-ppt-confirm",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="生成 pptx",
        ),
    ]

    pending = agent._pending_action_context_for_short_confirmation(session, "生成 pptx")

    assert "幻灯片大纲" in pending
    assert "生成 pptx" in pending


def test_file_deliverable_task_does_not_trigger_patch_first_gate():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.channel = "feishu"

    assert not agent._should_run_patch_first_gate(
        session=session,
        task_text="做一个关于 AI Agent 的幻灯片",
        changed_files=set(),
        already_repaired=False,
        is_final_iteration=False,
        force_final_answer=False,
        soft_deadline_reached=False,
    )


@pytest.mark.asyncio
async def test_cron_agent_soft_deadline_stops_research_and_synthesizes(monkeypatch):
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.return_value = {"content": "已基于已有结果总结。", "__tool_calls__": False}

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-soft-deadline"
    session.channel = "cron"
    session.channel_user_id = "job_1"
    session.user_id = "job_1"
    session.messages = []
    session.metadata = {"soft_deadline_seconds": 1}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    real_monotonic = time.monotonic
    times = [real_monotonic(), real_monotonic() + 2]
    monkeypatch.setattr("pyclaw.core.agent.time.monotonic", lambda: times.pop(0) if times else real_monotonic() + 2)

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-cron-soft-deadline", channel="cron", channel_user_id="job_1", session_id="s-cron-soft-deadline",
        type=MessageType.TEXT, role=MessageRole.USER, content="执行一个快超时的定时任务"
    )

    response = await agent.process_message(user_msg)

    assert response.content == "已基于已有结果总结。"
    tools.get_all_specs.assert_called()
    tools.execute_tool_calls.assert_not_called()
    assert any("cron research budget is exhausted" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_cron_soft_deadline_allows_one_delivery_tool(monkeypatch):
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(return_value=[
        {
            "role": "tool",
            "tool_call_id": "mail1",
            "name": "163email__send_email",
            "content": "email sent",
            "success": True,
            "metadata": {},
        }
    ])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [
        {"name": "web_read", "description": "read", "parameters": {}},
        {"name": "163email__send_email", "description": "send", "parameters": {}},
    ]

    model.chat.side_effect = [
        {
            "content": "准备发送邮件",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "mail1",
                "function": {"name": "163email__send_email", "arguments": '{"to":"u@example.com"}'},
            }],
        },
        {"content": "邮件已发送，任务完成。", "__tool_calls__": False},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-delivery"
    session.channel = "cron"
    session.channel_user_id = "job_1"
    session.user_id = "job_1"
    session.messages = []
    session.metadata = {"soft_deadline_seconds": 1}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    real_monotonic = time.monotonic
    times = [real_monotonic(), real_monotonic() + 2]
    monkeypatch.setattr("pyclaw.core.agent.time.monotonic", lambda: times.pop(0) if times else real_monotonic() + 2)

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-cron-delivery", channel="cron", channel_user_id="job_1", session_id="s-cron-delivery",
        type=MessageType.TEXT, role=MessageRole.USER, content="发送一封定时邮件"
    )

    response = await agent.process_message(user_msg)

    assert response.content == "邮件已发送，任务完成。"
    tools.execute_tool_calls.assert_awaited_once()
    first_chat_tools = model.chat.await_args_list[0].kwargs["tools"]
    assert first_chat_tools == [{"name": "163email__send_email", "description": "send", "parameters": {}}]
    assert any("收尾交付动作已执行" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_cron_soft_deadline_blocks_more_research_tools(monkeypatch):
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [
        {"name": "web_read", "description": "read", "parameters": {}},
        {"name": "163email__send_email", "description": "send", "parameters": {}},
    ]

    model.chat.side_effect = [
        {
            "content": "还想继续读网页",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "read1",
                "function": {"name": "web_read", "arguments": '{"url":"https://example.com"}'},
            }],
        },
        {"content": "基于已有信息总结。", "__tool_calls__": False},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-block-research"
    session.channel = "cron"
    session.channel_user_id = "job_1"
    session.user_id = "job_1"
    session.messages = []
    session.metadata = {"soft_deadline_seconds": 1}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    real_monotonic = time.monotonic
    times = [real_monotonic(), real_monotonic() + 2]
    monkeypatch.setattr("pyclaw.core.agent.time.monotonic", lambda: times.pop(0) if times else real_monotonic() + 2)

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-cron-block-research", channel="cron", channel_user_id="job_1", session_id="s-cron-block-research",
        type=MessageType.TEXT, role=MessageRole.USER, content="执行一个快超时的定时任务"
    )

    response = await agent.process_message(user_msg)

    assert response.content == "基于已有信息总结。"
    tools.execute_tool_calls.assert_not_called()
    assert any("只允许一次邮件/消息等交付动作" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_sanitizes_internal_timeout_preamble_from_final_answer():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.return_value = {
        "content": (
            "工具调用已达到执行时限，不再继续搜索。基于已有数据整理如下：\n"
            "🏆 今日赛程\n"
            "- A vs B：待官方确认"
        ),
        "__tool_calls__": False,
    }

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-timeout-preamble"
    session.channel = "cron"
    session.channel_user_id = "job_1"
    session.user_id = "job_1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-timeout-preamble", channel="cron", channel_user_id="job_1", session_id="s-timeout-preamble",
        type=MessageType.TEXT, role=MessageRole.USER, content="整理赛事消息",
    )

    response = await agent.process_message(user_msg)

    assert "工具调用已达到执行时限" not in response.content
    assert "不再继续搜索" not in response.content
    assert response.content.startswith("基于已有数据整理如下")


@pytest.mark.asyncio
async def test_cron_agent_uses_session_iteration_budget_and_reports_activity():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    responses = [
        {"content": f"第{i}轮", "__tool_calls__": False}
        for i in range(12)
    ]
    model.chat.side_effect = responses

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-budget"
    session.channel = "cron"
    session.channel_user_id = "job_1"
    session.user_id = "job_1"
    session.messages = []
    session.metadata = {"max_iterations": 12}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-cron-budget", channel="cron", channel_user_id="job_1", session_id="s-cron-budget",
        type=MessageType.TEXT, role=MessageRole.USER, content="执行需要较多轮的任务",
    )

    response = await agent.process_message(user_msg)
    activity = agent.get_activity_summary()

    assert response.content == "第0轮"
    assert model.chat.call_count == 1
    assert activity["activity_seq"] > 0
    assert activity["session_id"] == "s-cron-budget"
    assert activity["last_event"] == "final_answer"


@pytest.mark.asyncio
async def test_agent_retries_transient_llm_timeout_before_success(monkeypatch):
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = [TimeoutError("Request timed out."), {"content": "重试后成功。", "__tool_calls__": False}]
    monkeypatch.setattr("pyclaw.core.agent.asyncio.sleep", AsyncMock())

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-llm-retry"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-llm-retry", channel="feishu", channel_user_id="u1", session_id="s-llm-retry",
        type=MessageType.TEXT, role=MessageRole.USER, content="查一下最新消息",
    )

    response = await agent.process_message(user_msg)

    assert response.content == "重试后成功。"
    assert model.chat.call_count == 2


@pytest.mark.asyncio
async def test_cron_llm_timeout_returns_non_provider_error_after_retries(monkeypatch):
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = TimeoutError("Request timed out.")
    monkeypatch.setattr("pyclaw.core.agent.asyncio.sleep", AsyncMock())

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-llm-timeout"
    session.channel = "cron"
    session.channel_user_id = "job_1"
    session.user_id = "job_1"
    session.messages = []
    session.metadata = {"llm_retry_attempts": 2}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-cron-llm-timeout", channel="cron", channel_user_id="job_1", session_id="s-cron-llm-timeout",
        type=MessageType.TEXT, role=MessageRole.USER, content="执行早报",
    )

    response = await agent.process_message(user_msg)

    assert "模型请求连续超时" in response.content
    assert "Request timed out" not in response.content
    assert model.chat.call_count == 2


@pytest.mark.asyncio
async def test_agent_forces_final_answer_after_successful_side_effect_tool():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "执行通知命令...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "notify",
                "function": {"name": "terminal", "arguments": '{"command": "notify"}'}
            }]
        },
        {
            "content": "通知已完成。",
            "__tool_calls__": False,
        },
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "notify",
            "name": "terminal",
            "content": "notification sent",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-side-effect"
    session.channel = "cron"
    session.channel_user_id = "job_1"
    session.user_id = "job_1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-side-effect", channel="cron", channel_user_id="job_1", session_id="s-side-effect",
        type=MessageType.TEXT, role=MessageRole.USER, content="发送一次通知"
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "通知已完成。"
    assert "notification sent" not in response.content
    assert "副作用工具重复调用" not in response.content
    assert model.chat.await_args_list[-1].kwargs["tools"] is None
    assert any("副作用工具已经成功执行" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_allows_distinct_terminal_commands_without_repeat_guard():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "打开页面...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "open",
                "function": {
                    "name": "terminal",
                    "arguments": json.dumps({
                        "command": "opencli browser open https://example.larkoffice.com/wiki/WIKI_TOKEN_EXAMPLE",
                        "timeout": 30,
                    }),
                },
            }],
        },
        {
            "content": "查看页面状态...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "state",
                "function": {
                    "name": "terminal",
                    "arguments": json.dumps({"command": "opencli browser state", "timeout": 30}),
                },
            }],
        },
        {"content": "页面状态已读取。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "terminal",
            "name": "terminal",
            "content": "ok",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-distinct"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-distinct", channel="telegram", channel_user_id="u1", session_id="s-terminal-distinct",
        type=MessageType.TEXT, role=MessageRole.USER, content="打开并读取页面",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.content == "页面状态已读取。"
    assert "副作用工具重复调用" not in response.content


@pytest.mark.asyncio
async def test_agent_allows_multiple_distinct_file_edits_in_one_coding_turn():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    first_edit_args = json.dumps({"path": "MainActivity.java", "old": "a", "new": "b"})
    second_edit_args = json.dumps({"path": "activity_main.xml", "old": "x", "new": "y"})
    model.chat.side_effect = [
        {
            "content": "先改 Java。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "edit1", "function": {"name": "edit_file", "arguments": first_edit_args}}],
        },
        {
            "content": "再改布局。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "edit2", "function": {"name": "edit_file", "arguments": second_edit_args}}],
        },
        {
            "content": "运行验证。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "test1",
                "function": {"name": "terminal", "arguments": json.dumps({"command": "python -m py_compile MainActivity.java"})},
            }],
        },
        {"content": "已完成。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "edit1", "name": "edit_file", "content": "File edited: MainActivity.java", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "edit2", "name": "edit_file", "content": "File edited: activity_main.xml", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "test1", "name": "terminal", "content": "Command: python -m py_compile MainActivity.java\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-multiple-file-edits"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-multiple-file-edits", channel="feishu", channel_user_id="u1", session_id="s-multiple-file-edits",
        type=MessageType.TEXT, role=MessageRole.USER, content="请实现功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 3
    assert not any("副作用工具此前已经成功执行" in m.content and "edit_file" in m.content for m in session.messages)
    assert "已完成" in response.content
    assert "任务清单" in response.content

@pytest.mark.asyncio
async def test_agent_filters_duplicate_terminal_calls_in_same_batch():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    terminal_args = json.dumps({"command": "osascript -e 'display notification \"ok\"'", "timeout": 10})
    model.chat.side_effect = [
        {
            "content": "执行一次桌面通知...",
            "__tool_calls__": True,
            "tool_calls": [
                {"id": "notify1", "function": {"name": "terminal", "arguments": terminal_args}},
                {"id": "notify2", "function": {"name": "terminal", "arguments": terminal_args}},
            ],
        },
        {"content": "通知已完成。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "notify1",
            "name": "terminal",
            "content": "notification sent",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-same-batch"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-same-batch", channel="telegram", channel_user_id="u1", session_id="s-terminal-same-batch",
        type=MessageType.TEXT, role=MessageRole.USER, content="发送一次通知",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    executed_payload = json.loads(tools.execute_tool_calls.await_args.args[0])
    assert [tc["id"] for tc in executed_payload["tool_calls"]] == ["notify1"]
    assert response.content == "通知已完成。"
    assert "副作用工具重复调用" not in response.content
    assert any("已跳过重复项" in m.content and "terminal:" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_filters_repeated_lock_call_but_runs_distinct_wake_call():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    lock_args = json.dumps({"command": "pmset displaysleepnow", "approved": True})
    wake_args = json.dumps({"command": "caffeinate -u -t 1", "approved": True})
    model.chat.side_effect = [
        {
            "content": "先锁屏...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock1", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {
            "content": "再唤醒屏幕...",
            "__tool_calls__": True,
            "tool_calls": [
                {"id": "lock2", "function": {"name": "terminal", "arguments": lock_args}},
                {"id": "wake1", "function": {"name": "terminal", "arguments": wake_args}},
            ],
        },
        {"content": "锁屏和唤醒命令已执行。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [
            {
                "role": "tool",
                "tool_call_id": "lock1",
                "name": "terminal",
                "content": "display sleep requested",
                "success": True,
                "metadata": {},
            }
        ],
        [
            {
                "role": "tool",
                "tool_call_id": "wake1",
                "name": "terminal",
                "content": "wake requested",
                "success": True,
                "metadata": {},
            }
        ],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-lock-wake"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-lock-wake", channel="wechat", channel_user_id="u1", session_id="s-terminal-lock-wake",
        type=MessageType.TEXT, role=MessageRole.USER, content="帮我锁屏然后唤醒 Mac",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    second_payload = json.loads(tools.execute_tool_calls.await_args_list[1].args[0])
    assert [tc["id"] for tc in second_payload["tool_calls"]] == ["wake1"]
    assert response.content == "锁屏和唤醒命令已执行。"
    assert "副作用工具重复调用" not in response.content
    assert any(
        "已跳过重复项" in m.content and "terminal:mac_desktop_control:display_sleep" in m.content
        for m in session.messages
    )


@pytest.mark.asyncio
async def test_agent_synthesizes_after_repeated_executed_terminal_side_effect():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    notify_args = json.dumps({"command": "osascript -e 'display notification \"ok\"'", "approved": True})
    model.chat.side_effect = [
        {
            "content": "执行通知...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "notify1", "function": {"name": "terminal", "arguments": notify_args}}],
        },
        {
            "content": "再执行一次通知...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "notify2", "function": {"name": "terminal", "arguments": notify_args}}],
        },
        {"content": "已发送通知。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "notify1",
            "name": "terminal",
            "content": "notification sent",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-lock-repeat"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-lock-repeat", channel="wechat", channel_user_id="u1", session_id="s-terminal-lock-repeat",
        type=MessageType.TEXT, role=MessageRole.USER, content="帮我发送一次通知",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "已发送通知。"
    assert "副作用工具重复调用" not in response.content
    assert any("本轮只有重复的副作用工具调用" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_allows_repeated_mac_desktop_control_terminal_calls():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    lock_args = json.dumps({"command": "pmset displaysleepnow", "approved": True})
    model.chat.side_effect = [
        {
            "content": "第一次锁屏...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock1", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {
            "content": "按要求再锁屏一次...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock2", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {"content": "已执行两次锁屏命令。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [
            {
                "role": "tool",
                "tool_call_id": "lock1",
                "name": "terminal",
                "content": "display sleep requested",
                "success": True,
                "metadata": {},
            }
        ],
        [
            {
                "role": "tool",
                "tool_call_id": "lock2",
                "name": "terminal",
                "content": "display sleep requested",
                "success": True,
                "metadata": {},
            }
        ],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-lock-repeat-allowed"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-lock-repeat-allowed", channel="wechat", channel_user_id="u1", session_id="s-terminal-lock-repeat-allowed",
        type=MessageType.TEXT, role=MessageRole.USER, content="帮我连续锁屏两次 Mac",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.content == "已执行两次锁屏命令。"
    assert "副作用工具重复调用" not in response.content
    assert not any("本轮只有重复的副作用工具调用" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_stops_repeated_mac_lock_for_simple_wechat_request():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    lock_args = json.dumps({"command": "pmset displaysleepnow", "approved": True})
    model.chat.side_effect = [
        {
            "content": "锁屏...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock1", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {
            "content": "再试一次锁屏...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock2", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {"content": "已锁屏。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "lock1",
            "name": "terminal",
            "content": "display sleep requested",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-lock-simple"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-lock-simple", channel="wechat", channel_user_id="u1", session_id="s-terminal-lock-simple",
        type=MessageType.TEXT, role=MessageRole.USER, content="锁屏",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "已锁屏。"
    assert "副作用工具重复调用" not in response.content
    assert any(
        "本轮只有重复的副作用工具调用" in m.content
        and "terminal:mac_desktop_control:display_sleep" in m.content
        for m in session.messages
    )


@pytest.mark.asyncio
async def test_agent_counts_failed_mac_lock_attempt_against_repeat_budget():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    lock_args = json.dumps({"command": "pmset displaysleepnow"})
    model.chat.side_effect = [
        {
            "content": "锁屏...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock1", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {
            "content": "再试一次锁屏...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock2", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {"content": "已锁屏。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "lock1",
            "name": "terminal",
            "content": "Command: pmset displaysleepnow\nExit code: 1",
            "success": False,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-lock-failed-attempt"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-lock-failed-attempt",
        channel="wechat",
        channel_user_id="u1",
        session_id="s-terminal-lock-failed-attempt",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="锁屏",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "已锁屏。"
    assert "副作用工具重复调用" not in response.content
    assert any(
        "本轮只有重复的副作用工具调用" in m.content
        and "terminal:mac_desktop_control:display_sleep" in m.content
        for m in session.messages
    )


@pytest.mark.asyncio
async def test_agent_counts_failed_executed_terminal_side_effect_against_repeat_budget():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    lock_script_args = json.dumps({"command": "bash ~/.pyclaw/skills/mac-lock-unlock/lock.sh"})
    model.chat.side_effect = [
        {
            "content": "执行锁屏脚本...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock1", "function": {"name": "terminal", "arguments": lock_script_args}}],
        },
        {
            "content": "再执行一次锁屏脚本...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock2", "function": {"name": "terminal", "arguments": lock_script_args}}],
        },
        {"content": "锁屏脚本已尝试执行，检测未通过。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "lock1",
            "name": "terminal",
            "content": "Command: bash ~/.pyclaw/skills/mac-lock-unlock/lock.sh\nExit code: 1\nSTDOUT:\n❌ 锁屏失败",
            "success": False,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-script-failed-attempt"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-script-failed-attempt",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-terminal-script-failed-attempt",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="锁屏",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "锁屏脚本已尝试执行，检测未通过。"
    assert "副作用工具重复调用" not in response.content
    assert any("本轮只有重复的副作用工具调用" in m.content and "terminal:" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_recovers_failed_executed_terminal_side_effect_from_tool_history():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    lock_script_args = json.dumps({"command": "bash ~/.pyclaw/skills/mac-lock-unlock/lock.sh"})
    model.chat.side_effect = [
        {
            "content": "再次执行锁屏脚本...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock2", "function": {"name": "terminal", "arguments": lock_script_args}}],
        },
        {"content": "锁屏脚本前次已执行失败，不再重复刷屏。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "lock2",
            "name": "terminal",
            "content": "Command: bash ~/.pyclaw/skills/mac-lock-unlock/lock.sh\nExit code: 1\nSTDOUT:\n❌ 锁屏失败",
            "success": False,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-script-failed-rehydrated"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = [
        Message(
            id="m-terminal-script-failed-rehydrated",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-terminal-script-failed-rehydrated",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="锁屏",
        ),
        Message(
            id="tool-lock1-s-terminal-script-failed-rehydrated",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-terminal-script-failed-rehydrated",
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content="<error_context>\nOBSERVATION from terminal (FAILED):\nCommand: bash ~/.pyclaw/skills/mac-lock-unlock/lock.sh\nExit code: 1\nSTDOUT:\n❌ 锁屏失败\n</error_context>",
            metadata={"tool_name": "terminal", "tool_call_id": "lock1"},
        ),
    ]
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)

    response_content, _ = await agent._agent_loop(session)

    assert tools.execute_tool_calls.call_count == 0
    assert response_content == "锁屏脚本前次已执行失败，不再重复刷屏。"
    assert "副作用工具重复调用" not in response_content
    assert any("本轮只有重复的副作用工具调用" in m.content and "terminal:" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_counts_mac_lock_attempt_even_without_tool_result():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    lock_args = json.dumps({"command": "pmset displaysleepnow"})
    model.chat.side_effect = [
        {
            "content": "锁屏...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock1", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {
            "content": "再次锁屏...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock2", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {"content": "已锁屏。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = []

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-lock-no-result"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-lock-no-result",
        channel="wechat",
        channel_user_id="u1",
        session_id="s-terminal-lock-no-result",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="锁屏",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "已锁屏。"
    assert "副作用工具重复调用" not in response.content
    assert any(
        "本轮只有重复的副作用工具调用" in m.content
        and "terminal:mac_desktop_control:display_sleep" in m.content
        for m in session.messages
    )


@pytest.mark.asyncio
async def test_agent_recovers_mac_lock_attempt_from_saved_assistant_message():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    lock_args = json.dumps({"command": "pmset displaysleepnow", "approved": True})
    model.chat.side_effect = [
        {
            "content": "再次尝试锁屏...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock2", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {"content": "已锁屏。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "lock2",
            "name": "terminal",
            "content": "display sleep requested",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-lock-rehydrated"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = [
        Message(
            id="m-terminal-lock-rehydrated",
            channel="wechat",
            channel_user_id="u1",
            session_id="s-terminal-lock-rehydrated",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="锁屏",
        ),
        Message(
            id="assistant-toolcall-existing-s-terminal-lock-rehydrated",
            channel="wechat",
            channel_user_id="u1",
            session_id="s-terminal-lock-rehydrated",
            type=MessageType.TEXT,
            role=MessageRole.ASSISTANT,
            content="锁屏...",
            metadata={
                "tool_calls": [
                    {"id": "lock1", "function": {"name": "terminal", "arguments": lock_args}}
                ]
            },
        ),
    ]
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)

    response_content, _ = await agent._agent_loop(session)

    assert tools.execute_tool_calls.call_count == 0
    assert response_content == "已锁屏。"
    assert "副作用工具重复调用" not in response_content
    assert any(
        "本轮只有重复的副作用工具调用" in m.content
        and "terminal:mac_desktop_control:display_sleep" in m.content
        for m in session.messages
    )


def test_mac_desktop_control_commands_use_semantic_side_effect_keys():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())

    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": "pmset displaysleepnow", "approved": True})
    ) == "terminal:mac_desktop_control:display_sleep"
    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": "caffeinate -u", "approved": True})
    ) == "terminal:mac_desktop_control:wake"
    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": "caffeinate -u -t 1", "approved": True})
    ) == "terminal:mac_desktop_control:wake"
    assert agent._terminal_side_effect_call_key(
        json.dumps({
            "command": '"/System/Library/CoreServices/Menu Extras/User.menu/Contents/Resources/CGSession" -suspend',
            "approved": True,
        })
    ) == "terminal:mac_desktop_control:lock_screen"
    assert agent._terminal_side_effect_call_key(
        json.dumps({
            "command": "osascript -e 'tell application \"System Events\" to keystroke \"q\" using {control down, command down}'",
            "approved": True,
        })
    ) == "terminal:mac_desktop_control:lock_shortcut"
    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": "~/.pyclaw/bin/unlock.sh", "approved": True})
    ) == "terminal:mac_desktop_control:unlock"
    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": "bash ~/.pyclaw/bin/unlock.sh", "approved": True})
    ) == "terminal:mac_desktop_control:unlock"
    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": "bash ~/.pyclaw/skills/mac-wake-unlock/unlock.sh", "approved": True})
    ) == "terminal:mac_desktop_control:unlock"
    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": "echo start && bash ~/.pyclaw/skills/mac-wake-unlock/unlock.sh; echo done", "approved": True})
    ) == "terminal:mac_desktop_control:unlock"
    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": "bash ~/.pyclaw/skills/mac-lock-unlock/lock.sh", "approved": True})
    ) == "terminal:mac_desktop_control:lock_screen"

    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": "osascript -e 'display notification \"ok\"'", "approved": True})
    ) is not None


def test_user_intent_terminal_commands_use_semantic_side_effect_key():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    first = (
        'mkdir -p ~/.pyclaw/screenshots && '
        'f=~/.pyclaw/screenshots/screen_$(date +%Y%m%d_%H%M%S).png && '
        'screencapture -x "$f" && ls -lh "$f" && echo "PATH=$f"'
    )
    variant = (
        'mkdir -p ~/.pyclaw/screenshots && '
        'FILE=~/.pyclaw/screenshots/截图_$(date +%Y%m%d_%H%M%S).png && '
        'screencapture -x "$FILE" && ls -la "$FILE" && echo "PATH=$FILE"'
    )
    photo = 'mkdir -p ~/.pyclaw/photos && imagesnap ~/.pyclaw/photos/photo.jpg'

    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": first})
    ) == "terminal:semantic:capture_screenshot"
    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": variant})
    ) == "terminal:semantic:capture_screenshot"
    assert agent._terminal_side_effect_call_key(
        json.dumps({"command": photo})
    ) == "terminal:semantic:capture_photo"


def test_agent_auto_approves_terminal_call_when_latest_user_explicitly_requests_intent():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.messages = [
        Message(
            id="u1",
            channel="wechat",
            channel_user_id="u1",
            session_id="s1",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="截屏",
        )
    ]
    command = (
        'mkdir -p ~/.pyclaw/screenshots && '
        'f=~/.pyclaw/screenshots/screen_$(date +%Y%m%d_%H%M%S).png && '
        'screencapture -x "$f" && ls -lh "$f" && echo "PATH=$f"'
    )
    tool_calls = [{
        "id": "shot1",
        "function": {"name": "terminal", "arguments": json.dumps({"command": command})},
    }]

    updated = agent._auto_approve_explicit_terminal_calls(tool_calls, session=session)
    args = json.loads(updated[0]["function"]["arguments"])

    assert args["approved"] is True


@pytest.mark.asyncio
async def test_agent_executes_screenshot_with_injected_approval_for_real_user_turn(tmp_path):
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.allowed_paths = []
    tools.get_all_specs.return_value = []

    command = (
        'mkdir -p ~/.pyclaw/screenshots && '
        'f=~/.pyclaw/screenshots/截图_$(date +%Y%m%d_%H%M%S).png && '
        'screencapture -x "$f" && ls -lh "$f" && echo "PATH=$f"'
    )
    model.chat.side_effect = [
        {
            "content": "我来截屏。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "shot1",
                "function": {"name": "terminal", "arguments": json.dumps({"command": command}, ensure_ascii=False)},
            }],
        },
        {"content": "截图已完成。", "__tool_calls__": False},
    ]
    artifact_dir = tmp_path / ".pyclaw" / "screenshots"
    artifact_dir.mkdir(parents=True)
    artifact = artifact_dir / "截图.png"
    artifact.write_bytes(b"png")
    tools.execute_tool_calls.return_value = [{
        "role": "tool",
        "tool_call_id": "shot1",
        "name": "terminal",
        "content": f"Command: screencapture\nExit code: 0\nPATH={artifact}",
        "success": True,
        "metadata": {},
    }]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-screenshot-approval"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=3)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = os.path.abspath(os.path.join(os.getcwd(), "tmp_test_artifacts", "persist_contract"))
    user_msg = Message(
        id="m-screenshot-approval",
        channel="telegram",
        channel_user_id="u1",
        session_id="s-screenshot-approval",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="截个屏",
    )

    response = await agent.process_message(user_msg)

    assert "截屏已完成" in response.content
    assert response.metadata["pending_files"] == [{
        "file_path": str(artifact),
        "description": "截屏已完成 📸",
    }]
    payload = json.loads(tools.execute_tool_calls.call_args.args[0])
    executed_args = json.loads(payload["tool_calls"][0]["function"]["arguments"])
    assert executed_args["command"] == command
    assert executed_args["approved"] is True


@pytest.mark.asyncio
async def test_agent_auto_delivers_successful_screenshot_artifact_and_stops(tmp_path):
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.allowed_paths = []
    tools.get_all_specs.return_value = []

    artifact_dir = tmp_path / ".pyclaw" / "screenshots"
    artifact_dir.mkdir(parents=True)
    artifact = artifact_dir / "shot_20260704_234246.png"
    artifact.write_bytes(b"png")
    command = (
        'mkdir -p ~/.pyclaw/screenshots && '
        'F=~/.pyclaw/screenshots/shot_$(date +%Y%m%d_%H%M%S).png && '
        'screencapture -x "$F" && echo "OK:$F" && ls -la "$F"'
    )
    repeated_command = (
        'mkdir -p ~/.pyclaw/screenshots && '
        'f=~/.pyclaw/screenshots/screen_$(date +%Y%m%d_%H%M%S).png && '
        'screencapture -x "$f" && ls -la "$f" && echo "PATH=$f"'
    )
    model.chat.side_effect = [
        {
            "content": "我来截屏。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "shot1",
                "function": {"name": "terminal", "arguments": json.dumps({"command": command, "approved": True})},
            }],
        },
        {
            "content": "再截一次。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "shot2",
                "function": {"name": "terminal", "arguments": json.dumps({"command": repeated_command, "approved": True})},
            }],
        },
    ]
    tools.execute_tool_calls.return_value = [{
        "role": "tool",
        "tool_call_id": "shot1",
        "name": "terminal",
        "content": f"Command: screencapture\nExit code: 0\nSTDOUT:\nOK:{artifact}\n-rw-r--r--  1 mac staff 3 Jul 4 23:42 {artifact}",
        "success": True,
        "metadata": {},
    }]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-screenshot-auto-deliver"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m-screenshot-auto-deliver",
        channel="wechat",
        channel_user_id="u1",
        session_id="s-screenshot-auto-deliver",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="截个屏",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert model.chat.call_count == 1
    assert "截屏已完成" in response.content
    assert response.metadata["pending_files"] == [{
        "file_path": str(artifact),
        "description": "截屏已完成 📸",
    }]


def test_agent_configures_default_capture_artifact_paths_on_terminal_tools():
    registry = ToolRegistry(work_dir="/tmp/pyclaw-work", allowed_paths=[])
    terminal = TerminalTool()
    registry.register(terminal)

    Agent(AsyncMock(), registry, AsyncMock(), work_dir="/tmp/pyclaw-work")

    assert "~/.pyclaw/screenshots" in registry.allowed_paths
    assert "~/.pyclaw/photos" in registry.allowed_paths
    assert "~/.pyclaw/recordings" in terminal.allowed_paths


def test_agent_does_not_auto_approve_terminal_call_for_mismatched_user_intent():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = MagicMock()
    session.messages = [
        Message(
            id="u1",
            channel="wechat",
            channel_user_id="u1",
            session_id="s1",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="看一下目录",
        )
    ]
    command = 'mkdir -p ~/.pyclaw/screenshots && f=~/.pyclaw/screenshots/screen.png && screencapture -x "$f"'
    tool_calls = [{
        "id": "shot1",
        "function": {"name": "terminal", "arguments": json.dumps({"command": command})},
    }]

    updated = agent._auto_approve_explicit_terminal_calls(tool_calls, session=session)
    args = json.loads(updated[0]["function"]["arguments"])

    assert "approved" not in args


@pytest.mark.asyncio
async def test_agent_stops_screenshot_path_variants_after_bounded_repairs():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.allowed_paths = []
    tools.get_all_specs.return_value = []

    desktop_args = json.dumps({
        "command": 'FILE=~/Desktop/截图_$(date +%Y%m%d_%H%M%S).png && screencapture -x "$FILE" && echo "$FILE"',
    }, ensure_ascii=False)
    pyclaw_args = json.dumps({
        "command": 'mkdir -p ~/.pyclaw/screenshots && f=~/.pyclaw/screenshots/截图_$(date +%Y%m%d_%H%M%S).png && screencapture -x "$f" && echo "$f"',
        "approved": True,
    }, ensure_ascii=False)
    pyclaw_variant_args = json.dumps({
        "command": 'mkdir -p ~/.pyclaw/screenshots && f=~/.pyclaw/screenshots/screen_$(date +%Y%m%d_%H%M%S).png && screencapture -x "$f" && ls -lh "$f"',
        "approved": True,
    }, ensure_ascii=False)

    model.chat.side_effect = [
        {
            "content": "先截到桌面。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "shot1", "function": {"name": "terminal", "arguments": desktop_args}}],
        },
        {
            "content": "改到 PyClaw 目录再试。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "shot2", "function": {"name": "terminal", "arguments": pyclaw_args}}],
        },
        {
            "content": "再换一个文件名。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "shot3", "function": {"name": "terminal", "arguments": pyclaw_variant_args}}],
        },
        {"content": "截图失败：系统截图权限或沙箱路径仍阻止执行。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{
            "role": "tool",
            "tool_call_id": "shot1",
            "name": "terminal",
            "content": "⚠️ 拦截到非法路径访问: `~/Desktop/截图_20260704.png`。\n指令: `FILE=~/Desktop/截图_$(date +%Y%m%d_%H%M%S).png && screencapture -x \"$FILE\" && echo \"$FILE\"`",
            "success": False,
            "metadata": {},
        }],
        [{
            "role": "tool",
            "tool_call_id": "shot2",
            "name": "terminal",
            "content": "Command: mkdir -p ~/.pyclaw/screenshots && f=~/.pyclaw/screenshots/截图_$(date +%Y%m%d_%H%M%S).png && screencapture -x \"$f\" && echo \"$f\"\nExit code: 1\nSTDERR:\nscreencapture: cannot write file",
            "success": False,
            "metadata": {},
        }],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-screenshot-bounded-repair"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-screenshot-bounded-repair",
        channel="telegram",
        channel_user_id="u1",
        session_id="s-screenshot-bounded-repair",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="截屏",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.content == "截图失败：系统截图权限或沙箱路径仍阻止执行。"
    assert "副作用工具重复调用" not in response.content
    # The third screenshot variant should be blocked internally by the semantic
    # capture budget, while the user-facing answer should be a concrete blocker
    # rather than a raw duplicate-side-effect guardrail message.
    executed_ids = [
        json.loads(call.args[0])["tool_calls"][0]["id"]
        for call in tools.execute_tool_calls.call_args_list
    ]
    assert executed_ids == ["shot1", "shot2"]


@pytest.mark.asyncio
async def test_agent_stops_photo_capture_variants_after_bounded_attempts():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.allowed_paths = []
    tools.get_all_specs.return_value = []

    first_args = json.dumps({
        "command": "mkdir -p ~/.pyclaw/photos && /usr/local/bin/imagesnap -w 1 ~/.pyclaw/photos/photo_$(date +%Y%m%d_%H%M%S).jpg 2>&1 || /opt/homebrew/bin/imagesnap -w 1 ~/.pyclaw/photos/photo_$(date +%Y%m%d_%H%M%S).jpg 2>&1",
        "approved": True,
    }, ensure_ascii=False)
    second_args = json.dumps({
        "command": "mkdir -p ~/.pyclaw/photos && imagesnap -w 1 ~/.pyclaw/photos/snap_$(date +%Y%m%d_%H%M%S).jpg && ls -t ~/.pyclaw/photos/*.jpg | head -1",
    }, ensure_ascii=False)
    third_args = json.dumps({
        "command": "mkdir -p ~/.pyclaw/photos && f=~/.pyclaw/photos/snap_$(date +%Y%m%d_%H%M%S).jpg && imagesnap -w 1 \"$f\" >/dev/null 2>&1 && echo \"$f\"",
    }, ensure_ascii=False)

    model.chat.side_effect = [
        {
            "content": "我来拍照。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "photo1", "function": {"name": "terminal", "arguments": first_args}}],
        },
        {
            "content": "换个 imagesnap 路径再试。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "photo2", "function": {"name": "terminal", "arguments": second_args}}],
        },
        {
            "content": "再换一个文件名。",
            "__tool_calls__": True,
            "tool_calls": [{"id": "photo3", "function": {"name": "terminal", "arguments": third_args}}],
        },
        {"content": "拍照失败：系统相机权限或 imagesnap 环境阻止执行。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{
            "role": "tool",
            "tool_call_id": "photo1",
            "name": "terminal",
            "content": "Command: imagesnap\nExit code: 127\nSTDERR:\nimagesnap: command not found",
            "success": False,
            "metadata": {},
        }],
        [{
            "role": "tool",
            "tool_call_id": "photo2",
            "name": "terminal",
            "content": "Command: imagesnap\nExit code: 1\nSTDERR:\nNo video device found",
            "success": False,
            "metadata": {},
        }],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-photo-bounded-repair"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-photo-bounded-repair",
        channel="telegram",
        channel_user_id="u1",
        session_id="s-photo-bounded-repair",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="拍个照",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.content == "拍照失败：系统相机权限或 imagesnap 环境阻止执行。"
    second_payload = json.loads(tools.execute_tool_calls.call_args_list[1].args[0])
    second_executed_args = json.loads(second_payload["tool_calls"][0]["function"]["arguments"])
    assert second_executed_args["approved"] is True
    assert "副作用工具重复调用" not in response.content
    # The third photo command is a near-identical capture variant and should not
    # execute after the bounded repair attempt has already failed.
    executed_ids = [
        json.loads(call.args[0])["tool_calls"][0]["id"]
        for call in tools.execute_tool_calls.call_args_list
    ]
    assert executed_ids == ["photo1", "photo2"]


@pytest.mark.asyncio
async def test_agent_stops_unlock_script_variants_after_one_desktop_control_attempt():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    first_args = json.dumps({"command": "bash ~/.pyclaw/skills/mac-wake-unlock/unlock.sh", "approved": True})
    variant_args = json.dumps({"command": "echo start && bash ~/.pyclaw/skills/mac-wake-unlock/unlock.sh; echo done", "approved": True})
    model.chat.side_effect = [
        {
            "content": "执行解锁脚本...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "unlock1", "function": {"name": "terminal", "arguments": first_args}}],
        },
        {
            "content": "换个命令形态再试...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "unlock2", "function": {"name": "terminal", "arguments": variant_args}}],
        },
        {"content": "解锁脚本已执行一次，当前检测为已解锁。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "unlock1",
            "name": "terminal",
            "content": "Command: bash ~/.pyclaw/skills/mac-wake-unlock/unlock.sh\nExit code: 0\nSTDOUT:\n✅ 当前已解锁，跳过",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-unlock-script-variant-repeat"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-unlock-script-variant-repeat",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-unlock-script-variant-repeat",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="解锁",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "解锁脚本已执行一次，当前检测为已解锁。"
    assert "副作用工具重复调用" not in response.content


@pytest.mark.asyncio
async def test_agent_corrects_desktop_control_final_that_denies_successful_unlock_observation():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    unlock_args = json.dumps({"command": "bash ~/.pyclaw/skills/mac-wake-unlock/unlock.sh"})
    model.chat.side_effect = [
        {
            "content": "执行解锁脚本...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "unlock1", "function": {"name": "terminal", "arguments": unlock_args}}],
        },
        {"content": "这次也没能真正落到你机器上，请手动跑 unlock.sh。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "unlock1",
            "name": "terminal",
            "content": "Command: bash ~/.pyclaw/skills/mac-wake-unlock/unlock.sh\nExit code: 0\nSTDOUT:\n当前状态: UNLOCKED\n✅ 当前已解锁，跳过",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-unlock-final-correction"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-unlock-final-correction",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-unlock-final-correction",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="解锁",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "解锁脚本已执行；检测到当前已解锁，所以没有输入密码。"
    assert "手动" not in response.content
    assert any("Mac 桌面控制工具已经执行" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_corrects_lock_action_sent_unconfirmed_final():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    lock_args = json.dumps({"command": "bash ~/.pyclaw/skills/mac-lock-unlock/lock.sh"})
    model.chat.side_effect = [
        {
            "content": "执行锁屏脚本...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "lock1", "function": {"name": "terminal", "arguments": lock_args}}],
        },
        {"content": "这轮锁屏没落到机器上，请手动跑 lock.sh。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "lock1",
            "name": "terminal",
            "content": "Command: bash ~/.pyclaw/skills/mac-lock-unlock/lock.sh\nExit code: 0\nSTDOUT:\nACTION_SENT_UNCONFIRMED: 锁屏命令已发送，但本机状态检测不可用，无法自动确认。",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-lock-final-correction"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-lock-final-correction",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-lock-final-correction",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="锁屏",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "锁屏命令已发送；本机状态检测不可用，无法自动确认，但不会再重复执行。"
    assert "手动" not in response.content


def test_desktop_control_finalizer_ignores_stale_lock_observation_for_ppt_task():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    session = Session(session_id="s-stale-lock-ppt", user_id="u1", channel="feishu")

    session.add_message(Message(
        id="m-old-lock",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="锁屏",
        timestamp=datetime.fromisoformat("2026-07-07T18:00:00"),
    ))
    session.add_message(Message(
        id="tool-old-lock",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.TOOL,
        content=(
            "OBSERVATION from terminal:\n"
            "Command: bash ~/.pyclaw/skills/mac-lock-unlock/lock.sh\n"
            "Exit code: 0\n"
            "STDOUT:\n已锁屏"
        ),
        timestamp=datetime.fromisoformat("2026-07-07T18:01:00"),
        metadata={"tool_name": "terminal", "tool_call_id": "lock1"},
    ))
    session.add_message(Message(
        id="m-new-ppt",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="走完整的baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        timestamp=datetime.fromisoformat("2026-07-07T18:52:19"),
    ))
    session.add_message(Message(
        id="tool-current-artifact-dir",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.TOOL,
        content=(
            "OBSERVATION from terminal:\n"
            "Command: mkdir -p /Users/bytedance/.pyclaw/artifacts/rag_ppt && ls -la\n"
            "Exit code: 0"
        ),
        timestamp=datetime.fromisoformat("2026-07-07T18:53:48"),
        metadata={"tool_name": "terminal", "tool_call_id": "mkdir1"},
    ))

    draft = "未观察到目标文件已生成并通过 send_file_to_user 发送，任务未完成。"

    assert agent._prepare_desktop_control_final_content(session, draft) == draft


def test_read_only_diagnostics_with_stderr_redirect_are_navigation_not_side_effect():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    command = 'ioreg -n Root -d1 -a 2>/dev/null | grep -i "CGSSessionScreenIsLocked" || echo "NO MATCH"'

    assert agent._looks_like_terminal_navigation(command)
    assert agent._terminal_command_semantic_kind(json.dumps({"command": command})) == "navigation"
    assert agent._terminal_side_effect_call_key(json.dumps({"command": command})) is None


def test_opencli_remote_query_is_navigation_not_side_effect():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    command = (
        "opencli vephone pod-detail 7602589948898384666 -f json | "
        "python3 -c 'import sys,json; print(json.load(sys.stdin).get(\"node\"))' | head -30"
    )

    assert agent._looks_like_terminal_navigation(command)
    assert agent._terminal_command_semantic_kind(json.dumps({"command": command})) == "navigation"
    assert agent._terminal_side_effect_call_key(json.dumps({"command": command})) is None


def test_opencli_browser_open_still_counts_as_side_effect():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    command = "opencli browser open https://example.com"

    assert not agent._looks_like_terminal_navigation(command)
    assert agent._terminal_side_effect_call_key(json.dumps({"command": command})) == "terminal:semantic:open"


def test_direct_execution_of_changed_script_counts_as_validation_result():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())
    changed_files = {"/Users/bytedance/.pyclaw/skills/mac-lock-unlock/lock.sh"}
    validation_results: list[str] = []
    build_results: list[str] = []

    agent._record_coding_tool_effects(
        tool_results=[{
            "role": "tool",
            "tool_call_id": "run1",
            "name": "terminal",
            "content": "Command: cd ~/.pyclaw/skills/mac-lock-unlock && bash lock.sh 2>&1\nExit code: 1\nSTDOUT:\nfailed",
            "success": False,
            "metadata": {},
        }],
        changed_files=changed_files,
        validation_results=validation_results,
        build_results=build_results,
    )

    assert validation_results == ["FAIL: cd ~/.pyclaw/skills/mac-lock-unlock && bash lock.sh 2>&1"]
    assert build_results == []


@pytest.mark.asyncio
async def test_agent_stops_repeated_mac_unlock_script_for_simple_wechat_request():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    unlock_args = json.dumps({"command": "~/.pyclaw/bin/unlock.sh", "approved": True})
    model.chat.side_effect = [
        {
            "content": "解锁...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "unlock1", "function": {"name": "terminal", "arguments": unlock_args}}],
        },
        {
            "content": "再试一次解锁...",
            "__tool_calls__": True,
            "tool_calls": [{"id": "unlock2", "function": {"name": "terminal", "arguments": unlock_args}}],
        },
        {"content": "已触发解锁。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "unlock1",
            "name": "terminal",
            "content": "Command: ~/.pyclaw/bin/unlock.sh\nExit code: 0",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-unlock-simple"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-terminal-unlock-simple",
        channel="wechat",
        channel_user_id="u1",
        session_id="s-terminal-unlock-simple",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="解锁",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "已触发解锁。"
    assert "副作用工具重复调用" not in response.content
    assert any(
        "本轮只有重复的副作用工具调用" in m.content
        and "terminal:mac_desktop_control:unlock" in m.content
        for m in session.messages
    )


@pytest.mark.asyncio
async def test_agent_allows_corrected_retry_after_failed_side_effect_attempt():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    command = 'mkdir -p "$HOME/.pyclaw/cron_history" && osascript -e "display notification \\"ok\\""'
    model.chat.side_effect = [
        {
            "content": "先执行通知命令...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "notify1",
                "function": {
                    "name": "terminal",
                    "arguments": json.dumps({"command": command, "timeout": 10}),
                },
            }],
        },
        {
            "content": "修正 approval 后重试...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "notify2",
                "function": {
                    "name": "terminal",
                    "arguments": json.dumps({"command": command, "timeout": 10, "approved": True}),
                },
            }],
        },
        {"content": "通知已完成。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [
            {
                "role": "tool",
                "tool_call_id": "notify1",
                "name": "terminal",
                "content": "⚠️ 检测到有副作用的指令，请设置 approved=True 后重试。",
                "success": False,
                "metadata": {},
            }
        ],
        [
            {
                "role": "tool",
                "tool_call_id": "notify2",
                "name": "terminal",
                "content": "notification sent",
                "success": True,
                "metadata": {},
            }
        ],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-side-effect-retry"
    session.channel = "cron"
    session.channel_user_id = "job_1"
    session.user_id = "job_1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-side-effect-retry", channel="cron", channel_user_id="job_1", session_id="s-side-effect-retry",
        type=MessageType.TEXT, role=MessageRole.USER, content="发送一次通知",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.content == "通知已完成。"
    assert "副作用工具重复调用" not in response.content


@pytest.mark.asyncio
async def test_agent_allows_distinct_generic_side_effect_calls():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "发送第一条消息。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send1",
                "function": {
                    "name": "send_message",
                    "arguments": json.dumps({"to": "user-1", "text": "第一条"}, ensure_ascii=False),
                },
            }],
        },
        {
            "content": "发送第二条消息。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send2",
                "function": {
                    "name": "send_message",
                    "arguments": json.dumps({"to": "user-1", "text": "第二条"}, ensure_ascii=False),
                },
            }],
        },
        {"content": "两条消息都已发送。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "send1", "name": "send_message", "content": "ok1", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "send2", "name": "send_message", "content": "ok2", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-generic-side-effect"
    session.channel = "wechat"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-generic-side-effect", channel="wechat", channel_user_id="u1", session_id="s-generic-side-effect",
        type=MessageType.TEXT, role=MessageRole.USER, content="分别发送两条不同消息",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.content == "两条消息都已发送。"
    assert "副作用工具重复调用" not in response.content
    assert not any("试图重复执行" in m.content and "send_message" in m.content for m in session.messages)


def test_sanitize_user_facing_content_strips_side_effect_guardrail_prefixes():
    agent = Agent(MagicMock(), MagicMock(), MagicMock())

    cleaned = agent._sanitize_user_facing_content(
        "⚠️  检测到副作用工具重复调用（terminal:abc），我已停止继续执行。\n\n"
        "已根据现有结果完成处理。"
    )

    assert cleaned == "已根据现有结果完成处理。"


@pytest.mark.asyncio
async def test_agent_allows_repeated_read_only_cronjob_list():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "查看任务...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "list1",
                "function": {"name": "cronjob", "arguments": '{"action": "list"}'}
            }],
        },
        {
            "content": "再次确认任务...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "list2",
                "function": {"name": "cronjob", "arguments": '{"action": "list"}'}
            }],
        },
        {"content": "任务状态已确认。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "list",
            "name": "cronjob",
            "content": "[]",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-list"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-cron-list", channel="telegram", channel_user_id="u1", session_id="s-cron-list",
        type=MessageType.TEXT, role=MessageRole.USER, content="看看任务状态"
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.content == "任务状态已确认。"


@pytest.mark.asyncio
async def test_agent_requires_extract_after_search_for_current_events():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [
        {"name": "web_search", "description": "search", "parameters": {}},
        {"name": "web_extract", "description": "extract", "parameters": {}},
    ]

    model.chat.side_effect = [
        {
            "content": "先搜索最新赛果。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "search1",
                "function": {"name": "web_search", "arguments": '{"query":"世界杯最新赛果"}'},
            }],
        },
        {"content": "搜索结果显示 A 队赢了。", "__tool_calls__": False},
        {
            "content": "需要读取来源确认。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "extract1",
                "function": {"name": "web_extract", "arguments": '{"urls":["https://example.com/match"]}'},
            }],
        },
        {"content": "基于权威来源确认：A 队获胜。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [
            {
                "role": "tool",
                "tool_call_id": "search1",
                "name": "web_search",
                "content": "https://example.com/match",
                "success": True,
                "metadata": {},
            }
        ],
        [
            {
                "role": "tool",
                "tool_call_id": "extract1",
                "name": "web_extract",
                "content": "官方赛果：A 队获胜。",
                "success": True,
                "metadata": {},
            }
        ],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-current-extract"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-current-extract", channel="feishu", channel_user_id="u1", session_id="s-current-extract",
        type=MessageType.TEXT, role=MessageRole.USER, content="世界杯最新赛事",
    )

    response = await agent.process_message(user_msg)

    assert response.content == "基于权威来源确认：A 队获胜。"
    executed_names = []
    for call in tools.execute_tool_calls.await_args_list:
        payload = json.loads(call.args[0])
        executed_names.extend(tc["function"]["name"] for tc in payload["tool_calls"])
    assert executed_names == ["web_search", "web_extract"]
    assert any("have not extracted any source page" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_quality_gate_repairs_unresolved_requested_facts():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [
        {"name": "web_search", "description": "search", "parameters": {}},
        {"name": "web_extract", "description": "extract", "parameters": {}},
    ]

    pending_draft = """
🏆 世界杯晚报
🇵🇹 葡萄牙 vs 🇨🇩 刚果（金） — K组
⏰ 01:00 CST 已完赛 | ⚠️ 比分待确认
""".strip()
    model.chat.side_effect = [
        {
            "content": "先搜索赛果。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "search1",
                "function": {"name": "web_search", "arguments": '{"query":"世界杯今日赛果"}'},
            }],
        },
        {
            "content": "读取权威来源。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "extract1",
                "function": {"name": "web_extract", "arguments": '{"urls":["https://example.com/scores"]}'},
            }],
        },
        {"content": pending_draft, "__tool_calls__": False},
        {
            "content": "定向查询缺失比分。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "search2",
                "function": {"name": "web_search", "arguments": '{"query":"Portugal vs DR Congo final score"}'},
            }],
        },
        {"content": "确认：葡萄牙 2-0 刚果（金）。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "search1", "name": "web_search", "content": "https://example.com/scores", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "extract1", "name": "web_extract", "content": "赛程页缺少比分", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "search2", "name": "web_search", "content": "Portugal 2-0 DR Congo FT", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-sports-pending"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-sports-pending", channel="telegram", channel_user_id="u1", session_id="s-sports-pending",
        type=MessageType.TEXT, role=MessageRole.USER, content="世界杯晚报，给我今日赛果",
    )

    response = await agent.process_message(user_msg)

    assert response.content == "确认：葡萄牙 2-0 刚果（金）。"
    executed_names = []
    for call in tools.execute_tool_calls.await_args_list:
        payload = json.loads(call.args[0])
        executed_names.extend(tc["function"]["name"] for tc in payload["tool_calls"])
    assert executed_names == ["web_search", "web_extract", "web_search"]
    assert any("unresolved_requested_facts" in m.content for m in session.messages)
    assert any("targeted lookups/extractions" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_quality_gate_repairs_stale_cutoff_for_live_tasks():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [
        {"name": "web_search", "description": "search", "parameters": {}},
    ]

    stale_draft = """
## 最新赛事汇总

## ✅ 已完赛比分（截至6月18日）
| 组别 | 比分 |
|---|---|
| A组 | A 1-0 B |
""".strip()
    model.chat.side_effect = [
        {
            "content": "先查最新数据。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "search1",
                "function": {"name": "web_search", "arguments": '{"query":"today latest scores"}'},
            }],
        },
        {"content": stale_draft, "__tool_calls__": False},
        {
            "content": "按当前日期定向修复。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "search2",
                "function": {"name": "web_search", "arguments": '{"query":"2026-06-22 final scores"}'},
            }],
        },
        {"content": "确认：6月22日最新赛果已更新。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "search1", "name": "web_search", "content": "old page as of Jun 18", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "search2", "name": "web_search", "content": "fresh Jun 22 results", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-stale-cutoff"
    session.channel = "cron"
    session.channel_user_id = "job_1"
    session.user_id = "job_1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-stale-cutoff", channel="cron", channel_user_id="job_1", session_id="s-stale-cutoff",
        type=MessageType.TEXT, role=MessageRole.USER,
        content="当前执行时间：2026-06-22 18:00:00 CST+0800。请整理今日最新赛果和明日赛程。",
    )

    response = await agent.process_message(user_msg)

    assert response.content == "确认：6月22日最新赛果已更新。"
    executed_names = []
    for call in tools.execute_tool_calls.await_args_list:
        payload = json.loads(call.args[0])
        executed_names.extend(tc["function"]["name"] for tc in payload["tool_calls"])
    assert executed_names == ["web_search", "web_search"]
    assert any("stale_cutoff_date" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_quality_gate_is_not_sports_specific():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [
        {"name": "web_search", "description": "search", "parameters": {}},
        {"name": "web_extract", "description": "extract", "parameters": {}},
    ]

    model.chat.side_effect = [
        {
            "content": "先查价格。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "search-price",
                "function": {"name": "web_search", "arguments": '{"query":"产品 X 今日价格"}'},
            }],
        },
        {
            "content": "读取来源。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "extract-price",
                "function": {"name": "web_extract", "arguments": '{"urls":["https://example.com/price"]}'},
            }],
        },
        {"content": "产品 X 今日价格：暂未获取到完整报价。", "__tool_calls__": False},
        {
            "content": "定向查缺失报价。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "search-price-2",
                "function": {"name": "web_search", "arguments": '{"query":"产品 X 官方价格 今日"}'},
            }],
        },
        {"content": "产品 X 今日官方价格确认：199 元。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "search-price", "name": "web_search", "content": "https://example.com/price", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "extract-price", "name": "web_extract", "content": "页面缺少报价", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "search-price-2", "name": "web_search", "content": "官方价格 199 元", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-price-pending"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-price-pending", channel="telegram", channel_user_id="u1", session_id="s-price-pending",
        type=MessageType.TEXT, role=MessageRole.USER, content="查一下产品 X 今日价格",
    )

    response = await agent.process_message(user_msg)

    assert response.content == "产品 X 今日官方价格确认：199 元。"
    assert any("unresolved_requested_facts" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_does_not_require_extract_for_non_current_search_task():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [
        {"name": "web_search", "description": "search", "parameters": {}},
        {"name": "web_extract", "description": "extract", "parameters": {}},
    ]

    model.chat.side_effect = [
        {
            "content": "搜索概念资料。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "search1",
                "function": {"name": "web_search", "arguments": '{"query":"ReAct paper"}'},
            }],
        },
        {"content": "ReAct 是一种推理与行动交替的方法。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "search1",
            "name": "web_search",
            "content": "ReAct paper search result",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-non-current-search"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-non-current-search", channel="feishu", channel_user_id="u1", session_id="s-non-current-search",
        type=MessageType.TEXT, role=MessageRole.USER, content="解释 ReAct agent 是什么",
    )

    response = await agent.process_message(user_msg)

    assert response.content == "ReAct 是一种推理与行动交替的方法。"
    assert model.chat.call_count == 2
    assert not any("have not extracted any source page" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_allows_multiple_read_only_lark_cli_terminal_calls():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "先查 wiki 元数据...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "wiki_meta",
                "function": {
                    "name": "terminal",
                    "arguments": json.dumps({
                        "command": "lark-cli wiki spaces get_node --params '{\"token\":\"WIKI_TOKEN_EXAMPLE\"}'",
                        "timeout": 15,
                    }),
                },
            }],
        },
        {
            "content": "再读取正文...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "doc_fetch",
                "function": {
                    "name": "terminal",
                    "arguments": json.dumps({
                        "command": "lark-cli docs +fetch --api-version v2 --doc https://example.larkoffice.com/wiki/WIKI_TOKEN_EXAMPLE --doc-format markdown",
                        "timeout": 30,
                    }),
                },
            }],
        },
        {"content": "文章已读完。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "lark",
            "name": "terminal",
            "content": "ok",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-lark-read"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-lark-read", channel="telegram", channel_user_id="u1", session_id="s-lark-read",
        type=MessageType.TEXT, role=MessageRole.USER,
        content="你给我阅读这篇文章 https://example.larkoffice.com/wiki/WIKI_TOKEN_EXAMPLE",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.content == "文章已读完。"
    assert "副作用工具重复调用" not in response.content


@pytest.mark.asyncio
async def test_agent_allows_triggering_multiple_distinct_cron_jobs_once():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "触发多个任务...",
            "__tool_calls__": True,
            "tool_calls": [
                {
                    "id": "trigger1",
                    "function": {"name": "cronjob", "arguments": '{"action": "trigger", "job_id": "job-a"}'},
                },
                {
                    "id": "trigger2",
                    "function": {"name": "cronjob", "arguments": '{"action": "trigger", "job_id": "job-b"}'},
                },
            ],
        },
        {"content": "两个任务都已触发。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {"role": "tool", "tool_call_id": "trigger1", "name": "cronjob", "content": "job-a triggered", "success": True, "metadata": {}},
        {"role": "tool", "tool_call_id": "trigger2", "name": "cronjob", "content": "job-b triggered", "success": True, "metadata": {}},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-trigger-batch"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-cron-trigger-batch", channel="telegram", channel_user_id="u1", session_id="s-cron-trigger-batch",
        type=MessageType.TEXT, role=MessageRole.USER, content="触发所有任务"
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "两个任务都已触发。"


@pytest.mark.asyncio
async def test_agent_synthesizes_after_repeated_executed_cron_side_effect():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "触发任务...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "trigger",
                "function": {"name": "cronjob", "arguments": '{"action": "trigger", "job_id": "job-a"}'},
            }],
        },
        {
            "content": "重复触发同一任务...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "trigger-again",
                "function": {"name": "cronjob", "arguments": '{"action": "trigger", "job_id": "job-a"}'},
            }],
        },
        {"content": "任务已触发完成。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {"role": "tool", "tool_call_id": "trigger", "name": "cronjob", "content": "job-a triggered", "success": True, "metadata": {}},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-retrigger"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-cron-retrigger", channel="telegram", channel_user_id="u1", session_id="s-cron-retrigger",
        type=MessageType.TEXT, role=MessageRole.USER, content="触发任务"
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "任务已触发完成。"
    assert "副作用工具重复调用" not in response.content
    assert model.chat.await_args_list[-1].kwargs["tools"] is None
    assert any("本轮只有重复的副作用工具调用" in m.content and "cronjob:trigger:job-a" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_agent_synthesizes_after_repeated_executed_cron_update():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    update_args = json.dumps({"action": "update", "job_id": "90e15343", "prompt": "new prompt"})
    model.chat.side_effect = [
        {
            "content": "更新任务...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "update1",
                "function": {"name": "cronjob", "arguments": update_args},
            }],
        },
        {
            "content": "再次确认更新...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "update2",
                "function": {"name": "cronjob", "arguments": update_args},
            }],
        },
        {"content": "任务已更新完成。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {"role": "tool", "tool_call_id": "update1", "name": "cronjob", "content": "job updated", "success": True, "metadata": {}},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-update-repeat"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-cron-update-repeat", channel="feishu", channel_user_id="u1", session_id="s-cron-update-repeat",
        type=MessageType.TEXT, role=MessageRole.USER, content="更新这个定时任务",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert response.content == "任务已更新完成。"
    assert "副作用工具重复调用" not in response.content
    assert "cronjob:update:90e15343" not in response.content


@pytest.mark.asyncio
async def test_agent_distinguishes_distinct_cron_create_payloads():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "创建第一个任务...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "create1",
                "function": {
                    "name": "cronjob",
                    "arguments": json.dumps({
                        "action": "create",
                        "name": "daily codex push",
                        "schedule": "0 8 * * *",
                        "prompt": "Push one Codex article.",
                    }),
                },
            }],
        },
        {
            "content": "创建另一个任务...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "create2",
                "function": {
                    "name": "cronjob",
                    "arguments": json.dumps({
                        "action": "create",
                        "name": "daily agent lesson",
                        "schedule": "0 9 * * *",
                        "prompt": "Push one AI agent lesson.",
                    }),
                },
            }],
        },
        {"content": "两个任务已创建。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {"role": "tool", "tool_call_id": "create", "name": "cronjob", "content": "created", "success": True, "metadata": {}},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-cron-create-distinct"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=30)
    user_msg = Message(
        id="m-cron-create-distinct", channel="telegram", channel_user_id="u1", session_id="s-cron-create-distinct",
        type=MessageType.TEXT, role=MessageRole.USER, content="创建两个不同任务"
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 2
    assert response.content == "两个任务已创建。"
    assert "cronjob:create:<no-job-id>" not in response.content
    assert "副作用工具重复调用" not in response.content


@pytest.mark.asyncio
async def test_agent_converts_tool_executor_exception_to_observation():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(side_effect=RuntimeError("executor down"))
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    first_resp = {
        "content": "Need a tool",
        "__tool_calls__": True,
        "tool_calls": [{
            "id": "call1",
            "function": {"name": "terminal", "arguments": '{"command": "pwd"}'}
        }]
    }
    second_resp = {"content": "Recovered after tool executor failure."}
    model.chat.side_effect = [first_resp, second_resp]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-tool-error"
    session.channel = "t"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions)
    user_msg = Message(
        id="m-tool-error", channel="t", channel_user_id="u1", session_id="s-tool-error",
        type=MessageType.TEXT, role=MessageRole.USER, content="Run pwd"
    )

    response = await agent.process_message(user_msg)

    assert model.chat.call_count == 2
    second_call_messages = model.chat.call_args_list[1][1]["messages"]
    tool_msg = next(m for m in second_call_messages if m["role"] == "tool")
    assert "Tool execution framework error: RuntimeError: executor down" in tool_msg["content"]
    assert response.content == "Recovered after tool executor failure."


@pytest.mark.asyncio
async def test_agent_adds_current_task_boundary_after_history_summary():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-boundary"
    session.channel = "telegram"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {"history_summary": "Previous pending task: implement /reset."}
    session.get_history.side_effect = lambda limit=10: [
        {"role": "system", "content": "I am PyClaw"},
        {
            "role": "system",
            "content": (
                "<read_only_conversation_summary>\n"
                "Previous pending task: implement /reset.\n"
                "</read_only_conversation_summary>"
            ),
        },
        *[m.to_llm_format() for m in session.messages],
    ]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session
    model.chat.return_value = {"content": "这是在回答当前问题。"}

    agent = Agent(model, tools, sessions)
    user_msg = Message(
        id="m-boundary", channel="telegram", channel_user_id="u1", session_id="s-boundary",
        type=MessageType.TEXT, role=MessageRole.USER, content="这个问题又是为啥呢？"
    )

    await agent.process_message(user_msg)

    first_call_messages = model.chat.call_args_list[0][1]["messages"]
    boundary = first_call_messages[-1]
    assert boundary["role"] == "system"
    assert "<current_task_boundary>" in boundary["content"]
    assert "Only the latest user message below defines the current task" in boundary["content"]
    assert "这个问题又是为啥呢？" in boundary["content"]
    assert "implement /reset" not in boundary["content"]

@pytest.mark.asyncio
async def test_agent_clear_session_on_new_command():
    model = AsyncMock()
    tools = MagicMock()
    tools.skills_dirs = []
    
    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s3"
    session.channel = "t"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    
    sessions.get_or_create.return_value = session
    
    agent = Agent(model, tools, sessions)
    
    # User sends /new command
    user_msg = Message(
        id="m3", channel="t", channel_user_id="u1", session_id="s3",
        type=MessageType.TEXT, role=MessageRole.USER, content="/new"
    )
    
    response = await agent.process_message(user_msg)
    
    # clear_session should have been called
    sessions.clear_session.assert_called_once_with(session)
    
    # The model should not have been called (LLM loop skipped)
    model.chat.assert_not_called()
    
    # We should get a reset confirmation reply
    assert "会话已重置" in response.content


@pytest.mark.asyncio
async def test_agent_clear_session_on_reset_command():
    model = AsyncMock()
    tools = MagicMock()
    tools.skills_dirs = []

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-reset"
    session.channel = "t"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}

    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions)
    user_msg = Message(
        id="m-reset", channel="t", channel_user_id="u1", session_id="s-reset",
        type=MessageType.TEXT, role=MessageRole.USER, content="/reset"
    )

    response = await agent.process_message(user_msg)

    sessions.clear_session.assert_called_once_with(session)
    model.chat.assert_not_called()
    assert "会话已重置" in response.content


def test_agent_reset_command_parser_accepts_lark_mentions():
    agent = Agent(AsyncMock(), MagicMock(), AsyncMock())

    assert agent._is_session_reset_command("/new")
    assert agent._is_session_reset_command("/reset")
    assert agent._is_session_reset_command("/new@PyClawBot")
    assert agent._is_session_reset_command("@PyClaw /new")
    assert agent._is_session_reset_command("@PyClaw /reset@PyClawBot")
    assert not agent._is_session_reset_command("已经发过 /new 了")
    assert not agent._is_session_reset_command("请解释 /reset 的作用")


@pytest.mark.asyncio
async def test_telegram_new_command_is_forwarded_to_agent():
    from pyclaw.channels.telegram import TelegramChannel

    channel = TelegramChannel(token="test-token")
    handled_messages = []

    async def handle_message(message):
        handled_messages.append(message)

    channel.on_message(handle_message)

    update = MagicMock()
    update.message.text = "/new"
    update.message.document = None
    update.message.caption = None
    update.effective_user.id = 12345
    context = MagicMock()

    await channel._on_command(update, context)

    assert len(handled_messages) == 1
    assert handled_messages[0].channel == "telegram"
    assert handled_messages[0].channel_user_id == "12345"
    assert handled_messages[0].session_id == "telegram:12345"
    assert handled_messages[0].content == "/new"
    update.message.reply_text.assert_not_called()


@pytest.mark.asyncio
async def test_telegram_bot_qualified_new_command_is_forwarded_to_agent():
    from pyclaw.channels.telegram import TelegramChannel

    channel = TelegramChannel(token="test-token")
    handled_messages = []

    async def handle_message(message):
        handled_messages.append(message)

    channel.on_message(handle_message)

    update = MagicMock()
    update.message.text = "/new@PyClawBot"
    update.message.document = None
    update.message.caption = None
    update.effective_user.id = 12345
    context = MagicMock()

    await channel._on_command(update, context)

    assert len(handled_messages) == 1
    assert handled_messages[0].content == "/new"
    update.message.reply_text.assert_not_called()


@pytest.mark.asyncio
async def test_telegram_reset_command_is_forwarded_to_agent():
    from pyclaw.channels.telegram import TelegramChannel

    channel = TelegramChannel(token="test-token")
    handled_messages = []

    async def handle_message(message):
        handled_messages.append(message)

    channel.on_message(handle_message)

    update = MagicMock()
    update.message.text = "/reset@PyClawBot"
    update.message.document = None
    update.message.caption = None
    update.effective_user.id = 12345
    context = MagicMock()

    await channel._on_command(update, context)

    assert len(handled_messages) == 1
    assert handled_messages[0].content == "/reset"
    update.message.reply_text.assert_not_called()


def test_telegram_readable_formatter_compacts_single_line_code_blocks():
    from pyclaw.channels.telegram import TelegramChannel

    channel = TelegramChannel(token="test-token")
    source = """我也已经执行过语法检查：

```bash
python3 -m py_compile pyclaw/core/agent.py pyclaw/channels/telegram.py
```

结果通过。

现在重启 Bot / 服务后，发送：

```text
/reset
```
"""

    readable = channel._format_telegram_readable(source)

    assert "```" not in readable
    assert "`python3 -m py_compile pyclaw/core/agent.py pyclaw/channels/telegram.py`" in readable
    assert "`/reset`" in readable
    assert "结果：✅ 通过" in readable


def test_telegram_markdown_does_not_italicize_underscores():
    from pyclaw.channels.telegram import TelegramChannel

    channel = TelegramChannel(token="test-token")

    formatted = channel._format_markdown("检查 foo_bar_baz 和 `inline_code`。")

    assert "foo_bar_baz" in formatted
    assert "<i>" not in formatted
    assert "<code>inline_code</code>" in formatted


def test_telegram_formatter_wraps_code_explanation_lines():
    from pyclaw.channels.telegram import TelegramChannel

    channel = TelegramChannel(token="test-token")
    source = '''## 拆解这行代码

print(f"🔎 [web_search] provider={provider.name} failed: {type(e).__name__}: {e}")

### 1. {provider.name}

provider.name

假设：

provider.name = "brave"

---

### 2. {type(e).__name__}
'''

    readable = channel._format_telegram_readable(source)
    formatted = channel._format_markdown(readable)

    assert "<b>拆解这行代码</b>" in formatted
    assert "<b>1. {provider.name}</b>" in formatted
    assert '<code>print(f"🔎 [web_search] provider={provider.name} failed: {type(e).__name__}: {e}")</code>' in formatted
    assert "<code>provider.name</code>" in formatted
    assert '<code>provider.name = "brave"</code>' in formatted
    assert "type(e).<b>name</b>" not in formatted
    assert "##" not in formatted
    assert "###" not in formatted


@pytest.mark.asyncio
async def test_patch_first_gate_rejects_implementation_without_diff():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = [
        {"content": "方案：需要修改文件。", "__tool_calls__": False},
        {"content": "已修改并完成。", "__tool_calls__": False},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-patch-first"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m-patch-first", channel="feishu", channel_user_id="u1", session_id="s-patch-first",
        type=MessageType.TEXT, role=MessageRole.USER, content="请帮我实现这个功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert response.content == "已修改并完成。"
    assert any("Patch-first quality gate failed" in m.content for m in session.messages)
    assert "方案：需要修改文件。" not in response.content


@pytest.mark.asyncio
async def test_verification_gate_requires_validation_after_code_change():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = [
        {
            "content": "修改文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "edit1",
                "function": {"name": "edit_file", "arguments": '{"path":"app.py","old":"a","new":"b"}'},
            }],
        },
        {"content": "已完成修改。", "__tool_calls__": False},
        {
            "content": "运行测试。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "test1",
                "function": {"name": "terminal", "arguments": '{"command":"pytest tests/test_app.py -q"}'},
            }],
        },
        {
            "content": "运行编译。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "build1",
                "function": {"name": "terminal", "arguments": '{"command":"python -m py_compile app.py"}'},
            }],
        },
        {"content": "已完成修改。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "edit1", "name": "edit_file", "content": "File edited: app.py\n--- diff", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "test1", "name": "terminal", "content": "Command: pytest tests/test_app.py -q\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "build1", "name": "terminal", "content": "Command: python -m py_compile app.py\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-verification"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=8)
    user_msg = Message(
        id="m-verification", channel="feishu", channel_user_id="u1", session_id="s-verification",
        type=MessageType.TEXT, role=MessageRole.USER, content="请实现功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert any("Verification gate failed" in m.content for m in session.messages)
    assert "验证结果：PASS: pytest tests/test_app.py -q; PASS: python -m py_compile app.py" in response.content



@pytest.mark.asyncio
async def test_coding_task_status_is_added_to_final_answer_after_validation():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = [
        {
            "content": "修改文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "edit1",
                "function": {"name": "edit_file", "arguments": '{"path":"app.py","old":"a","new":"b"}'},
            }],
        },
        {
            "content": "运行测试。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "test1",
                "function": {"name": "terminal", "arguments": '{"command":"pytest tests/test_app.py -q"}'},
            }],
        },
        {
            "content": "运行编译。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "build1",
                "function": {"name": "terminal", "arguments": '{"command":"python -m py_compile app.py"}'},
            }],
        },
        {"content": "已完成修改。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "edit1", "name": "edit_file", "content": "File edited: app.py\n--- diff", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "test1", "name": "terminal", "content": "Command: pytest tests/test_app.py -q\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "build1", "name": "terminal", "content": "Command: python -m py_compile app.py\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-task-status"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-task-status", channel="feishu", channel_user_id="u1", session_id="s-task-status",
        type=MessageType.TEXT, role=MessageRole.USER, content="请实现功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert "任务清单：" in response.content
    assert "[x] 完成代码修改" in response.content
    assert "[x] 运行最小验证" in response.content
    assert "[x] 尝试编译/构建" in response.content
    assert "PASS: python -m py_compile app.py" in response.content


@pytest.mark.asyncio
async def test_build_gate_requires_compile_after_validation_passes():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = [
        {
            "content": "修改文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "edit1",
                "function": {"name": "edit_file", "arguments": '{"path":"app.py","old":"a","new":"b"}'},
            }],
        },
        {
            "content": "运行测试。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "test1",
                "function": {"name": "terminal", "arguments": '{"command":"pytest tests/test_app.py -q"}'},
            }],
        },
        {"content": "已完成修改。", "__tool_calls__": False},
        {
            "content": "运行编译。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "build1",
                "function": {"name": "terminal", "arguments": '{"command":"python -m py_compile app.py"}'},
            }],
        },
        {"content": "已完成修改。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "edit1", "name": "edit_file", "content": "File edited: app.py\n--- diff", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "test1", "name": "terminal", "content": "Command: pytest tests/test_app.py -q\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "build1", "name": "terminal", "content": "Command: python -m py_compile app.py\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-build-gate"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-build-gate", channel="feishu", channel_user_id="u1", session_id="s-build-gate",
        type=MessageType.TEXT, role=MessageRole.USER, content="请实现功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert any("Build gate failed" in m.content for m in session.messages)
    assert "PASS: python -m py_compile app.py" in response.content


@pytest.mark.asyncio
async def test_coding_repeated_navigation_pivots_to_edit_instead_of_finalizing():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "继续定位代码 1。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "read1",
                "function": {"name": "read_lines", "arguments": '{"path":"app.py","start_line":1,"end_line":20}'},
            }],
        },
        {
            "content": "继续定位代码 2。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "read2",
                "function": {"name": "read_lines", "arguments": '{"path":"app.py","start_line":21,"end_line":40}'},
            }],
        },
        {
            "content": "继续定位代码 3。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "read3",
                "function": {"name": "read_lines", "arguments": '{"path":"app.py","start_line":41,"end_line":60}'},
            }],
        },
        {
            "content": "开始修改。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "edit1",
                "function": {"name": "edit_file", "arguments": '{"path":"app.py","old":"a","new":"b"}'},
            }],
        },
        {
            "content": "运行测试。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "test1",
                "function": {"name": "terminal", "arguments": '{"command":"pytest tests/test_app.py -q"}'},
            }],
        },
        {
            "content": "运行编译。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "build1",
                "function": {"name": "terminal", "arguments": '{"command":"python -m py_compile app.py"}'},
            }],
        },
        {"content": "已完成修改。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "read1", "name": "read_lines", "content": "code 1", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "read2", "name": "read_lines", "content": "code 2", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "edit1", "name": "edit_file", "content": "File edited: app.py\n--- diff", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "test1", "name": "terminal", "content": "Command: pytest tests/test_app.py -q\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "build1", "name": "terminal", "content": "Command: python -m py_compile app.py\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-navigation-pivot"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {"coding_repeated_tool_limit": 2}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-navigation-pivot", channel="feishu", channel_user_id="u1", session_id="s-navigation-pivot",
        type=MessageType.TEXT, role=MessageRole.USER, content="请实现功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert any("Repeated code navigation detected" in m.content for m in session.messages)
    assert not any("Tool usage must stop now" in m.content and "read_lines" in m.content for m in session.messages)
    executed_names = [call.args[0] for call in tools.execute_tool_calls.call_args_list]
    assert any('"name": "edit_file"' in payload for payload in executed_names)
    assert tools.execute_tool_calls.call_count == 5
    assert "验证结果：PASS: pytest tests/test_app.py -q; PASS: python -m py_compile app.py" in response.content


@pytest.mark.asyncio
async def test_terminal_navigation_repetition_pivots_for_coding_task():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "用 rg 定位代码 1。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "rg1",
                "function": {"name": "terminal", "arguments": '{"command":"rg -n Button app/src"}'},
            }],
        },
        {
            "content": "用 sed 看代码 2。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "sed1",
                "function": {"name": "terminal", "arguments": '{"command":"sed -n 1,80p app.py"}'},
            }],
        },
        {
            "content": "用 find 再找代码 3。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "find1",
                "function": {"name": "terminal", "arguments": '{"command":"find . -name app.py"}'},
            }],
        },
        {
            "content": "开始修改。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "edit1",
                "function": {"name": "edit_file", "arguments": '{"path":"app.py","old":"a","new":"b"}'},
            }],
        },
        {
            "content": "运行测试。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "test1",
                "function": {"name": "terminal", "arguments": '{"command":"pytest tests/test_app.py -q"}'},
            }],
        },
        {
            "content": "运行编译。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "build1",
                "function": {"name": "terminal", "arguments": '{"command":"python -m py_compile app.py"}'},
            }],
        },
        {"content": "已完成修改。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "rg1", "name": "terminal", "content": "Command: rg -n Button app/src\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "sed1", "name": "terminal", "content": "Command: sed -n 1,80p app.py\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "edit1", "name": "edit_file", "content": "File edited: app.py\n--- diff", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "test1", "name": "terminal", "content": "Command: pytest tests/test_app.py -q\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "build1", "name": "terminal", "content": "Command: python -m py_compile app.py\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-navigation-pivot"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {"coding_repeated_tool_limit": 2}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-terminal-navigation-pivot", channel="feishu", channel_user_id="u1", session_id="s-terminal-navigation-pivot",
        type=MessageType.TEXT, role=MessageRole.USER, content="请实现功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert any("Repeated code navigation detected" in m.content and "terminal" in m.content for m in session.messages)
    assert not any("Tool usage must stop now" in m.content and "terminal" in m.content for m in session.messages)
    assert tools.execute_tool_calls.call_count == 5
    assert "验证结果：PASS: pytest tests/test_app.py -q; PASS: python -m py_compile app.py" in response.content


@pytest.mark.asyncio
async def test_terminal_validation_can_rerun_after_code_change():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    pytest_args = json.dumps({"command": "pytest tests/test_app.py -q"})
    model.chat.side_effect = [
        {"content": "先跑测试。", "__tool_calls__": True, "tool_calls": [{"id": "test1", "function": {"name": "terminal", "arguments": pytest_args}}]},
        {"content": "修改文件。", "__tool_calls__": True, "tool_calls": [{"id": "edit1", "function": {"name": "edit_file", "arguments": '{"path":"app.py","old":"a","new":"b"}'}}]},
        {"content": "重跑同一个测试。", "__tool_calls__": True, "tool_calls": [{"id": "test2", "function": {"name": "terminal", "arguments": pytest_args}}]},
        {"content": "运行编译。", "__tool_calls__": True, "tool_calls": [{"id": "build1", "function": {"name": "terminal", "arguments": '{"command":"python -m py_compile app.py"}'}}]},
        {"content": "已完成修改。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "test1", "name": "terminal", "content": "Command: pytest tests/test_app.py -q\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "edit1", "name": "edit_file", "content": "File edited: app.py\n--- diff", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "test2", "name": "terminal", "content": "Command: pytest tests/test_app.py -q\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "build1", "name": "terminal", "content": "Command: python -m py_compile app.py\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-validation-rerun"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-validation-rerun", channel="feishu", channel_user_id="u1", session_id="s-validation-rerun",
        type=MessageType.TEXT, role=MessageRole.USER, content="请实现功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 4
    assert "副作用工具重复调用" not in response.content
    assert not any("试图重复执行" in m.content and "terminal" in m.content for m in session.messages)
    assert "PASS: pytest tests/test_app.py -q" in response.content


@pytest.mark.asyncio
async def test_coding_ledger_persists_across_continue_turns():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = [
        {"content": "继续后先验证。", "__tool_calls__": True, "tool_calls": [{"id": "test1", "function": {"name": "terminal", "arguments": '{"command":"pytest tests/test_app.py -q"}'}}]},
        {"content": "再编译。", "__tool_calls__": True, "tool_calls": [{"id": "build1", "function": {"name": "terminal", "arguments": '{"command":"python -m py_compile app.py"}'}}]},
        {"content": "已完成修改。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "test1", "name": "terminal", "content": "Command: pytest tests/test_app.py -q\nExit code: 0", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "build1", "name": "terminal", "content": "Command: python -m py_compile app.py\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-continue-ledger"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {
        "coding_changed_files": ["app.py"],
        "coding_validation_results": [],
        "coding_build_results": [],
        "coding_task_status": {
            "kind": "coding_task_status",
            "task_text": "请实现功能并修改代码",
            "tasks": [
                {"id": "understand", "title": "理解需求与约束", "status": "completed"},
                {"id": "locate", "title": "定位相关代码", "status": "completed"},
                {"id": "patch", "title": "完成代码修改", "status": "completed"},
                {"id": "validate", "title": "运行最小验证", "status": "pending"},
                {"id": "build", "title": "尝试编译/构建", "status": "pending"},
                {"id": "report", "title": "汇总变更与验证结果", "status": "pending"},
            ],
        },
    }
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-continue-ledger", channel="feishu", channel_user_id="u1", session_id="s-continue-ledger",
        type=MessageType.TEXT, role=MessageRole.USER, content="继续",
    )

    response = await agent.process_message(user_msg)

    assert not any("Patch-first quality gate failed" in m.content for m in session.messages)
    assert tools.execute_tool_calls.call_count == 2
    assert session.metadata["coding_changed_files"] == ["app.py"]
    assert session.metadata["coding_validation_results"]
    assert session.metadata["coding_build_results"]
    assert "验证结果：PASS: pytest tests/test_app.py -q; PASS: python -m py_compile app.py" in response.content


@pytest.mark.asyncio
async def test_unrelated_message_does_not_resume_stale_coding_ledger():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.return_value = {"content": "按你的新问题回复。", "__tool_calls__": False}

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-stale-ledger"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {
        "coding_changed_files": ["old.py"],
        "coding_task_status": {
            "kind": "coding_task_status",
            "task_text": "请修改旧代码",
            "tasks": [{"id": "patch", "title": "旧代码修改", "status": "pending"}],
        },
    }
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m-new-topic", channel="feishu", channel_user_id="u1", session_id="s-stale-ledger",
        type=MessageType.TEXT, role=MessageRole.USER, content="华为手机微信能加密码吗",
    )

    response = await agent.process_message(user_msg)

    assert response.content == "按你的新问题回复。"
    assert tools.execute_tool_calls.call_count == 0
    first_call_messages = model.chat.call_args_list[0][1]["messages"]
    assert not any("旧代码修改" in str(m.get("content", "")) for m in first_call_messages)


@pytest.mark.asyncio
async def test_unverified_coding_final_is_downgraded_after_changed_files():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = [
        {
            "content": "修改文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "edit1",
                "function": {"name": "edit_file", "arguments": '{"path":"app.py","old":"a","new":"b"}'},
            }],
        },
        {"content": "全量开发完成，全部完成。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {"role": "tool", "tool_call_id": "edit1", "name": "edit_file", "content": "File edited: app.py\n--- diff", "success": True, "metadata": {}},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-unverified-final"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=2)
    user_msg = Message(
        id="m-unverified-final", channel="feishu", channel_user_id="u1", session_id="s-unverified-final",
        type=MessageType.TEXT, role=MessageRole.USER, content="请实现功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert "全量开发完成" not in response.content
    assert "全部完成" not in response.content
    assert "验证结果：未运行" in response.content
    assert "不能视为完整验证通过的交付" in response.content
    assert "[ ] 运行最小验证" in response.content
    assert "[ ] 尝试编译/构建" in response.content

@pytest.mark.asyncio
async def test_agent_uses_larger_read_file_repeat_budget_for_coding_tasks():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    repeated_read = {
        "content": "继续读取代码...",
        "__tool_calls__": True,
        "tool_calls": [{
            "id": "read-code",
            "function": {"name": "read_file", "arguments": '{"path":"MainActivity.java"}'},
        }],
    }
    model.chat.side_effect = [repeated_read] * 10 + [
        {"content": "已完成实现。", "__tool_calls__": False}
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "read-code",
            "name": "read_file",
            "content": "code chunk",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-coding-repeat-budget"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=20)
    user_msg = Message(
        id="m-coding-repeat-budget", channel="feishu", channel_user_id="u1", session_id="s-coding-repeat-budget",
        type=MessageType.TEXT, role=MessageRole.USER, content="请分析这个代码问题并读取代码",
    )

    response = await agent.process_message(user_msg)

    # Exact duplicate calls still trigger the existing reflection guard every other turn.
    # The important regression check is that the read_file name-budget no longer stops
    # an interactive coding task after the historical limit of 8 total read_file calls.
    assert tools.execute_tool_calls.call_count == 5
    assert response.content == "已完成实现。"
    assert not any("只读/查询类工具重复调用过多" in m.content for m in session.messages)
    assert not any("工具调用次数已达到上限" in m.content for m in session.messages)
    assert not any("Patch-first quality gate failed" in m.content for m in session.messages)


@pytest.mark.asyncio
async def test_patch_first_repair_points_to_full_code_navigation_toolchain():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = [
        {"content": "方案：需要修改代码。", "__tool_calls__": False},
        {"content": "这次没有完成代码修改：当前没有产生任何文件 diff。", "__tool_calls__": False},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-code-nav-repair"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m-code-nav-repair", channel="feishu", channel_user_id="u1", session_id="s-code-nav-repair",
        type=MessageType.TEXT, role=MessageRole.USER, content="给我实现：把代码导航从 read_file 迁到 grep_code/read_lines/list_symbols/find_refs/goto_def",
    )

    await agent.process_message(user_msg)

    repair_messages = [m.content for m in session.messages if "Patch-first quality gate failed" in m.content]
    assert repair_messages
    assert "grep_code/read_lines/list_symbols/find_refs/goto_def" in repair_messages[-1]


@pytest.mark.asyncio
async def test_coding_turn_gets_larger_tool_budget_than_regular_chat():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "继续定位代码...",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": f"read-{i}",
                "function": {"name": "read_lines", "arguments": f'{{"path":"big.py","start_line":{i},"end_line":{i+1}}}'},
            } for i in range(21)],
        },
        {"content": "需要我继续完成代码修改吗？", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": f"read-{i}",
            "name": "read_lines",
            "content": "code",
            "success": True,
            "metadata": {},
        } for i in range(21)
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-coding-tool-budget"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=6)
    user_msg = Message(
        id="m-coding-tool-budget", channel="feishu", channel_user_id="u1", session_id="s-coding-tool-budget",
        type=MessageType.TEXT, role=MessageRole.USER, content="请帮我实现这个功能并修改代码",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert any("Patch-first quality gate failed" in m.content for m in session.messages)
    assert not any("工具调用次数已达到上限" in m.content for m in session.messages)
    assert "请确认" not in response.content
    assert "需要我继续" not in response.content
    assert "没有完成代码修改" in response.content


@pytest.mark.asyncio
async def test_unfinished_implementation_confirmation_is_rewritten():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []
    model.chat.side_effect = [
        {"content": "方案：需要修改文件。", "__tool_calls__": False},
        {"content": "需要我继续完成这两处文件的代码修改吗？如果确认，我会继续。", "__tool_calls__": False},
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-confirm-rewrite"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.messages = []
    session.metadata = {}
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m-confirm-rewrite", channel="feishu", channel_user_id="u1", session_id="s-confirm-rewrite",
        type=MessageType.TEXT, role=MessageRole.USER, content="你尝试帮我修改完成这个代码功能",
    )

    response = await agent.process_message(user_msg)

    assert "没有完成代码修改" in response.content
    assert "需要我继续" not in response.content
    assert "如果确认" not in response.content

@pytest.mark.asyncio
async def test_short_write_confirmation_resumes_pending_file_action():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "写入文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "write1",
                "function": {
                    "name": "write_file",
                    "arguments": '{"path":"~/.pyclaw/skills/mac-lock-unlock/lock.sh","content":"#!/bin/bash\\necho lock\\n"}',
                },
            }],
        },
        {
            "content": "运行语法检查。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "check1",
                "function": {"name": "terminal", "arguments": '{"command":"bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh"}'},
            }],
        },
        {"content": "已写入。", "__tool_calls__": False},
        {"content": "脚本无独立构建步骤，已完成。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "write1", "name": "write_file", "content": "File written: ~/.pyclaw/skills/mac-lock-unlock/lock.sh", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "check1", "name": "terminal", "content": "Command: bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-short-write-confirmation"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = [
        Message(
            id="prev-assistant", channel="feishu", channel_user_id="u1", session_id="s-short-write-confirmation",
            type=MessageType.TEXT, role=MessageRole.ASSISTANT,
            content="新版 lock.sh 已定稿。你回复『写入』，我就直接写入 ~/.pyclaw/skills/mac-lock-unlock/lock.sh 并验证。",
        )
    ]
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=8)
    user_msg = Message(
        id="m-write-confirm", channel="feishu", channel_user_id="u1", session_id="s-short-write-confirmation",
        type=MessageType.TEXT, role=MessageRole.USER, content="写入",
    )

    response = await agent.process_message(user_msg)

    first_call_messages = model.chat.call_args_list[0][1]["messages"]
    assert any("PENDING_ACTION_CONTEXT" in str(m.get("content", "")) for m in first_call_messages)
    assert tools.execute_tool_calls.call_count == 2
    assert any('"name": "write_file"' in call.args[0] for call in tools.execute_tool_calls.call_args_list)
    assert "当前没有完成代码修改" not in response.content
    assert "验证结果：PASS: bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh" in response.content


@pytest.mark.asyncio
async def test_terminal_approval_failures_nudge_to_file_tools_before_max_depth():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {"content": "先备份。", "__tool_calls__": True, "tool_calls": [{"id": "cp1", "function": {"name": "terminal", "arguments": '{"command":"cp ~/.pyclaw/skills/mac-lock-unlock/lock.sh ~/.pyclaw/skills/mac-lock-unlock/lock.sh.bak 2>/dev/null; ls ~/.pyclaw/skills/mac-lock-unlock/"}'}}]},
        {"content": "换个写法备份。", "__tool_calls__": True, "tool_calls": [{"id": "cp2", "function": {"name": "terminal", "arguments": '{"command":"cp ~/.pyclaw/skills/mac-lock-unlock/lock.sh ~/.pyclaw/skills/mac-lock-unlock/lock.sh.bak 2>/dev/null; echo done"}'}}]},
        {"content": "改用文件工具。", "__tool_calls__": True, "tool_calls": [{"id": "write1", "function": {"name": "write_file", "arguments": '{"path":"~/.pyclaw/skills/mac-lock-unlock/lock.sh","content":"#!/bin/bash\\necho lock\\n"}'}}]},
        {"content": "验证。", "__tool_calls__": True, "tool_calls": [{"id": "check1", "function": {"name": "terminal", "arguments": '{"command":"bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh"}'}}]},
        {"content": "已完成。", "__tool_calls__": False},
        {"content": "脚本无独立构建步骤，已完成。", "__tool_calls__": False},
    ]
    blocked = "⚠️ 检测到有副作用的指令: `cp ~/.pyclaw/skills/mac-lock-unlock/lock.sh ~/.pyclaw/skills/mac-lock-unlock/lock.sh.bak`\n请添加 approved=True"
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "cp1", "name": "terminal", "content": blocked, "success": False, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "cp2", "name": "terminal", "content": blocked, "success": False, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "write1", "name": "write_file", "content": "File written: ~/.pyclaw/skills/mac-lock-unlock/lock.sh", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "check1", "name": "terminal", "content": "Command: bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-approval-pivot"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = [
        Message(
            id="prev-assistant", channel="feishu", channel_user_id="u1", session_id="s-terminal-approval-pivot",
            type=MessageType.TEXT, role=MessageRole.ASSISTANT,
            content="新版 lock.sh 已定稿。回复写入我就直接写入 ~/.pyclaw/skills/mac-lock-unlock/lock.sh。",
        )
    ]
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-write", channel="feishu", channel_user_id="u1", session_id="s-terminal-approval-pivot",
        type=MessageType.TEXT, role=MessageRole.USER, content="写入",
    )

    response = await agent.process_message(user_msg)

    assert any("repeatedly failed terminal mutations" in m.content for m in session.messages)
    assert tools.execute_tool_calls.call_count == 4
    assert "达到最大思考深度" not in response.content
    assert "当前没有完成代码修改" not in response.content
    assert "验证结果：PASS: bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh" in response.content


@pytest.mark.asyncio
async def test_repeated_write_file_same_target_is_stopped_after_success_and_keeps_validation_failure():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    write_args = json.dumps({
        "path": "~/.pyclaw/skills/mac-lock-unlock/lock.sh",
        "content": "#!/bin/bash\necho lock\n",
    })
    model.chat.side_effect = [
        {"content": "写入。", "__tool_calls__": True, "tool_calls": [{"id": "write1", "function": {"name": "write_file", "arguments": write_args}}]},
        {"content": "又写一次。", "__tool_calls__": True, "tool_calls": [{"id": "write2", "function": {"name": "write_file", "arguments": write_args}}]},
        {"content": "已写入，等待继续。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.return_value = [
        {
            "role": "tool",
            "tool_call_id": "write1",
            "name": "write_file",
            "content": "File written: /Users/bytedance/.pyclaw/skills/mac-lock-unlock/lock.sh",
            "success": True,
            "metadata": {},
        }
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-repeat-write-same-target"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m-repeat-write",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-repeat-write-same-target",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="继续，直接覆盖 lock.sh",
    )

    response = await agent.process_message(user_msg)

    assert tools.execute_tool_calls.call_count == 1
    assert any("本轮只有重复的副作用工具调用" in m.content and "write_file:" in m.content for m in session.messages)
    assert "File written" not in response.content
    assert "验证结果：未运行" in response.content


@pytest.mark.asyncio
async def test_failed_script_execution_is_reported_as_failed_validation_not_unrun():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "写入脚本。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "write1",
                "function": {"name": "write_file", "arguments": '{"path":"~/.pyclaw/skills/mac-lock-unlock/lock.sh","content":"#!/bin/bash\\necho lock\\n"}'},
            }],
        },
        {
            "content": "运行脚本验证。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "run1",
                "function": {"name": "terminal", "arguments": '{"command":"cd ~/.pyclaw/skills/mac-lock-unlock && bash lock.sh 2>&1"}'},
            }],
        },
        {"content": "已写入。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "write1", "name": "write_file", "content": "File written: /Users/bytedance/.pyclaw/skills/mac-lock-unlock/lock.sh", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "run1", "name": "terminal", "content": "Command: cd ~/.pyclaw/skills/mac-lock-unlock && bash lock.sh 2>&1\nExit code: 1\nSTDOUT:\nfailed", "success": False, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-failed-script-validation"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    user_msg = Message(
        id="m-failed-script-validation",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-failed-script-validation",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="继续，直接覆盖 lock.sh 并验证",
    )

    response = await agent.process_message(user_msg)

    assert "验证结果：FAIL: cd ~/.pyclaw/skills/mac-lock-unlock && bash lock.sh 2>&1" in response.content
    assert "验证结果：未运行" not in response.content
    assert "[!] 运行最小验证" in response.content


def test_session_latest_user_message_ignores_internal_notices():
    session = Session(session_id="s-latest-real-user", user_id="u1", channel="feishu")
    real = Message(
        id="u-real-latest",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-latest-real-user",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )
    session.messages = [
        real,
        Message(
            id="internal-latest",
            channel="feishu",
            channel_user_id="u1",
            session_id="s-latest-real-user",
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="NOTICE: Tool usage must stop now.",
            metadata={"internal_notice": True},
        ),
    ]

    assert session.get_latest_user_message() == real


@pytest.mark.asyncio
async def test_file_deliverable_contract_metadata_is_persisted_at_loop_start():
    model = AsyncMock()
    model.chat.return_value = {"content": "我没法生成文件。", "__tool_calls__": False}

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-persist-contract-loop"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    class FakeDB:
        def __init__(self):
            self.executed = []
        async def execute(self, *args):
            self.executed.append(args)
        async def commit(self):
            pass

    class FakeDBContext:
        def __init__(self, db):
            self.db = db
        async def __aenter__(self):
            return self.db
        async def __aexit__(self, exc_type, exc, tb):
            return False

    fake_db = FakeDB()
    def db_connect():
        return FakeDBContext(fake_db)

    sessions.db_connect = db_connect

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=3)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = os.path.abspath(os.path.join(os.getcwd(), "tmp_test_artifacts", "persist_contract"))
    user_msg = Message(
        id="m-persist-contract-loop",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-persist-contract-loop",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    await agent.process_message(user_msg)

    assert (
        "current_completion_contract" in session.metadata
        or "completed_completion_contract" in session.metadata
    )
    assert fake_db.executed


@pytest.mark.asyncio
async def test_terminal_created_artifact_satisfies_file_deliverable_contract(tmp_path):
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    output_file = tmp_path / "AI_Agent_Slides.pptx"
    _write_test_pptx(output_file, 5)
    model.chat.side_effect = [
        {
            "content": "生成 PPT 文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "run-build-ppt",
                "function": {
                    "name": "terminal",
                    "arguments": json.dumps({"command": f"mkdir -p {tmp_path} && touch {output_file} && echo {output_file}"}),
                },
            }],
        },
        {"content": "AI Agent 幻灯片已生成。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{
            "role": "tool",
            "tool_call_id": "run-build-ppt",
            "name": "terminal",
            "content": f"Command: touch {output_file}\nExit code: 0\nSTDOUT:\n{output_file}\n",
            "success": True,
            "metadata": {},
        }],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-artifact-evidence"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=4)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-terminal-artifact-evidence",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-terminal-artifact-evidence",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    assert output_file.exists()
    assert response.metadata["pending_files"] == [{
        "file_path": str(output_file),
        "description": "已生成文件：AI_Agent_Slides.pptx",
    }]
    assert "未观察到目标文件" not in response.content

@pytest.mark.asyncio
async def test_file_deliverable_final_fallback_generates_pptx_when_model_never_delivers(tmp_path):
    model = AsyncMock()
    model.chat.side_effect = [
        {"content": "这是 AI Agent 幻灯片大纲。", "__tool_calls__": False},
        {"content": "仍然只有大纲，没有文件。", "__tool_calls__": False},
        {"content": "我没有办法继续生成文件。", "__tool_calls__": False},
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-synthesize-pptx"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-synthesize-pptx",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-synthesize-pptx",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    output = pending_files[0]["file_path"]
    assert output.endswith(".pptx")
    assert output.startswith(str(tmp_path))
    assert os.path.exists(output)
    assert "未观察到目标文件" not in response.content


@pytest.mark.asyncio
async def test_file_deliverable_dependency_checks_do_not_fake_complete(tmp_path):
    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "先确认 PPT 依赖。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "check-pptx",
                "function": {
                    "name": "terminal",
                    "arguments": json.dumps({
                        "command": "mkdir -p ~/.pyclaw/artifacts/ai-agent-ppt && python3 -c \"import pptx; print(pptx.__version__)\""
                    }),
                },
            }],
        },
        {
            "content": "继续确认依赖。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "check-pptx-again",
                "function": {
                    "name": "terminal",
                    "arguments": json.dumps({
                        "command": f"mkdir -p {tmp_path} && python3 -c \"import pptx; print(pptx.__version__)\""
                    }),
                },
            }],
        },
        {"content": "已生成并发送文件：AI_Agent_Slides.pptx", "__tool_calls__": False},
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(side_effect=[
        [{
            "role": "tool",
            "tool_call_id": "check-pptx",
            "name": "terminal",
            "content": "Command: mkdir -p ~/.pyclaw/artifacts/ai-agent-ppt && python3 -c \"import pptx; print(pptx.__version__)\"\nExit code: 0\nSTDOUT:\n1.0.2\n",
            "success": True,
            "metadata": {},
        }],
        [{
            "role": "tool",
            "tool_call_id": "check-pptx-again",
            "name": "terminal",
            "content": f"Command: mkdir -p {tmp_path} && python3 -c \"import pptx; print(pptx.__version__)\"\nExit code: 0\nSTDOUT:\n1.0.2\n",
            "success": True,
            "metadata": {},
        }],
    ])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-dependency-check-no-fake-complete"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=4)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-dependency-check-no-fake-complete",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-dependency-check-no-fake-complete",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    assert response.metadata.get("pending_files", []) == []
    assert "已生成并发送文件" not in response.content
    assert "未观察到目标文件" in response.content
    assert any(
        "no generated deliverable file has been observed" in str(msg.content)
        for msg in session.messages
        if getattr(msg, "metadata", {}).get("internal_notice")
    )


def test_short_ppt_confirmation_artifact_path_uses_previous_user_request(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path))

    session = Session(session_id="s-ppt-confirm-clean-path", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="ppt-original-request",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    ))
    session.messages.append(Message(
        id="ppt-outline-assistant",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.ASSISTANT,
        content="AI Agent 幻灯片 当前进展汇报 已确认就绪的内容 12 页大纲 可直接落成 PPT 页码。你可以说生成 pptx。",
    ))
    session.messages.append(Message(
        id="ppt-confirm",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="生成 pptx",
    ))

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert "做一个关于_ai_agent_的幻灯片" in contract.artifact_dir.lower()
    assert "当前进展汇报" not in contract.artifact_dir
    assert "12_页大纲" not in contract.artifact_dir


def test_completion_contract_constraints_ignore_assistant_promised_slide_count(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path))

    session = Session(session_id="s-ignore-assistant-count", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="ppt-user-request-no-count",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    ))
    session.messages.append(Message(
        id="ppt-assistant-promises-count",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.ASSISTANT,
        content="AI Agent 幻灯片当前进展：已准备 12 页大纲，你回复『生成 pptx』我就生成。",
    ))
    session.messages.append(Message(
        id="ppt-user-confirm",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="生成 pptx",
    ))

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.task_text == "做一个关于 AI Agent 的幻灯片"
    spec = agent.artifact_acceptance.infer_spec(contract.task_text)
    assert spec.min_slides is None


@pytest.mark.asyncio
async def test_file_deliverable_success_rewrites_failed_text_when_file_is_accepted(tmp_path):
    artifact = tmp_path / "AI_Agent_5_pages.pptx"
    _write_test_pptx(artifact, 5)

    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "发送生成的 PPT。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send-five-page-ppt",
                "function": {
                    "name": "send_file_to_user",
                    "arguments": json.dumps({"file_path": str(artifact), "description": "AI Agent 幻灯片"}),
                },
            }],
        },
        {
            "content": (
                "关于 AI Agent 的幻灯片这一版没能顺利落地文件，"
                "本轮实际写出的 `.pptx` 只有 5 页，未达到 12 页要求，不发送这份残缺文件。"
            ),
            "__tool_calls__": False,
        },
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(return_value=[{
        "role": "tool",
        "tool_call_id": "send-five-page-ppt",
        "name": "send_file_to_user",
        "content": f"prepared file: {artifact}",
        "success": True,
        "metadata": {"is_file_transfer": True, "file_path": str(artifact), "description": "AI Agent 幻灯片"},
    }])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "send_file_to_user"}]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-rewrite-failed-text"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=3)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-rewrite-failed-text",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-rewrite-failed-text",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    assert response.metadata["pending_files"] == [{"file_path": str(artifact), "description": "AI Agent 幻灯片"}]
    assert response.content == "已生成并发送文件：AI_Agent_5_pages.pptx\n验收：PPTX 可打开，共 5 页。"
    assert "没能顺利落地" not in response.content
    assert "不发送" not in response.content


@pytest.mark.asyncio
async def test_file_deliverable_collects_artifact_from_non_terminal_tool(tmp_path):
    output_file = tmp_path / "deck.pptx"
    _write_test_pptx(output_file, 5)

    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "调用幻灯片工具生成文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "slides-tool",
                "function": {"name": "create_slides", "arguments": json.dumps({"topic": "AI Agent"})},
            }],
        },
        {"content": "幻灯片已生成。", "__tool_calls__": False},
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(return_value=[{
        "role": "tool",
        "tool_call_id": "slides-tool",
        "name": "create_slides",
        "content": f"created file: {output_file}",
        "success": True,
        "metadata": {},
    }])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "create_slides"}]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-non-terminal-artifact"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=4)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-non-terminal-artifact",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-non-terminal-artifact",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    assert response.metadata["pending_files"] == [{
        "file_path": str(output_file),
        "description": "已生成文件：deck.pptx",
    }]
    assert "未观察到目标文件" not in response.content


@pytest.mark.asyncio
async def test_session_persists_internal_notice_metadata_after_reload(tmp_path):
    from pyclaw.core.session import SessionManager

    manager = SessionManager(str(tmp_path / "pyclaw.db"))
    await manager.init_db()
    session = await manager.get_or_create("feishu", "u1")

    real = Message(
        id="persist-real-user",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )
    notice = Message(
        id="persist-internal-notice",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="NOTICE: Completion contract failed.",
        metadata={"internal_notice": True},
    )
    await manager.save_message(session, real)
    await manager.save_message(session, notice)

    reloaded_manager = SessionManager(str(tmp_path / "pyclaw.db"))
    reloaded = await reloaded_manager.get_or_create("feishu", "u1")

    assert reloaded.messages[-1].metadata == {"internal_notice": True}
    assert reloaded.get_latest_user_message().id == "persist-real-user"


def test_artifact_manager_returns_expanded_absolute_task_dir():
    from pyclaw.core.artifacts import ArtifactManager

    manager = ArtifactManager(root="~/.pyclaw/artifacts")
    task_dir = manager.task_dir(session_id="s-test", task_text="做一个关于 AI Agent 的幻灯片")

    assert task_dir.startswith(os.path.expanduser("~/.pyclaw/artifacts"))
    assert not task_dir.startswith("~")
    assert os.path.isabs(task_dir)


def test_polluted_persisted_completion_contract_is_not_reused_for_clean_task(tmp_path):
    model = AsyncMock()
    tools = MagicMock()
    sessions = MagicMock()
    agent = Agent(model, tools, sessions)
    agent.artifacts = MagicMock()
    clean_artifact_dir = str(tmp_path / "clean_ai_agent")
    agent.artifacts.task_dir.return_value = clean_artifact_dir

    session = Session(session_id="s-polluted-contract", user_id="u1", channel="feishu")
    session.metadata["current_completion_contract"] = {
        "kind": "file_deliverable",
        "task_text": "做一个关于 AI Agent 的幻灯片 已停止继续执行 避免重复触发副作用操作 验证结果 未观察到目标文件",
        "artifact_dir": str(tmp_path / "已停止继续执行_避免重复触发副作用操作"),
        "required_evidence": ["file_created", "send_file_to_user"],
        "max_repair_attempts": 2,
    }
    session.messages.append(Message(
        id="real-clean-task",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    ))

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.task_text == "做一个关于 AI Agent 的幻灯片"
    assert contract.artifact_dir == clean_artifact_dir
    assert "避免重复触发" not in contract.artifact_dir
    assert session.metadata["current_completion_contract"]["source_message_id"] == "real-clean-task"
    assert session.metadata["current_completion_contract"]["task_fingerprint"] == agent._task_fingerprint("做一个关于 AI Agent 的幻灯片")


def test_stale_completion_contract_source_id_mismatch_is_not_reused(tmp_path):
    model = AsyncMock()
    tools = MagicMock()
    sessions = MagicMock()
    agent = Agent(model, tools, sessions)
    agent.artifacts = MagicMock()
    rag_artifact_dir = str(tmp_path / "rag_enterprise_kb")
    agent.artifacts.task_dir.return_value = rag_artifact_dir

    session = Session(session_id="s-stale-contract-topic", user_id="u1", channel="feishu")
    session.metadata["current_completion_contract"] = {
        "kind": "file_deliverable",
        "task_text": "做一个关于 AI Agent 的幻灯片",
        "artifact_dir": str(tmp_path / "old_ai_agent"),
        "required_evidence": ["file_created", "send_file_to_user"],
        "max_repair_attempts": 2,
        "source_message_id": "rag-msg",
        "task_fingerprint": agent._task_fingerprint("做一个关于 AI Agent 的幻灯片"),
    }
    session.messages.append(Message(
        id="rag-msg",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 RAG 企业知识库的幻灯片",
    ))

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.task_text == "做一个关于 RAG 企业知识库的幻灯片"
    assert contract.artifact_dir == rag_artifact_dir
    assert session.metadata["current_completion_contract"]["source_message_id"] == "rag-msg"
    assert session.metadata["current_completion_contract"]["task_fingerprint"] == agent._task_fingerprint(
        "做一个关于 RAG 企业知识库的幻灯片"
    )


def test_same_text_new_user_message_creates_fresh_explicit_skill_contract(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = MagicMock()
    old_artifact_dir = str(tmp_path / "old_rag_deck")
    new_artifact_dir = str(tmp_path / "new_rag_deck")
    agent.artifacts.task_dir.return_value = new_artifact_dir

    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    old_contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=old_artifact_dir,
        required_evidence=("file_created", "send_file_to_user"),
        max_repair_attempts=2,
        source_message_id="m-old",
        task_fingerprint=agent._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time() - 3600,
    )
    session = Session(session_id="s-repeat-same-skill-task", user_id="u1", channel="feishu")
    session.metadata["current_completion_contract"] = old_contract.to_metadata()
    session.messages.extend([
        Message(
            id="m-old",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content=task,
        ),
        Message(
            id="a-old",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.ASSISTANT,
            content="已生成并发送文件：旧文件.pptx",
        ),
        Message(
            id="m-new",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content=task,
        ),
    ])

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.source_message_id == "m-new"
    assert contract.task_text == task
    assert contract.artifact_dir == new_artifact_dir
    assert contract.required_skills == ("baoyu-design",)
    assert contract.created_at > old_contract.created_at


def test_completed_completion_contract_is_not_recreated_in_same_turn(tmp_path):
    model = AsyncMock()
    tools = MagicMock()
    sessions = MagicMock()
    agent = Agent(model, tools, sessions)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path / "rag_enterprise_kb")

    session = Session(session_id="s-completed-contract", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="rag-msg",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 RAG 企业知识库的幻灯片",
    ))

    contract = agent._infer_completion_contract(session)
    assert contract is not None
    agent._clear_completion_contract(session, contract)

    assert agent._infer_completion_contract(session) is None
    assert "current_completion_contract" not in session.metadata
    assert session.metadata["completed_completion_contract"]["source_message_id"] == "rag-msg"


@pytest.mark.asyncio
async def test_file_deliverable_uses_observed_artifact_under_trusted_root_even_when_contract_path_polluted(tmp_path):
    artifacts_root = tmp_path / "artifacts"
    polluted_dir = artifacts_root / "已生成并发送文件_做一个关于_ai_agent_的幻灯片_已停止继续执行_避免重复触发副作用操作"
    clean_dir = artifacts_root / "ai_agent_slides"
    clean_dir.mkdir(parents=True)
    observed_file = clean_dir / "AI_Agent_Slides.pptx"
    _write_test_pptx(observed_file, 5)

    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "生成 PPT 文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "make-ppt",
                "function": {"name": "python_interpreter", "arguments": json.dumps({"code": "print('ok')"})},
            }],
        },
        {"content": "已生成并发送文件：AI_Agent_Slides.pptx", "__tool_calls__": False},
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(return_value=[{
        "role": "tool",
        "tool_call_id": "make-ppt",
        "name": "python_interpreter",
        "content": f"OBSERVATION from python_interpreter:\nSTDOUT:\n{observed_file} True 45725\n",
        "success": True,
        "metadata": {},
    }])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "python_interpreter"}]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-observed-artifact"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {
        "current_completion_contract": {
            "kind": "file_deliverable",
            "task_text": "做一个关于 AI Agent 的幻灯片 已停止继续执行 避免重复触发副作用操作 验证结果",
            "artifact_dir": str(polluted_dir),
            "required_evidence": ["file_created", "send_file_to_user"],
            "max_repair_attempts": 2,
        },
        "coding_changed_files": [str(artifacts_root / "ai-agent-ppt" / "gen.py")],
        "coding_validation_results": [],
        "coding_build_results": [],
    }
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=4)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(artifacts_root / "做一个关于_ai_agent_的幻灯片_s_observed")
    agent.artifacts.root_path.return_value = str(artifacts_root)

    user_msg = Message(
        id="m-observed-artifact",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-observed-artifact",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    assert response.metadata["pending_files"] == [{
        "file_path": str(observed_file),
        "description": "已生成文件：AI_Agent_Slides.pptx",
    }]
    assert "未观察到目标文件" not in response.content
    assert "验证结果" not in response.content
    assert "代码已产生修改" not in response.content
    assert "避免重复触发" not in response.content


@pytest.mark.asyncio
async def test_polluted_contract_fallback_generates_clean_absolute_pptx(tmp_path):
    model = AsyncMock()
    model.chat.side_effect = [
        {"content": "这里只是一个幻灯片大纲。", "__tool_calls__": False},
        {"content": "仍然没有生成文件。", "__tool_calls__": False},
        {"content": "工具调用已达到执行时限，稍后重试。", "__tool_calls__": False},
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-clean-fallback"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {
        "current_completion_contract": {
            "kind": "file_deliverable",
            "task_text": "做一个关于 AI Agent 的幻灯片 已停止继续执行 避免重复触发副作用操作",
            "artifact_dir": str(tmp_path / "已停止继续执行_避免重复触发副作用操作"),
            "required_evidence": ["file_created", "send_file_to_user"],
            "max_repair_attempts": 2,
        }
    }
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=3)
    agent.artifacts = MagicMock()
    clean_artifact_dir = str(tmp_path / "做一个关于_ai_agent_的幻灯片_s_clean")
    agent.artifacts.task_dir.return_value = clean_artifact_dir
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-clean-fallback",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-clean-fallback",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    output = pending_files[0]["file_path"]
    assert os.path.isabs(output)
    assert output.startswith(clean_artifact_dir)
    assert output.endswith(".pptx")
    assert os.path.exists(output)
    assert "避免重复触发" not in output
    assert "未观察到目标文件" not in response.content


def test_terminal_one_shot_capture_batch_limited_to_one_call():
    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    session = Session(session_id="s-capture-batch", user_id="u1", channel="telegram")
    tool_calls = [
        {"id": "shot-1", "function": {"name": "terminal", "arguments": json.dumps({"command": "screencapture -x /tmp/a.png"})}},
        {"id": "shot-2", "function": {"name": "terminal", "arguments": json.dumps({"command": "screencapture -x /tmp/b.png"})}},
    ]

    assert agent._terminal_one_shot_batch_limit(session, tool_calls) == 1

@pytest.mark.asyncio
async def test_file_deliverable_success_rewrites_stale_incomplete_final_and_dedupes(tmp_path):
    artifact = tmp_path / "AI_Agent_Slides.pptx"
    _write_test_pptx(artifact, 5)

    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "生成并发送 PPT。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send-ppt",
                "function": {
                    "name": "send_file_to_user",
                    "arguments": json.dumps({"file_path": str(artifact), "description": "AI Agent 幻灯片"}),
                },
            }],
        },
        {
            "content": (
                "# AI Agent 幻灯片 — 当前进展\n\n"
                "## ⏸ 待下一轮落地\n"
                "`.pptx` 文件本轮尚未实际写出。下一次你说一句「继续生成」，我会生成并发送。"
            ),
            "__tool_calls__": False,
        },
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(return_value=[{
        "role": "tool",
        "tool_call_id": "send-ppt",
        "name": "send_file_to_user",
        "content": f"sent file: {artifact}",
        "success": True,
        "metadata": {
            "is_file_transfer": True,
            "file_path": str(artifact),
            "description": "AI Agent 幻灯片已生成，共 12 页。",
        },
    }])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "send_file_to_user"}]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-stale-final"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=4)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-stale-final",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-stale-final",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    assert response.metadata["pending_files"] == [{
        "file_path": str(artifact),
        "description": "AI Agent 幻灯片已生成，共 12 页。",
    }]
    assert response.content == "已生成并发送文件：AI_Agent_Slides.pptx\n验收：PPTX 可打开，共 5 页。"
    assert "尚未实际写出" not in response.content
    assert "待下一轮" not in response.content
    assert "继续生成" not in response.content


def test_dedupe_pending_files_normalizes_path_and_keeps_richer_description(tmp_path):
    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    artifact = tmp_path / "deck.pptx"
    pending_files = [
        {"file_path": str(artifact), "description": "短"},
        {"file_path": os.path.join(str(tmp_path), ".", "deck.pptx"), "description": "更完整的幻灯片文件描述"},
    ]

    agent._dedupe_pending_files(pending_files)

    assert pending_files == [{
        "file_path": str(artifact),
        "description": "更完整的幻灯片文件描述",
    }]



def _write_skill_deck_html(path, topic: str = "RAG 企业知识库", slides: int = 12) -> None:
    sections = []
    titles = [
        topic, "业务背景", "核心概念", "数据接入", "切分与索引", "检索策略",
        "重排与上下文", "生成与引用", "权限安全", "效果评测", "生产架构", "落地路线",
    ]
    bullets = [
        ["面向企业私有知识的检索增强生成", "目标是让答案可追溯、可评测、可治理"],
        ["知识分散在文档、代码、会议和业务系统", "传统搜索难以理解自然语言问题"],
        ["先检索证据，再基于证据生成回答", "把来源引用作为答案质量的一部分"],
        ["连接飞书、PDF、网页、代码仓库和数据库", "同步权限、作者、时间和来源元数据"],
        ["按语义层级切分 chunk", "向量索引结合关键词索引提升召回"],
        ["向量召回处理语义相似", "BM25 覆盖专有名词、编号和精确匹配"],
        ["reranker 提升 TopK 证据质量", "合并相邻片段并控制 token 预算"],
        ["提示词约束模型仅基于证据回答", "关键结论附文档标题和段落引用"],
        ["检索阶段执行用户级 ACL", "敏感访问写入审计日志"],
        ["构建 golden set 评测召回和引用准确率", "线上收集反馈持续改进"],
        ["采集管道负责增量同步", "Query Pipeline 负责改写、检索、重排和生成"],
        ["先选高频场景小闭环", "把引用、权限、评测作为上线门槛"],
    ]
    for i in range(slides):
        title = titles[i % len(titles)]
        item_html = "".join(f"<li>{item}</li>" for item in bullets[i % len(bullets)])
        sections.append(f'<section data-label="{title}"><h1>{title}</h1><ul>{item_html}</ul></section>')
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        '<!doctype html><html><body><deck-stage width="1920" height="1080">'
        + "\n".join(sections)
        + '</deck-stage></body></html>',
        encoding="utf-8",
    )

def _write_test_pptx(path, slides: int) -> None:
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService

    service = ArtifactSynthesisService()
    topics = [
        ("AI Agent 总览", ["AI Agent 核心概念", "从对话走向任务执行", "以完成证据衡量结果"]),
        ("目标理解", ["识别用户意图和约束", "澄清缺失输入", "形成可执行任务契约"]),
        ("规划能力", ["拆解多步任务", "选择工具和数据源", "动态调整执行路径"]),
        ("工具调用", ["连接搜索、文件和业务系统", "用权限边界控制副作用", "记录可审计观察结果"]),
        ("记忆机制", ["保存稳定偏好", "压缩长期上下文", "避免把猜测写成事实"]),
        ("反馈闭环", ["观察工具结果", "发现错误后自动修复", "验证后再交付"]),
        ("安全治理", ["高风险操作需要审批", "敏感信息避免泄露", "输出与事实保持一致"]),
        ("可靠性交付", ["文件可打开", "内容匹配主题", "拒绝空白和占位页面"]),
        ("应用场景", ["代码助手", "研究报告", "办公自动化"]),
        ("工程挑战", ["长任务偏航", "工具失败", "环境差异"]),
        ("落地路径", ["先做小闭环", "定义验收标准", "持续回归测试"]),
        ("总结", ["Agent 是闭环执行系统", "控制器负责证据和质量", "交付必须真实可用"]),
    ]
    service._write_pptx(str(path), [
        topics[(i - 1) % len(topics)]
        for i in range(1, slides + 1)
    ])


def _write_process_report_pptx(path, slides: int = 12) -> None:
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService

    service = ArtifactSynthesisService()
    deck = [
        ("AI Agent", ["根据当前任务自动生成的交付版幻灯片"]),
        ("当前进展", ["12 页大纲已定稿", "本轮实际写出的 .pptx 只有 5 页，未达到 12 页要求，不发送这份残缺文件"]),
        ("下一步建议", ["下轮我会重新生成", "通过验证后再 send_file_to_user 推送"]),
    ]
    while len(deck) < slides:
        index = len(deck) + 1
        deck.append((f"补充要点 {index}", ["围绕主题补充可演示的关键观点", "后续可替换为更具体的数据、案例或图片素材"]))
    service._write_pptx(str(path), deck[:slides])


def _write_wrong_topic_pptx(path, slides: int = 8) -> None:
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService

    service = ArtifactSynthesisService()
    deck = [
        ("Kubernetes 集群治理", ["控制平面组件", "节点调度策略", "容器网络模型"]),
        ("服务发现机制", ["Service 抽象", "Ingress 流量入口", "DNS 解析链路"]),
        ("弹性伸缩策略", ["HPA 指标", "副本数调整", "资源配额管理"]),
        ("运维监控体系", ["日志采集", "指标告警", "容量规划"]),
        ("网络与安全", ["NetworkPolicy 隔离", "证书轮转", "访问控制"]),
        ("存储编排", ["PersistentVolume", "StorageClass", "状态服务迁移"]),
        ("发布治理", ["滚动发布", "灰度发布", "回滚策略"]),
        ("成本优化", ["资源画像", "节点池规划", "闲置资源回收"]),
    ]
    service._write_pptx(str(path), deck[:slides])


def _pptx_text(path) -> str:
    import html
    import re
    import zipfile

    parts = []
    with zipfile.ZipFile(path) as zf:
        for name in sorted(zf.namelist()):
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", name):
                raw = zf.read(name).decode("utf-8", errors="ignore")
                parts.extend(html.unescape(item) for item in re.findall(r"<a:t>(.*?)</a:t>", raw))
    return "\n".join(parts)


def test_artifact_acceptance_rejects_pptx_process_report_and_placeholders(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService

    artifact = tmp_path / "bad_process_report.pptx"
    _write_process_report_pptx(artifact, 12)
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="做一个关于 AI Agent 的幻灯片",
        artifact_dir=str(tmp_path),
    )

    result = ArtifactAcceptanceService().evaluate(contract, [{"file_path": str(artifact)}])

    assert not result.accepted
    assert "过程/失败汇报" in result.summary or "占位" in result.summary


def test_pptx_verifier_rejects_wrong_topic_without_substring_false_positive(tmp_path):
    from pyclaw.core.artifact_verification import ArtifactSpec, ArtifactVerifierRegistry

    artifact = tmp_path / "storage_topic.pptx"
    _write_wrong_topic_pptx(artifact, 8)
    spec = ArtifactSpec(expected_kind="pptx", min_slides=3, topic_keywords=("RAG",))

    result = ArtifactVerifierRegistry().verify(spec, str(artifact))

    assert not result.accepted
    assert "主题不匹配" in result.reason


def test_pptx_verifier_rejects_peripheral_topic_mention(tmp_path):
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService
    from pyclaw.core.artifact_verification import ArtifactSpec, ArtifactVerifierRegistry

    artifact = tmp_path / "agent_with_one_rag_mention.pptx"
    ArtifactSynthesisService()._write_pptx(str(artifact), [
        ("AI Agent 总览", ["从对话走向任务执行", "包含规划、工具和记忆"]),
        ("Agent Loop", ["Plan", "Act", "Observe", "Repair"]),
        ("工具系统", ["搜索、文件、终端和业务 API", "某些场景可以接入 RAG"]),
        ("权限与安全", ["副作用工具审批", "沙箱边界"]),
        ("总结", ["AI Agent 是闭环执行系统", "控制器负责验收"]),
    ])
    spec = ArtifactSpec(expected_kind="pptx", min_slides=3, topic_keywords=("RAG", "企业知识库"))

    result = ArtifactVerifierRegistry().verify(spec, str(artifact))

    assert not result.accepted
    assert "主线" in result.reason or "限定关键词" in result.reason


def test_generic_verifier_accepts_nonempty_unknown_artifact_without_quality_constraints(tmp_path):
    from pyclaw.core.artifact_verification import ArtifactSpec, ArtifactVerifierRegistry

    artifact = tmp_path / "notes.custom"
    artifact.write_text("done", encoding="utf-8")

    result = ArtifactVerifierRegistry().verify(ArtifactSpec(), str(artifact))

    assert result.accepted
    assert result.evidence.exists
    assert result.evidence.kind == "unknown"


def test_artifact_acceptance_service_uses_type_specific_pptx_verifier(tmp_path):
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.completion_contract import CompletionContract

    artifact = tmp_path / "short.pptx"
    _write_test_pptx(artifact, 2)
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="做一个关于 AI Agent 的幻灯片",
        artifact_dir=str(tmp_path),
    )

    result = ArtifactAcceptanceService().evaluate(contract, [{"file_path": str(artifact)}])

    assert not result.accepted
    assert result.evidence[0].kind == "pptx"
    assert "最低页数 3 页" in result.summary


def test_artifact_synthesis_ai_agent_fallback_uses_topic_content_not_failure_draft(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality

    contract = CompletionContract(
        kind="file_deliverable",
        task_text="做一个关于 AI Agent 的幻灯片",
        artifact_dir=str(tmp_path),
    )
    draft = "关于 AI Agent 的 12 页幻灯片这一版没能顺利落地文件。下一步建议：继续生成。"

    synthesized = ArtifactSynthesisService().synthesize(contract, draft=draft, quality=SynthesisQuality.BASIC)

    assert synthesized is not None
    text = _pptx_text(synthesized.file_path)
    assert "AI Agent" in text
    assert "工具" in text
    assert "可靠性" in text
    assert "没能顺利落地" not in text
    assert "当前进展" not in text
    assert "补充要点" not in text
    result = ArtifactAcceptanceService().evaluate(contract, [{"file_path": synthesized.file_path}])
    assert result.accepted


def test_artifact_synthesis_rag_fallback_uses_rag_topic_content(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality

    contract = CompletionContract(
        kind="file_deliverable",
        task_text="做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(tmp_path),
    )
    synthesized = ArtifactSynthesisService().synthesize(
        contract,
        draft="工具没能生成文件，不能把过程汇报当作 PPT。",
        quality=SynthesisQuality.BASIC,
    )

    assert synthesized is not None
    text = _pptx_text(synthesized.file_path)
    assert "RAG" in text
    assert "企业知识库" in text
    assert "向量" in text
    assert "检索" in text
    assert "AI Agent" not in text.splitlines()[0]
    result = ArtifactAcceptanceService().evaluate(contract, [{"file_path": synthesized.file_path}])
    assert result.accepted


def test_rich_html_deliverable_enforces_explicit_stack_sections_and_build_script(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality

    task = """
接着上次做 ChatGPT 图解教学网页的丰富版。
单文件 HTML，写到 ~/.pyclaw/artifacts/chatgpt-teach-rich/index.html，通过 send_file_to_user 交付给我。
10 个板块：时间线、Tokenizer、Transformer、注意力、Pretraining、SFT、RLHF、Sampling、Scaling Law、术语表。
技术栈 TailwindCSS CDN + Chart.js CDN + CSS animation + vanilla JS。
实现策略：先写 build.py，再生成 index.html，再 send_file_to_user。
"""
    contract = CompletionContract(kind="file_deliverable", task_text=task, artifact_dir=str(tmp_path))
    acceptance = ArtifactAcceptanceService()
    spec = acceptance.infer_spec(task)

    bad_html = tmp_path / "bad.html"
    bad_html.write_text("""<!doctype html><html><head><title>图解 ChatGPT</title></head><body>
<section>Tokenizer Transformer Pretraining SFT RLHF Sampling Scaling Law 术语表</section>
</body></html>""", encoding="utf-8")
    bad_result = acceptance.evaluate(contract, [{"file_path": str(bad_html)}])
    assert not bad_result.accepted
    assert "TailwindCSS" in bad_result.summary or "Chart.js" in bad_result.summary or "build.py" in bad_result.summary

    synthesized = ArtifactSynthesisService().synthesize(contract, draft="只是进展说明，尚未落地。", quality=SynthesisQuality.FULL)

    assert synthesized is not None
    html = Path(synthesized.file_path).read_text(encoding="utf-8")
    assert "cdn.tailwindcss.com" in html
    assert "cdn.jsdelivr.net/npm/chart.js" in html
    assert "new Chart" in html
    assert "@keyframes" in html
    assert "addEventListener" in html
    for section in ("时间线", "Tokenizer", "Transformer", "注意力", "Pretraining", "SFT", "RLHF", "Sampling", "Scaling Law", "术语表"):
        assert section in html
    assert (tmp_path / "build.py").is_file()
    result = acceptance.evaluate(contract, [{"file_path": synthesized.file_path}])
    assert result.accepted, result.summary


def test_tailwind_utility_class_is_not_misclassified_as_required_skill_for_html_deliverable(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager

    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "write_file"}, {"name": "send_file_to_user"}]
    agent = Agent(AsyncMock(), tools, MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path))

    target = tmp_path / "chatgpt-teach-rich" / "index.html"
    task = f"""
请帮我生成一份"从零学 ChatGPT / LLM"的富样式单页 HTML 教学文档，交付到 `{target}`，最后用 `send_file_to_user` 发给我。

硬性要求：
单个 HTML 文件，不要拆分、不要多文件、不要 zip。用一次 `write_file` 直接落盘，不要分片写。
必须通过 CDN 引入 TailwindCSS（`https://cdn.tailwindcss.com`），全站样式用 Tailwind utility class 完成，不要自己写大段 CSS。
额外通过 CDN 引入 Chart.js（`https://cdn.jsdelivr.net/npm/chart.js`）用于图表。
页面顶部要有固定侧边导航或顶部锚点导航，可跳转到下面 10 个板块。
必须包含以下 10 个板块，顺序如下，每个板块都要有小标题 + 讲解正文 + 至少 1 个可视化元素（图 / 表 / 卡片 / 时间线 / 代码块之一）：
1. 我们要解决什么问题（语言建模 & 自回归）
2. Tokenization（BPE、词表、token 计费直觉）
3. Embedding & 位置编码
4. Transformer 结构（Attention / MHA / FFN / 残差 / LayerNorm）
5. 训练目标：Next Token Prediction + 交叉熵
6. Pretrain：数据、算力、Scaling Law（配一张 loss vs compute 的 Chart.js 图）
7. SFT（指令微调）
8. RLHF / DPO（含奖励模型直觉 + 一张对比柱状图）
9. 推理阶段：temperature / top-k / top-p / repetition penalty（配一张采样分布示意图）
10. 术语表（表格形式，至少 15 个术语）
视觉风格：深色背景 + 强调色（蓝紫渐变），代码块用等宽字体和深色卡片样式，图表配色和整体风格一致。
交互：导航点击平滑滚动；术语表支持顶部输入框实时过滤。
页面自包含，双击 HTML 可直接在浏览器打开，不依赖本地文件。
请直接一次性生成完整文件并调用 `send_file_to_user` 交付，不要先给我大纲让我确认。
"""
    session = Session(session_id="s-tailwind-stack-not-skill", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-tailwind-stack-not-skill",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    assert agent._required_skills_for_task(task, session) == ()
    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.kind == "file_deliverable"
    assert contract.required_skills == ()
    assert contract.artifact_dir == str(target.parent)

    pending_files: list[dict[str, str]] = []
    final = agent._prepare_completion_contract_final_content(
        session=session,
        content="已生成的文件未通过 skill 工作流验收，任务未完成：未观察到已激活的必需 skill 上下文。",
        pending_files=pending_files,
    )

    assert pending_files == [{"file_path": str(target), "description": "已生成网页文件：index.html"}]
    assert "任务未完成" not in final
    assert target.is_file()
    html = target.read_text(encoding="utf-8")
    assert "cdn.tailwindcss.com" in html
    assert "cdn.jsdelivr.net/npm/chart.js" in html
    assert "new Chart" in html
    for section in ("Tokenization", "Embedding", "Transformer", "Next Token Prediction", "Pretrain", "SFT", "RLHF", "temperature", "术语表"):
        assert section in html
    result = agent.artifact_acceptance.evaluate(contract, pending_files)
    assert result.accepted, result.summary


def test_rich_html_deliverable_preserves_unnumbered_explicit_sections_in_order(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality

    task = f"""
请帮我生成一份"从零学 ChatGPT / LLM"的富样式单页 HTML 教学文档，交付到 `{tmp_path / 'chatgpt-teach-rich' / 'index.html'}`，最后用 `send_file_to_user` 发给我。

硬性要求：
单个 HTML 文件，不要拆分、不要多文件、不要 zip。用一次 `write_file` 直接落盘，不要分片写。
必须通过 CDN 引入 TailwindCSS（`https://cdn.tailwindcss.com`），全站样式用 Tailwind utility class 完成，不要自己写大段 CSS。
额外通过 CDN 引入 Chart.js（`https://cdn.jsdelivr.net/npm/chart.js`）用于图表。
页面顶部要有固定侧边导航或顶部锚点导航，可跳转到下面 10 个板块。
必须包含以下 10 个板块，顺序如下，每个板块都要有小标题 + 讲解正文 + 至少 1 个可视化元素（图 / 表 / 卡片 / 时间线 / 代码块之一）：
我们要解决什么问题（语言建模 & 自回归）
Tokenization（BPE、词表、token 计费直觉）
Embedding & 位置编码
Transformer 结构（Attention / MHA / FFN / 残差 / LayerNorm）
训练目标：Next Token Prediction + 交叉熵
Pretrain：数据、算力、Scaling Law（配一张 loss vs compute 的 Chart.js 图）
SFT（指令微调）
RLHF / DPO（含奖励模型直觉 + 一张对比柱状图）
推理阶段：temperature / top-k / top-p / repetition penalty（配一张采样分布示意图）
术语表（表格形式，至少 15 个术语）
视觉风格：深色背景 + 强调色（蓝紫渐变），代码块用等宽字体和深色卡片样式，图表配色和整体风格一致。
交互：导航点击平滑滚动；术语表支持顶部输入框实时过滤。
页面自包含，双击 HTML 可直接在浏览器打开，不依赖本地文件。
请直接一次性生成完整文件并调用 `send_file_to_user` 交付，不要先给我大纲让我确认。
"""
    required_sections = (
        "我们要解决什么问题", "Tokenization", "Embedding", "Transformer", "训练目标",
        "Pretrain", "SFT", "RLHF / DPO", "推理阶段", "术语表",
    )
    acceptance = ArtifactAcceptanceService()
    spec = acceptance.infer_spec(task)
    assert len(spec.required_html_sections) == 10
    assert spec.required_html_sections[0].startswith("我们要解决什么问题")

    generic_html = tmp_path / "generic.html"
    generic_html.write_text(
        '<!doctype html><html><head><script src="https://cdn.tailwindcss.com"></script>'
        '<script src="https://cdn.jsdelivr.net/npm/chart.js"></script></head><body>'
        '<section>Tokenization BPE Embedding Transformer Attention MHA FFN LayerNorm Next Token Prediction Pretrain SFT RLHF DPO temperature top-k top-p repetition penalty 术语表</section>'
        '<script>new Chart(document.createElement("canvas"), {type:"bar"}); document.body.addEventListener("click",()=>{});</script>'
        '</body></html>',
        encoding="utf-8",
    )
    contract = CompletionContract(kind="file_deliverable", task_text=task, artifact_dir=str(tmp_path))
    rejected = acceptance.evaluate(contract, [{"file_path": str(generic_html)}])
    assert not rejected.accepted
    assert rejected.summary

    synthesized = ArtifactSynthesisService().synthesize(contract, draft="进度说明，不是最终网页。", quality=SynthesisQuality.FULL)
    assert synthesized is not None
    html = Path(synthesized.file_path).read_text(encoding="utf-8")
    assert html.count('id="sec-') >= 10
    assert html.count("glossary-row") >= 15
    assert "cdn.tailwindcss.com" in html
    assert "cdn.jsdelivr.net/npm/chart.js" in html
    assert "new Chart" in html
    assert "<style" not in html.lower()
    cursor = -1
    for section in required_sections:
        position = html.find(section)
        assert position > cursor, section
        cursor = position
    for marker in ("语言建模", "自回归", "交叉熵", "loss vs compute", "repetition penalty"):
        assert marker in html
    accepted = acceptance.evaluate(contract, [{"file_path": synthesized.file_path}])
    assert accepted.accepted, accepted.summary


def test_rich_html_deliverable_contains_substantive_teaching_content(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality

    task = """
请帮我生成一份"从零学 ChatGPT / LLM"的富样式单页 HTML 教学文档，交付到 `~/.pyclaw/artifacts/chatgpt-teach-rich/index.html`，最后用 `send_file_to_user` 发给我。
必须通过 CDN 引入 TailwindCSS（`https://cdn.tailwindcss.com`），全站样式用 Tailwind utility class 完成，不要自己写大段 CSS。
额外通过 CDN 引入 Chart.js（`https://cdn.jsdelivr.net/npm/chart.js`）用于图表。
必须包含以下 10 个板块，顺序如下，每个板块都要有小标题 + 讲解正文 + 至少 1 个可视化元素：
我们要解决什么问题（语言建模 & 自回归）
Tokenization（BPE、词表、token 计费直觉）
Embedding & 位置编码
Transformer 结构（Attention / MHA / FFN / 残差 / LayerNorm）
训练目标：Next Token Prediction + 交叉熵
Pretrain：数据、算力、Scaling Law（配一张 loss vs compute 的 Chart.js 图）
SFT（指令微调）
RLHF / DPO（含奖励模型直觉 + 一张对比柱状图）
推理阶段：temperature / top-k / top-p / repetition penalty（配一张采样分布示意图）
术语表（表格形式，至少 15 个术语）
交互：导航点击平滑滚动；术语表支持顶部输入框实时过滤。
"""
    contract = CompletionContract(kind="file_deliverable", task_text=task, artifact_dir=str(tmp_path))
    synthesized = ArtifactSynthesisService().synthesize(contract, draft="进度报告", quality=SynthesisQuality.FULL)
    assert synthesized is not None
    html = Path(synthesized.file_path).read_text(encoding="utf-8")
    section_texts = []
    for body in re.findall(r"<article\b[^>]*id=\"sec-[^\"]+\"[^>]*>(.*?)</article>", html, flags=re.IGNORECASE | re.DOTALL):
        heading_match = re.search(r"<h2[^>]*>(.*?)</h2>", body, flags=re.IGNORECASE | re.DOTALL)
        paragraph_match = re.search(r"<p[^>]*leading-8[^>]*>(.*?)</p>", body, flags=re.IGNORECASE | re.DOTALL)
        assert heading_match is not None
        assert paragraph_match is not None
        heading_text = re.sub(r"<[^>]+>", " ", heading_match.group(1)).strip()
        paragraph_text = re.sub(r"<[^>]+>", " ", paragraph_match.group(1)).strip()
        assert len(paragraph_text) >= 45, heading_text
        assert paragraph_text != heading_text
        section_texts.append(paragraph_text)
    assert len(section_texts) >= 10
    joined = "\n".join(section_texts)
    for marker in ("给定前文 token", "条件概率", "BPE", "向量空间", "softmax", "奖励模型", "DPO", "repetition penalty"):
        assert marker in joined
    assert html.count('data-chart="loss"') >= 1
    assert html.count('data-chart="preference"') >= 1
    assert html.count('data-chart="sampling"') >= 1


def test_premium_html_request_infers_visual_polish_requirement():
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService

    task = "做一个精美、高颜值的可视化教学网页，图解 GPT-3 结构和训练流程"
    spec = ArtifactAcceptanceService().infer_spec(task)

    assert spec.expected_kind == "html"
    assert spec.requires_visual_polish is True


def test_premium_html_synthesis_passes_visual_polish_gate(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality

    task = """
做一个精美、高颜值的可视化教学网页，图解 GPT-3 结构和训练流程。
要求有高级设计感、交互式演示、响应式布局和数据可视化，最终生成单文件 HTML。
"""
    contract = CompletionContract(kind="file_deliverable", task_text=task, artifact_dir=str(tmp_path))
    synthesized = ArtifactSynthesisService().synthesize(contract, draft="先做了大纲，还没有落地文件。", quality=SynthesisQuality.FULL)

    assert synthesized is not None
    html_text = Path(synthesized.file_path).read_text(encoding="utf-8")
    assert "cdn.tailwindcss.com" in html_text
    assert "cdn.jsdelivr.net/npm/chart.js" in html_text
    assert "bg-gradient" in html_text or "radial-gradient" in html_text
    assert "backdrop-blur" in html_text
    assert "shadow" in html_text
    assert "rounded" in html_text
    assert "transition" in html_text
    assert "<nav" in html_text
    assert "<canvas" in html_text
    assert "addEventListener" in html_text
    assert len(html_text) >= 14_000

    accepted = ArtifactAcceptanceService().evaluate(contract, [{"file_path": synthesized.file_path}])
    assert accepted.accepted, accepted.summary


def test_premium_html_visual_polish_rejects_structural_page(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService

    task = "做一个精美、高颜值的可视化教学网页，讲解 RAG 企业知识库"
    artifact = tmp_path / "simple.html"
    artifact.write_text(
        "<!doctype html><html><head><title>RAG 企业知识库</title></head>"
        "<body><section><h1>RAG 企业知识库</h1><p>检索增强生成。</p></section></body></html>",
        encoding="utf-8",
    )
    contract = CompletionContract(kind="file_deliverable", task_text=task, artifact_dir=str(tmp_path))

    rejected = ArtifactAcceptanceService().evaluate(contract, [{"file_path": str(artifact)}])

    assert not rejected.accepted
    assert "视觉" in rejected.summary or "精美" in rejected.summary or "设计" in rejected.summary


@pytest.mark.asyncio
async def test_file_deliverable_replaces_process_report_pptx_with_topic_deck(tmp_path):
    bad_artifact = tmp_path / "AI_Agent_process_report.pptx"
    _write_process_report_pptx(bad_artifact, 12)

    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "先发送生成好的 PPT。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send-process-report-ppt",
                "function": {
                    "name": "send_file_to_user",
                    "arguments": json.dumps({"file_path": str(bad_artifact), "description": "AI Agent 幻灯片"}),
                },
            }],
        },
        {"content": "AI Agent 幻灯片已生成并发送。", "__tool_calls__": False},
        {"content": "AI Agent 幻灯片已生成并发送。", "__tool_calls__": False},
        {"content": "AI Agent 幻灯片已生成并发送。", "__tool_calls__": False},
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(return_value=[{
        "role": "tool",
        "tool_call_id": "send-process-report-ppt",
        "name": "send_file_to_user",
        "content": f"sent file: {bad_artifact}",
        "success": True,
        "metadata": {"is_file_transfer": True, "file_path": str(bad_artifact), "description": "AI Agent 幻灯片"},
    }])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "send_file_to_user"}]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-process-report-replacement"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=5)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-process-report-replacement",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-process-report-replacement",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的幻灯片",
    )

    response = await agent.process_message(user_msg)

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    replacement = pending_files[0]["file_path"]
    assert replacement != str(bad_artifact)
    assert replacement.endswith(".pptx")
    assert os.path.exists(replacement)
    text = _pptx_text(replacement)
    assert "AI Agent" in text
    assert "工具" in text
    assert "可靠性" in text
    assert "当前进展" not in text
    assert "没能顺利落地" not in text
    assert "补充要点" not in text
    assert "后续可替换" not in text
    acceptance = agent.artifact_acceptance.evaluate(agent._infer_completion_contract(session), pending_files)
    assert acceptance.accepted


@pytest.mark.asyncio
async def test_file_deliverable_rejects_short_pptx_and_repairs_to_required_slide_count(tmp_path):
    bad_artifact = tmp_path / "AI_Agent_short.pptx"
    good_artifact = tmp_path / "AI_Agent_12_pages.pptx"
    _write_test_pptx(bad_artifact, 5)
    _write_test_pptx(good_artifact, 12)

    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "先发送生成好的 PPT。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send-short-ppt",
                "function": {
                    "name": "send_file_to_user",
                    "arguments": json.dumps({"file_path": str(bad_artifact), "description": "AI Agent 幻灯片"}),
                },
            }],
        },
        {
            "content": "已重新生成 12 页版本。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send-good-ppt",
                "function": {
                    "name": "send_file_to_user",
                    "arguments": json.dumps({"file_path": str(good_artifact), "description": "AI Agent 幻灯片 12 页"}),
                },
            }],
        },
        {"content": "AI Agent 12 页幻灯片已生成并发送。", "__tool_calls__": False},
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(side_effect=[
        [{
            "role": "tool",
            "tool_call_id": "send-short-ppt",
            "name": "send_file_to_user",
            "content": f"sent file: {bad_artifact}",
            "success": True,
            "metadata": {"is_file_transfer": True, "file_path": str(bad_artifact), "description": "AI Agent 幻灯片"},
        }],
        [{
            "role": "tool",
            "tool_call_id": "send-good-ppt",
            "name": "send_file_to_user",
            "content": f"sent file: {good_artifact}",
            "success": True,
            "metadata": {"is_file_transfer": True, "file_path": str(good_artifact), "description": "AI Agent 幻灯片 12 页"},
        }],
    ])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "send_file_to_user"}]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-acceptance-repair"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=6)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-acceptance-repair",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-acceptance-repair",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的 12 页幻灯片",
    )

    response = await agent.process_message(user_msg)

    assert response.metadata["pending_files"] == [{
        "file_path": str(good_artifact),
        "description": "AI Agent 幻灯片 12 页",
    }]
    assert str(bad_artifact) not in json.dumps(response.metadata, ensure_ascii=False)
    assert "12" in response.content
    assert "验收" in response.content or "12 页" in response.content
    assert any(
        "只有 5 页" in str(msg.content) and "12 页" in str(msg.content)
        for msg in session.messages
        if getattr(msg, "metadata", {}).get("internal_notice")
    )


@pytest.mark.asyncio
async def test_file_deliverable_rejected_artifact_is_not_sent_on_forced_stop(tmp_path):
    bad_artifact = tmp_path / "AI_Agent_short.pptx"
    _write_test_pptx(bad_artifact, 5)

    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "发送 PPT。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send-short-ppt-stop",
                "function": {
                    "name": "send_file_to_user",
                    "arguments": json.dumps({"file_path": str(bad_artifact), "description": "AI Agent 幻灯片"}),
                },
            }],
        },
        {"content": "已生成并发送文件：AI_Agent_short.pptx", "__tool_calls__": False},
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(return_value=[{
        "role": "tool",
        "tool_call_id": "send-short-ppt-stop",
        "name": "send_file_to_user",
        "content": f"sent file: {bad_artifact}",
        "success": True,
        "metadata": {"is_file_transfer": True, "file_path": str(bad_artifact), "description": "AI Agent 幻灯片"},
    }])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "send_file_to_user"}]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-acceptance-stop"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=2)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-acceptance-stop",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-acceptance-stop",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的 12 页幻灯片",
    )

    response = await agent.process_message(user_msg)

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    assert pending_files[0]["file_path"] != str(bad_artifact)
    acceptance = agent.artifact_acceptance.evaluate(agent._infer_completion_contract(session), pending_files)
    assert acceptance.accepted
    assert "共 12 页" in acceptance.summary
    assert "未通过交付验收" not in response.content
    assert "验收：PPTX 可打开，共 12 页。" in response.content


@pytest.mark.asyncio
async def test_file_deliverable_synthesizes_replacement_for_rejected_explicit_slide_count(tmp_path):
    bad_artifact = tmp_path / "AI_Agent_short.pptx"
    _write_test_pptx(bad_artifact, 5)

    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "先发送当前 PPT。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "send-short-ppt-replace",
                "function": {
                    "name": "send_file_to_user",
                    "arguments": json.dumps({"file_path": str(bad_artifact), "description": "AI Agent 幻灯片"}),
                },
            }],
        },
        {"content": "AI Agent 12 页幻灯片已生成并发送。", "__tool_calls__": False},
    ]

    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock(return_value=[{
        "role": "tool",
        "tool_call_id": "send-short-ppt-replace",
        "name": "send_file_to_user",
        "content": f"sent file: {bad_artifact}",
        "success": True,
        "metadata": {"is_file_transfer": True, "file_path": str(bad_artifact), "description": "AI Agent 幻灯片"},
    }])
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "send_file_to_user"}]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-acceptance-replacement"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = []
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=3)
    agent.artifacts = MagicMock()
    agent.artifacts.task_dir.return_value = str(tmp_path)
    agent.artifacts.root = str(tmp_path)
    agent.artifacts.root_path.return_value = str(tmp_path)

    user_msg = Message(
        id="m-acceptance-replacement",
        channel="feishu",
        channel_user_id="u1",
        session_id="s-acceptance-replacement",
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的 12 页幻灯片",
    )

    response = await agent.process_message(user_msg)

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    replacement = pending_files[0]["file_path"]
    assert replacement != str(bad_artifact)
    assert replacement.endswith(".pptx")
    acceptance = agent.artifact_acceptance.evaluate(agent._infer_completion_contract(session), pending_files)
    assert acceptance.accepted
    assert "共 12 页" in acceptance.summary
    assert "未通过交付验收" not in response.content
    assert "验收：PPTX 可打开，共 12 页。" in response.content


def test_deliverable_workflow_rejects_wrong_topic_and_synthesizes_real_topic(tmp_path):
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.deliverable_workflow import DeliverableWorkflow

    wrong = tmp_path / "wrong_topic.pptx"
    _write_wrong_topic_pptx(wrong, 8)
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(tmp_path),
    )
    workflow = DeliverableWorkflow(
        acceptance=ArtifactAcceptanceService(),
        synthesis=ArtifactSynthesisService(),
    )
    pending_files = [{"file_path": str(wrong), "description": "RAG 幻灯片"}]

    evidence = workflow.evidence(contract=contract, pending_files=pending_files)
    assert evidence.artifact_acceptance is not None
    assert not evidence.artifact_acceptance.accepted
    assert evidence.artifact_acceptance.summary

    finalization = workflow.finalize(
        contract=contract,
        content="",
        pending_files=pending_files,
        synthesis_quality=SynthesisQuality.BASIC,
        force_repair_synthesis=True,
    )

    assert len(finalization.pending_files) == 1
    replacement = finalization.pending_files[0]["file_path"]
    assert replacement != str(wrong)
    text = _pptx_text(replacement)
    assert "RAG" in text
    assert "企业知识库" in text
    assert "Kubernetes" not in text
    assert finalization.acceptance is not None
    assert finalization.acceptance.accepted

@pytest.mark.asyncio
async def test_short_write_confirmation_resumes_pending_file_action():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {
            "content": "写入文件。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "write1",
                "function": {
                    "name": "write_file",
                    "arguments": '{"path":"~/.pyclaw/skills/mac-lock-unlock/lock.sh","content":"#!/bin/bash\\necho lock\\n"}',
                },
            }],
        },
        {
            "content": "运行语法检查。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "check1",
                "function": {"name": "terminal", "arguments": '{"command":"bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh"}'},
            }],
        },
        {"content": "已写入。", "__tool_calls__": False},
    ]
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "write1", "name": "write_file", "content": "File written: ~/.pyclaw/skills/mac-lock-unlock/lock.sh", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "check1", "name": "terminal", "content": "Command: bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-short-write-confirmation"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = [
        Message(
            id="prev-assistant", channel="feishu", channel_user_id="u1", session_id="s-short-write-confirmation",
            type=MessageType.TEXT, role=MessageRole.ASSISTANT,
            content="新版 lock.sh 已定稿。你回复『写入』，我就直接写入 ~/.pyclaw/skills/mac-lock-unlock/lock.sh 并验证。",
        )
    ]
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=8)
    user_msg = Message(
        id="m-write-confirm", channel="feishu", channel_user_id="u1", session_id="s-short-write-confirmation",
        type=MessageType.TEXT, role=MessageRole.USER, content="写入",
    )

    response = await agent.process_message(user_msg)

    first_call_messages = model.chat.call_args_list[0][1]["messages"]
    assert any("PENDING_ACTION_CONTEXT" in str(m.get("content", "")) for m in first_call_messages)
    assert tools.execute_tool_calls.call_count == 2
    assert any('"name": "write_file"' in call.args[0] for call in tools.execute_tool_calls.call_args_list)
    assert "当前没有完成代码修改" not in response.content
    assert "验证结果：PASS: bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh" in response.content


@pytest.mark.asyncio
async def test_terminal_approval_failures_nudge_to_file_tools_before_max_depth():
    model = AsyncMock()
    tools = MagicMock()
    tools.execute_tool_calls = AsyncMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = []

    model.chat.side_effect = [
        {"content": "先备份。", "__tool_calls__": True, "tool_calls": [{"id": "cp1", "function": {"name": "terminal", "arguments": '{"command":"cp ~/.pyclaw/skills/mac-lock-unlock/lock.sh ~/.pyclaw/skills/mac-lock-unlock/lock.sh.bak 2>/dev/null; ls ~/.pyclaw/skills/mac-lock-unlock/"}'}}]},
        {"content": "换个写法备份。", "__tool_calls__": True, "tool_calls": [{"id": "cp2", "function": {"name": "terminal", "arguments": '{"command":"cp ~/.pyclaw/skills/mac-lock-unlock/lock.sh ~/.pyclaw/skills/mac-lock-unlock/lock.sh.bak 2>/dev/null; echo done"}'}}]},
        {"content": "改用文件工具。", "__tool_calls__": True, "tool_calls": [{"id": "write1", "function": {"name": "write_file", "arguments": '{"path":"~/.pyclaw/skills/mac-lock-unlock/lock.sh","content":"#!/bin/bash\\necho lock\\n"}'}}]},
        {"content": "验证。", "__tool_calls__": True, "tool_calls": [{"id": "check1", "function": {"name": "terminal", "arguments": '{"command":"bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh"}'}}]},
        {"content": "已完成。", "__tool_calls__": False},
    ]
    blocked = "⚠️ 检测到有副作用的指令: `cp ~/.pyclaw/skills/mac-lock-unlock/lock.sh ~/.pyclaw/skills/mac-lock-unlock/lock.sh.bak`\n请添加 approved=True"
    tools.execute_tool_calls.side_effect = [
        [{"role": "tool", "tool_call_id": "cp1", "name": "terminal", "content": blocked, "success": False, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "cp2", "name": "terminal", "content": blocked, "success": False, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "write1", "name": "write_file", "content": "File written: ~/.pyclaw/skills/mac-lock-unlock/lock.sh", "success": True, "metadata": {}}],
        [{"role": "tool", "tool_call_id": "check1", "name": "terminal", "content": "Command: bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh\nExit code: 0", "success": True, "metadata": {}}],
    ]

    sessions = AsyncMock()
    session = MagicMock()
    session.session_id = "s-terminal-approval-pivot"
    session.channel = "feishu"
    session.channel_user_id = "u1"
    session.user_id = "u1"
    session.metadata = {}
    session.messages = [
        Message(
            id="prev-assistant", channel="feishu", channel_user_id="u1", session_id="s-terminal-approval-pivot",
            type=MessageType.TEXT, role=MessageRole.ASSISTANT,
            content="新版 lock.sh 已定稿。回复写入我就直接写入 ~/.pyclaw/skills/mac-lock-unlock/lock.sh。",
        )
    ]
    session.get_history.side_effect = lambda limit=10: [m.to_llm_format() for m in session.messages]

    async def save_msg_side_effect(sess, msg):
        if msg not in sess.messages:
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_msg_side_effect
    sessions.get_or_create.return_value = session

    agent = Agent(model, tools, sessions, max_iterations=10)
    user_msg = Message(
        id="m-write", channel="feishu", channel_user_id="u1", session_id="s-terminal-approval-pivot",
        type=MessageType.TEXT, role=MessageRole.USER, content="写入",
    )

    response = await agent.process_message(user_msg)

    assert any("repeatedly failed terminal mutations" in m.content for m in session.messages)
    assert tools.execute_tool_calls.call_count == 4
    assert "达到最大思考深度" not in response.content
    assert "当前没有完成代码修改" not in response.content
    assert "验证结果：PASS: bash -n ~/.pyclaw/skills/mac-lock-unlock/lock.sh" in response.content


def _make_baoyu_skill_fixture(skill_dir):
    (skill_dir / "built-in-skills").mkdir(parents=True)
    (skill_dir / "references").mkdir()
    (skill_dir / "SKILL.md").write_text(
        "---\nname: baoyu-design\ndescription: Rich design workflow.\n---\n\n"
        "Read [system-prompt.md](system-prompt.md).\n",
        encoding="utf-8",
    )
    (skill_dir / "system-prompt.md").write_text(
        "Read [references/codex.md](references/codex.md).\n"
        "Capabilities: [wireframe](built-in-skills/wireframe.md), "
        "[design system](built-in-skills/design-system-authoring-guide.md), "
        "[create design system](built-in-skills/create-design-system.md), "
        "[components](built-in-skills/design-components.md), "
        "[import figma](built-in-skills/import-from-figma.md), "
        "[import github](built-in-skills/import-from-github.md).\n",
        encoding="utf-8",
    )
    for rel in (
        "references/codex.md",
        "built-in-skills/make-a-deck.md",
        "built-in-skills/export-as-pptx-editable.md",
        "built-in-skills/wireframe.md",
        "built-in-skills/design-system-authoring-guide.md",
        "built-in-skills/create-design-system.md",
        "built-in-skills/design-components.md",
        "built-in-skills/import-from-figma.md",
        "built-in-skills/import-from-github.md",
    ):
        path = skill_dir / rel
        path.parent.mkdir(parents=True, exist_ok=True)
        body = "Use `<deck-stage width=\"1920\" height=\"1080\">` and `section data-label`." if rel.endswith("make-a-deck.md") else rel
        path.write_text(body, encoding="utf-8")


def test_baoyu_deck_skill_evidence_routes_to_deck_docs_only(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.skill_context import ActiveSkillContext
    from pyclaw.core.skill_evidence import SkillEvidenceService

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    ctx = ActiveSkillContext(
        name="baoyu-design",
        canonical_rel_path="baoyu-design/skills/baoyu-design",
        skill_md_path=str(skill_dir / "SKILL.md"),
        root_dir=str(skill_dir),
        content_sha256="sha",
    )
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(tmp_path / "artifact"),
        required_skills=("baoyu-design",),
    )

    requirement = SkillEvidenceService().infer_requirement(ctx, contract)

    assert "system-prompt.md" in requirement.required_paths
    assert "references/*.md" in requirement.required_paths
    assert "built-in-skills/make-a-deck.md" in requirement.required_paths
    assert "built-in-skills/export-as-pptx-editable.md" in requirement.required_paths
    assert "built-in-skills/wireframe.md" not in requirement.required_paths
    assert "built-in-skills/design-system-authoring-guide.md" not in requirement.required_paths
    assert "built-in-skills/create-design-system.md" not in requirement.required_paths
    assert "built-in-skills/design-components.md" not in requirement.required_paths


def test_baoyu_deck_stage_takes_precedence_over_html2pptx_helpers(tmp_path):
    from pyclaw.core.skill_evidence import SkillEvidenceService

    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    # Baoyu-style design bundles can contain PPTX/export helpers.  Their deck
    # authoring contract is still deck-stage, not standalone html2pptx.
    (skill_dir / "html2pptx.md").write_text("# legacy helper\n", encoding="utf-8")
    (skill_dir / "scripts").mkdir()
    (skill_dir / "scripts" / "html2pptx.js").write_text("// helper\n", encoding="utf-8")

    evidence = SkillEvidenceService()

    assert evidence.is_deck_stage_workflow(str(skill_dir))
    assert not evidence.is_html2pptx_workflow(str(skill_dir))


@pytest.mark.asyncio
async def test_baoyu_notice_ignores_unrelated_active_pptx_context(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    baoyu_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(baoyu_dir)
    pptx_dir = tmp_path / "skills" / "pptx"
    (pptx_dir / "scripts").mkdir(parents=True)
    (pptx_dir / "SKILL.md").write_text(
        "---\nname: pptx\n---\n[html2pptx](html2pptx.md)\n",
        encoding="utf-8",
    )
    (pptx_dir / "html2pptx.md").write_text("# html2pptx workflow\n", encoding="utf-8")
    (pptx_dir / "scripts" / "html2pptx.js").write_text("// converter\n", encoding="utf-8")
    session = Session(
        session_id="s-baoyu-notice-filter",
        user_id="u1",
        channel="feishu",
        metadata={"active_skill_contexts": [
            {
                "name": "baoyu-design",
                "canonical_rel_path": "baoyu-design/skills/baoyu-design",
                "skill_md_path": str(baoyu_dir / "SKILL.md"),
                "root_dir": str(tmp_path / "skills"),
            },
            {
                "name": "pptx",
                "canonical_rel_path": "pptx",
                "skill_md_path": str(pptx_dir / "SKILL.md"),
                "root_dir": str(tmp_path / "skills"),
            },
        ]},
    )

    async def save_message(sess, msg):
        sess.messages.append(msg)

    agent.sessions.save_message = AsyncMock(side_effect=save_message)
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(tmp_path / "artifact"),
        source_message_id="m-baoyu",
        task_fingerprint=agent._task_fingerprint("走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"),
        required_skills=("baoyu-design",),
    )

    await agent._ensure_skill_workflow_orchestration_notice(session, contract)

    notice = "\n".join(msg.content for msg in session.messages if msg.metadata.get("skill_workflow_orchestration"))
    assert "required_skills: baoyu-design" in notice
    assert "<deck-stage" in notice
    assert "html2pptx.md exactly" not in notice


def test_continue_keeps_active_deliverable_contract_and_original_artifact_dir(tmp_path):
    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    original_dir = str(tmp_path / "rag_enterprise_kb")
    session = Session(session_id="s-continue-contract", user_id="u1", channel="feishu")
    original_task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    original_contract = {
        "kind": "file_deliverable",
        "task_text": original_task,
        "artifact_dir": original_dir,
        "required_evidence": ["file_created", "send_file_to_user"],
        "max_repair_attempts": 2,
        "source_message_id": "m-original",
        "task_fingerprint": agent._task_fingerprint(original_task),
        "required_skills": ["baoyu-design"],
    }
    session.metadata["current_completion_contract"] = original_contract
    session.messages.extend([
        Message(id="m-original", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.USER, content=original_task),
        Message(id="a-failed", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.ASSISTANT, content="我又没有真正把文件写出来，任务未完成。"),
        Message(id="m-continue", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.USER, content="继续"),
    ])

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.task_text == original_task
    assert contract.artifact_dir == original_dir
    assert contract.required_skills == ("baoyu-design",)


def test_repeated_continue_artifact_task_skips_previous_confirmation(tmp_path):
    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    session = Session(session_id="s-repeated-continue", user_id="u1", channel="feishu")
    session.messages.extend([
        Message(id="m-original", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.USER, content="做一个关于 RAG 企业知识库的幻灯片"),
        Message(id="a1", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.ASSISTANT, content="需要继续生成。"),
        Message(id="m-c1", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.USER, content="继续"),
        Message(id="a2", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.ASSISTANT, content="仍未生成文件。"),
        Message(id="m-c2", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.USER, content="继续"),
    ])

    task_text = agent._artifact_task_text_for_completion_contract(session, "继续", "")

    assert task_text == "做一个关于 RAG 企业知识库的幻灯片"
    assert task_text != "继续"


def test_short_deck_continuation_without_metadata_recovers_previous_explicit_skill_contract(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    pptx_dir = skills_root / "pptx"
    (pptx_dir / "scripts").mkdir(parents=True)
    (pptx_dir / "SKILL.md").write_text("---\nname: pptx\n---\n[html2pptx](html2pptx.md)\n", encoding="utf-8")
    (pptx_dir / "html2pptx.md").write_text("# html2pptx workflow\n", encoding="utf-8")
    (pptx_dir / "scripts" / "html2pptx.js").write_text("// converter\n", encoding="utf-8")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools.skills_dirs = [skills_root]
    agent = Agent(AsyncMock(), tools, MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    session = Session(session_id="s-continue-without-metadata", user_id="u1", channel="feishu")
    original_task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    session.messages.extend([
        Message(id="m-original", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.USER, content=original_task),
        Message(id="a-failed", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.ASSISTANT, content="这次没能交付，回复「继续生成 deck」。"),
        Message(id="m-continue", channel="feishu", channel_user_id="u1", user_id="u1", session_id=session.session_id, type=MessageType.TEXT, role=MessageRole.USER, content="继续生成 deck"),
    ])

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.task_text == original_task
    assert contract.required_skills == ("baoyu-design",)
    assert "继续生成_deck" not in contract.artifact_dir


def test_short_deck_continuation_prefers_durable_explicit_skill_contract_without_history(tmp_path):
    """A compacted Feishu continuation must resume the original skill objective.

    Live failure pattern: history compaction/reload left only ``继续生成 deck``
    plus polluted current/last-incomplete contracts.  The controller must use
    the durable explicit-skill contract instead of routing to generic pptx.
    """
    from pyclaw.core.completion_contract import CompletionContract

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    original_task = "走完整的baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    original = CompletionContract(
        kind="file_deliverable",
        task_text=original_task,
        artifact_dir=str(tmp_path / "artifacts" / "rag_enterprise_kb"),
        required_evidence=("file_created", "send_file_to_user"),
        source_message_id="m-original-baoyu",
        task_fingerprint=agent._task_fingerprint(original_task),
        required_skills=("baoyu-design",),
    ).to_metadata()
    polluted = CompletionContract(
        kind="file_deliverable",
        task_text="继续生成 deck",
        artifact_dir=str(tmp_path / "artifacts" / "继续生成_deck_b1203608-4e1_om_x100b6be8"),
        required_evidence=("file_created", "send_file_to_user"),
        source_message_id="m-continue-deck",
        task_fingerprint=agent._task_fingerprint("继续生成 deck"),
        required_skills=("pptx",),
    ).to_metadata()
    session = Session(session_id="s-durable-skill-contract", user_id="u1", channel="feishu")
    session.metadata["last_explicit_skill_completion_contract"] = original
    session.metadata["current_completion_contract"] = polluted
    session.metadata["last_incomplete_completion_contract"] = polluted
    session.messages.append(Message(
        id="m-continue-deck",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="继续生成 deck",
    ))

    contract = agent._infer_completion_contract(session)

    assert contract is not None
    assert contract.task_text == original_task
    assert contract.required_skills == ("baoyu-design",)
    assert "继续生成_deck" not in contract.artifact_dir
    assert session.metadata["current_completion_contract"]["task_text"] == original_task


def test_short_deck_continuation_without_state_does_not_create_generic_pptx_contract(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    pptx_dir = skills_root / "pptx"
    (pptx_dir / "scripts").mkdir(parents=True)
    (pptx_dir / "SKILL.md").write_text("---\nname: pptx\n---\n[html2pptx](html2pptx.md)\n", encoding="utf-8")
    (pptx_dir / "html2pptx.md").write_text("# html2pptx workflow\n", encoding="utf-8")
    (pptx_dir / "scripts" / "html2pptx.js").write_text("// converter\n", encoding="utf-8")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools.skills_dirs = [skills_root]
    agent = Agent(AsyncMock(), tools, MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    session = Session(session_id="s-orphan-continuation", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-continue-deck",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="继续生成 deck",
    ))

    contract = agent._infer_completion_contract(session)

    assert contract is None
    assert "current_completion_contract" not in session.metadata


def test_controller_adopts_verified_workspace_artifact_when_pending_files_lost(tmp_path):
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    session = Session(session_id="s-adopt-workspace-artifact", user_id="u1", channel="feishu")
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    artifact_dir = tmp_path / "rag_deck_artifacts"
    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    session.metadata["active_skill_contexts"] = [{
        "name": "baoyu-design",
        "canonical_rel_path": "baoyu-design/skills/baoyu-design",
        "skill_md_path": str(skill_dir / "SKILL.md"),
        "root_dir": str(skill_dir),
        "content_sha256": "sha",
    }]
    session.metadata["active_skills"] = ["baoyu-design"]
    contract = agent.completion_contracts.infer(
        task_text=task,
        pending_context="",
        artifact_dir=str(artifact_dir),
    )
    assert contract is not None
    contract = contract.__class__(
        kind=contract.kind,
        task_text=task,
        artifact_dir=str(artifact_dir),
        required_evidence=contract.required_evidence,
        max_repair_attempts=contract.max_repair_attempts,
        source_message_id="m-rag-deck",
        task_fingerprint=agent._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time() - 5,
    )
    session.metadata["current_completion_contract"] = contract.to_metadata()
    _write_skill_deck_html(artifact_dir / "deck.html")
    session.messages.append(Message(
        id="m-rag-deck",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    synthesized = ArtifactSynthesisService().synthesize_skill_deck(
        contract,
        draft="",
        quality=SynthesisQuality.FULL,
    )
    assert synthesized is not None
    (artifact_dir / "skill-workflow-evidence.md").write_text(
        "\n".join([
            "# Skill Workflow Evidence",
            "```json",
            json.dumps({
                "producer": "model_tool_workflow",
                "task_text": task,
                "task_fingerprint": agent._task_fingerprint(task),
                "artifact_dir": str(artifact_dir),
                "required_skills": ["baoyu-design"],
                "outputs": [synthesized.file_path],
                "output_facts": [{"path": synthesized.file_path, "exists": True, "sha256": "test"}],
            }, ensure_ascii=False),
            "```",
            "SKILL.md",
            "system-prompt.md",
            "references/codex.md",
            "built-in-skills/make-a-deck.md",
            "built-in-skills/export-as-pptx-editable.md",
            "<deck-stage width=\"1920\" height=\"1080\">",
            "section data-label",
        ]),
        encoding="utf-8",
    )

    pending_files = []
    content = agent._prepare_completion_contract_final_content(
        session=session,
        content="没有产出任何文件，任务未完成。",
        pending_files=pending_files,
    )

    assert len(pending_files) == 1
    assert pending_files[0]["file_path"] == synthesized.file_path
    assert "已生成并发送文件" in content
    assert "没有产出任何文件" not in content
    assert "任务未完成" not in content
    assert "current_completion_contract" not in session.metadata


def test_repeated_explicit_skill_task_reuses_matching_manifest_artifact(tmp_path):
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.deliverable_workflow import DeliverableWorkflow
    from pyclaw.core.skill_context import SkillContextService
    from pyclaw.core.skill_evidence import SkillEvidenceService
    from pyclaw.core.skill_workspace import SkillWorkspaceService

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    artifact_dir = tmp_path / "rag_deck_artifacts"
    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)

    old_contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-old",
        task_fingerprint=agent._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time() - 7200,
    )
    _write_skill_deck_html(artifact_dir / "deck.html")
    synthesized = ArtifactSynthesisService().synthesize_skill_deck(
        old_contract,
        draft="",
        quality=SynthesisQuality.FULL,
    )
    assert synthesized is not None

    session = Session(
        session_id="s-reuse-skill-manifest",
        user_id="u1",
        channel="feishu",
        metadata={"active_skill_contexts": [{
            "name": "baoyu-design",
            "canonical_rel_path": "baoyu-design/skills/baoyu-design",
            "skill_md_path": str(skill_dir / "SKILL.md"),
            "root_dir": str(skill_dir),
            "content_sha256": "sha",
        }]},
    )
    contexts = SkillContextService()
    evidence_service = SkillEvidenceService(contexts)
    workspace = SkillWorkspaceService(contexts, evidence_service)
    workspace.write_manifest(
        session=session,
        contract=old_contract,
        pending_files=[{"file_path": synthesized.file_path, "description": "RAG deck"}],
    )
    manifest = artifact_dir / "skill-workflow-evidence.md"
    manifest.write_text(
        manifest.read_text(encoding="utf-8").replace(
            '"producer": "controller_workspace_adapter"',
            '"producer": "model_tool_workflow"',
        ),
        encoding="utf-8",
    )

    old_time = time.time() - 3600
    for path in (artifact_dir / "deck.html", artifact_dir / "skill-workflow-evidence.md", artifact_dir / synthesized.file_path):
        os.utime(path, (old_time, old_time))

    new_contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-new",
        task_fingerprint=agent._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time(),
    )
    session.messages.append(Message(
        id="m-new",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))
    for rel in ("system-prompt.md", "references/codex.md", "built-in-skills/make-a-deck.md", "built-in-skills/export-as-pptx-editable.md"):
        session.messages.append(Message(
            id=f"tool-new-{rel}",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content=f"OBSERVATION from read_file:\nFile: {skill_dir / rel}\n...",
            metadata={"tool_name": "read_file"},
        ))
    workflow = DeliverableWorkflow(skill_evidence=evidence_service, skill_workspace=workspace)
    pending_files: list[dict[str, str]] = []

    finalization = workflow.finalize(
        contract=new_contract,
        content="已生成的文件未通过 skill 工作流验收，任务未完成。",
        pending_files=pending_files,
        synthesis_quality=SynthesisQuality.BASIC,
        session=session,
    )

    assert len(finalization.pending_files) == 1
    assert finalization.pending_files[0]["file_path"] == synthesized.file_path
    assert finalization.acceptance is not None
    assert finalization.acceptance.accepted
    assert "已生成并发送文件" in finalization.content
    assert "任务未完成" not in finalization.content


def test_explicit_full_skill_request_rejects_controller_adapter_manifest_only(tmp_path):
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.deliverable_workflow import DeliverableWorkflow
    from pyclaw.core.skill_context import SkillContextService
    from pyclaw.core.skill_evidence import SkillEvidenceService
    from pyclaw.core.skill_workspace import SkillWorkspaceService

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    artifact_dir = tmp_path / "rag_deck_artifacts"
    skill_dir = tmp_path / "skills" / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)

    old_contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-old",
        task_fingerprint=agent._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time() - 7200,
    )
    _write_skill_deck_html(artifact_dir / "deck.html")
    synthesized = ArtifactSynthesisService().synthesize_skill_deck(
        old_contract,
        draft="",
        quality=SynthesisQuality.FULL,
    )
    assert synthesized is not None
    session = Session(
        session_id="s-controller-manifest-rejected",
        user_id="u1",
        channel="feishu",
        metadata={"active_skill_contexts": [{
            "name": "baoyu-design",
            "canonical_rel_path": "baoyu-design/skills/baoyu-design",
            "skill_md_path": str(skill_dir / "SKILL.md"),
            "root_dir": str(skill_dir),
            "content_sha256": "sha",
        }]},
    )
    contexts = SkillContextService()
    evidence_service = SkillEvidenceService(contexts)
    workspace = SkillWorkspaceService(contexts, evidence_service)
    workspace.write_manifest(
        session=session,
        contract=old_contract,
        pending_files=[{"file_path": synthesized.file_path, "description": "controller fallback deck"}],
    )
    old_time = time.time() - 3600
    for path in (artifact_dir / "deck.html", artifact_dir / "skill-workflow-evidence.md", artifact_dir / synthesized.file_path):
        os.utime(path, (old_time, old_time))

    new_contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-new",
        task_fingerprint=agent._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time(),
    )
    session.messages.append(Message(
        id="m-new",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))
    workflow = DeliverableWorkflow(skill_evidence=evidence_service, skill_workspace=workspace)
    pending_files: list[dict[str, str]] = []

    finalization = workflow.finalize(
        contract=new_contract,
        content="已生成的文件未通过 skill 工作流验收，任务未完成。",
        pending_files=pending_files,
        synthesis_quality=SynthesisQuality.DISABLED,
        allow_skill_synthesis=True,
        session=session,
    )

    assert finalization.pending_files == ()
    assert "已生成并发送文件" not in finalization.content
    assert "未通过 skill 工作流验收" in finalization.content


def test_repeated_explicit_skill_task_rejects_mismatched_manifest_artifact(tmp_path):
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.skill_evidence import SkillEvidenceService

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    artifact_dir = tmp_path / "rag_deck_artifacts"
    artifact_dir.mkdir()
    artifact = artifact_dir / "RAG.pptx"
    _write_test_pptx(artifact, 10)
    (artifact_dir / "deck.html").write_text(
        '<deck-stage width="1920" height="1080"><section data-label="RAG">RAG 企业知识库</section></deck-stage>',
        encoding="utf-8",
    )
    old_task = "走完整的 baoyu design skill 做一个关于 AI Agent 的幻灯片"
    (artifact_dir / "skill-workflow-evidence.md").write_text(
        "\n".join([
            "# Skill Workflow Evidence",
            "```json",
            json.dumps({
                "task_text": old_task,
                "task_fingerprint": agent._task_fingerprint(old_task),
                "artifact_dir": str(artifact_dir),
                "required_skills": ["baoyu-design"],
                "outputs": [str(artifact)],
            }, ensure_ascii=False),
            "```",
            "<deck-stage",
            "section data-label",
        ]),
        encoding="utf-8",
    )
    old_time = time.time() - 3600
    for path in (artifact, artifact_dir / "deck.html", artifact_dir / "skill-workflow-evidence.md"):
        os.utime(path, (old_time, old_time))

    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        task_fingerprint=agent._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time(),
    )

    acceptance = ArtifactAcceptanceService().evaluate(contract, [{"file_path": str(artifact)}])
    assert acceptance.accepted is False
    assert any("早于当前任务" in reason for reason in acceptance.reasons)
    assert SkillEvidenceService()._is_stale_for_contract(artifact_dir / "deck.html", contract) is True


def test_repeated_explicit_skill_task_rejects_same_task_manifest_without_producer(tmp_path):
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.skill_evidence import SkillEvidenceService

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    artifact_dir = tmp_path / "rag_deck_artifacts"
    artifact_dir.mkdir()
    artifact = artifact_dir / "RAG.pptx"
    _write_test_pptx(artifact, 10)
    (artifact_dir / "deck.html").write_text(
        '<deck-stage width="1920" height="1080"><section data-label="RAG">RAG 企业知识库</section></deck-stage>',
        encoding="utf-8",
    )
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    (artifact_dir / "skill-workflow-evidence.md").write_text(
        "\n".join([
            "# Skill Workflow Evidence",
            "```json",
            json.dumps({
                "task_text": task,
                "task_fingerprint": agent._task_fingerprint(task),
                "artifact_dir": str(artifact_dir),
                "required_skills": ["baoyu-design"],
                "outputs": [str(artifact)],
                "output_facts": [{"path": str(artifact), "exists": True, "sha256": "test"}],
            }, ensure_ascii=False),
            "```",
            "system-prompt.md",
            "references/codex.md",
            "built-in-skills/make-a-deck.md",
            "built-in-skills/export-as-pptx-editable.md",
            "<deck-stage",
            "section data-label",
        ]),
        encoding="utf-8",
    )
    old_time = time.time() - 3600
    for path in (artifact, artifact_dir / "deck.html", artifact_dir / "skill-workflow-evidence.md"):
        os.utime(path, (old_time, old_time))

    contract = CompletionContract(
        kind="file_deliverable",
        task_text=task,
        artifact_dir=str(artifact_dir),
        source_message_id="m-new",
        task_fingerprint=agent._task_fingerprint(task),
        required_skills=("baoyu-design",),
        created_at=time.time(),
    )

    acceptance = ArtifactAcceptanceService().evaluate(contract, [{"file_path": str(artifact)}])

    assert acceptance.accepted is False
    assert any("早于当前任务" in reason for reason in acceptance.reasons)
    assert SkillEvidenceService()._is_stale_for_contract(artifact_dir / "deck.html", contract) is True


def test_deliverable_workflow_rejects_wrong_topic_and_synthesizes_real_topic(tmp_path):
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService
    from pyclaw.core.artifact_synthesis import ArtifactSynthesisService, SynthesisQuality
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.deliverable_workflow import DeliverableWorkflow

    wrong = tmp_path / "wrong_topic.pptx"
    _write_wrong_topic_pptx(wrong, 8)
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="做一个关于 RAG 企业知识库的幻灯片",
        artifact_dir=str(tmp_path),
    )
    workflow = DeliverableWorkflow(
        acceptance=ArtifactAcceptanceService(),
        synthesis=ArtifactSynthesisService(),
    )
    pending_files = [{"file_path": str(wrong), "description": "RAG 幻灯片"}]

    evidence = workflow.evidence(contract=contract, pending_files=pending_files)
    assert evidence.artifact_acceptance is not None
    assert not evidence.artifact_acceptance.accepted
    assert "主题不匹配" in evidence.artifact_acceptance.summary

    finalization = workflow.finalize(
        contract=contract,
        content="",
        pending_files=pending_files,
        synthesis_quality=SynthesisQuality.BASIC,
        force_repair_synthesis=True,
    )

    assert len(finalization.pending_files) == 1
    replacement = finalization.pending_files[0]["file_path"]
    assert replacement != str(wrong)
    text = _pptx_text(replacement)
    assert "RAG" in text
    assert "企业知识库" in text
    assert "Kubernetes" not in text
    assert finalization.acceptance is not None
    assert finalization.acceptance.accepted


def test_required_skills_for_task_resolves_frontmatter_skill_name(tmp_path, monkeypatch):
    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "rich-deck" / "skills" / "rich-deck"
    skill_dir.mkdir(parents=True)
    (skill_dir / "SKILL.md").write_text(
        "---\nname: rich-deck\ndescription: Rich deck workflow.\n---\n",
        encoding="utf-8",
    )
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    agent = Agent(AsyncMock(), tools, AsyncMock())

    assert agent._required_skills_for_task("走完整的 rich deck skill 做一个关于 RAG 的幻灯片") == ("rich-deck",)
    assert agent._required_skills_for_task("使用 /rich-deck 做一个关于 RAG 的幻灯片") == ("rich-deck",)
    assert agent._required_skills_for_task("用 $rich-deck 生成 PPT") == ("rich-deck",)


def test_required_skills_for_task_uses_active_context_when_explicit_but_name_omitted(tmp_path):
    skill_dir = tmp_path / "skills" / "active-deck"
    skill_dir.mkdir(parents=True)
    (skill_dir / "SKILL.md").write_text("---\nname: active-deck\n---\n", encoding="utf-8")
    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    session = Session(
        session_id="s-active-skill-fallback",
        user_id="u1",
        channel="feishu",
        metadata={"active_skill_contexts": [{
            "name": "active-deck",
            "canonical_rel_path": "active-deck",
            "skill_md_path": str(skill_dir / "SKILL.md"),
            "root_dir": str(tmp_path / "skills"),
        }]},
    )

    assert agent._required_skills_for_task("按这个技能做一个关于 RAG 的幻灯片", session) == ("active-deck",)


@pytest.mark.asyncio
async def test_plain_ppt_deliverable_auto_routes_to_pptx_skill_and_hydrates_docs(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.skill_evidence import SkillEvidenceService

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "pptx"
    skill_dir.mkdir(parents=True)
    (skill_dir / "SKILL.md").write_text(
        "---\nname: pptx\ndescription: Presentation creation.\n---\n"
        "# PPTX\n"
        "When creating a new PowerPoint presentation from scratch, use the html2pptx workflow.\n"
        "MANDATORY: read html2pptx.md before creating slides.\n"
        "[html2pptx](html2pptx.md)\n",
        encoding="utf-8",
    )
    (skill_dir / "html2pptx.md").write_text(
        "# html2pptx workflow\nCreate HTML slides, convert with html2pptx, validate thumbnails.\n",
        encoding="utf-8",
    )
    (skill_dir / "scripts").mkdir()
    (skill_dir / "scripts" / "html2pptx.js").write_text("// html2pptx library\n", encoding="utf-8")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    sessions = AsyncMock()

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    agent = Agent(AsyncMock(), tools, sessions)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    task = "做一个关于 RAG 企业知识库的 12 页幻灯片"
    session = Session(session_id="s-auto-pptx-skill", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-auto-pptx-skill",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    contract = agent._infer_completion_contract(session)
    assert contract is not None
    assert contract.required_skills == ("pptx",)

    await agent._ensure_skill_workflow_hydration_observations(session, contract)
    await agent._ensure_skill_workflow_orchestration_notice(session, contract)

    hydrated = [msg for msg in session.messages if msg.metadata.get("controller_skill_hydration")]
    assert hydrated
    combined = "\n".join(msg.content for msg in hydrated)
    assert "Required path: SKILL.md" in combined
    assert "Required path: html2pptx.md" in combined
    assert "Required path: scripts/html2pptx.js" in combined
    assert "html2pptx workflow" in combined

    result = SkillEvidenceService(agent.skill_contexts).evaluate(session=session, contract=contract, pending_files=[])
    assert result is not None
    assert "SKILL.md" in result.observed_paths

    notice = "\n".join(
        msg.content for msg in session.messages
        if msg.role == MessageRole.USER and msg.metadata.get("skill_workflow_orchestration")
    )
    assert "required_skills: pptx" in notice
    assert "SKILL.md" in notice
    assert "html2pptx.md" in notice
    assert "scripts/html2pptx.js" in notice
    assert "<deck-stage" not in notice


def test_real_pptx_skill_requirement_is_html2pptx_not_deck_stage(tmp_path):
    from pyclaw.core.completion_contract import CompletionContract
    from pyclaw.core.skill_context import ActiveSkillContext
    from pyclaw.core.skill_evidence import SkillEvidenceService

    skill_dir = tmp_path / "skills" / "pptx"
    (skill_dir / "scripts").mkdir(parents=True)
    (skill_dir / "SKILL.md").write_text(
        "---\nname: pptx\n---\n"
        "# PPTX\n"
        "MANDATORY: read [html2pptx.md](html2pptx.md).\n"
        "Use [`html2pptx.js`](scripts/html2pptx.js) to convert slides.\n",
        encoding="utf-8",
    )
    (skill_dir / "html2pptx.md").write_text("# html2pptx workflow\n", encoding="utf-8")
    (skill_dir / "scripts" / "html2pptx.js").write_text("// converter\n", encoding="utf-8")
    ctx = ActiveSkillContext(
        name="pptx",
        canonical_rel_path="pptx",
        skill_md_path=str(skill_dir / "SKILL.md"),
        root_dir=str(tmp_path / "skills"),
    )
    contract = CompletionContract(
        kind="file_deliverable",
        task_text="做一个关于 RAG 的 12 页 PPT",
        artifact_dir=str(tmp_path / "artifacts"),
        required_skills=("pptx",),
    )

    requirement = SkillEvidenceService().infer_requirement(ctx, contract)

    assert requirement.required_paths == ("SKILL.md", "html2pptx.md", "scripts/html2pptx.js")
    assert "<deck-stage" not in requirement.required_output_markers
    assert "*.html:<deck-stage" not in requirement.required_file_patterns


def test_plain_ppt_deliverable_with_auto_skill_disables_generic_fallback(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.artifact_synthesis import SynthesisQuality

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "pptx"
    skill_dir.mkdir(parents=True)
    (skill_dir / "SKILL.md").write_text("---\nname: pptx\n---\n# PPTX\n", encoding="utf-8")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    agent = Agent(AsyncMock(), tools, AsyncMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    session = Session(session_id="s-auto-pptx-no-fallback", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-auto-pptx-no-fallback",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个关于 AI Agent 的 12 页 PPT",
    ))

    contract = agent._infer_completion_contract(session)
    assert contract is not None
    assert contract.required_skills == ("pptx",)
    assert agent._synthesis_quality_for_contract(session, contract, "## 大纲\n- a\n- b\n- c\n- d\n- e\n- f\n- g\n- h") == SynthesisQuality.DISABLED


@pytest.mark.asyncio
async def test_explicit_skill_file_contract_injects_controller_workflow_notice(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    sessions = AsyncMock()
    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)
    sessions.save_message.side_effect = save_message
    agent = Agent(AsyncMock(), tools, sessions)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的 12 页幻灯片"
    session = Session(session_id="s-orchestrate-skill", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-orchestrate-skill",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    contract = agent._infer_completion_contract(session)
    await agent._ensure_skill_workflow_orchestration_notice(session, contract)
    await agent._ensure_skill_workflow_orchestration_notice(session, contract)

    notices = [
        msg for msg in session.messages
        if msg.role == MessageRole.USER and msg.metadata.get("skill_workflow_orchestration")
    ]
    assert len(notices) == 1
    notice = notices[0].content
    assert "read_file" in notice
    assert "deck.html" in notice
    assert "<deck-stage" in notice
    assert "send_file_to_user" in notice
    assert contract is not None and contract.artifact_dir in notice
    assert "system-prompt.md" in notice
    assert "built-in-skills/make-a-deck.md" in notice
    assert "built-in-skills/export-as-pptx-editable.md" in notice


@pytest.mark.asyncio
async def test_explicit_skill_file_contract_hydrates_required_docs_as_current_turn_evidence(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.skill_evidence import SkillEvidenceService

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    sessions = AsyncMock()

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    agent = Agent(AsyncMock(), tools, sessions)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的 12 页幻灯片"
    session = Session(session_id="s-hydrate-skill", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-hydrate-skill",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    contract = agent._infer_completion_contract(session)
    assert contract is not None
    await agent._ensure_skill_workflow_hydration_observations(session, contract)
    await agent._ensure_skill_workflow_hydration_observations(session, contract)

    hydrated = [msg for msg in session.messages if msg.metadata.get("controller_skill_hydration")]
    assert len(hydrated) >= 4
    assert {msg.metadata.get("tool_name") for msg in hydrated} == {"read_file"}
    combined = "\n".join(msg.content for msg in hydrated)
    assert "system-prompt.md" in combined
    assert "built-in-skills/make-a-deck.md" in combined
    assert "built-in-skills/export-as-pptx-editable.md" in combined
    assert "<deck-stage" in combined

    result = SkillEvidenceService(agent.skill_contexts).evaluate(session=session, contract=contract, pending_files=[])
    assert result is not None
    assert "built-in-skills/make-a-deck.md" in result.observed_paths


@pytest.mark.asyncio
async def test_skill_hydration_key_rehydrates_when_observations_missing(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    sessions = AsyncMock()

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    agent = Agent(AsyncMock(), tools, sessions)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的 12 页幻灯片"
    session = Session(session_id="s-rehydrate-skill", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-rehydrate-skill",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    contract = agent._infer_completion_contract(session)
    assert contract is not None
    session.metadata["skill_workflow_hydration_key"] = (
        f"{contract.source_message_id}:{contract.task_fingerprint}:{','.join(contract.required_skills)}"
    )

    await agent._ensure_skill_workflow_hydration_observations(session, contract)

    hydrated = [msg for msg in session.messages if msg.metadata.get("controller_skill_hydration")]
    assert len(hydrated) >= 4
    assert agent._has_skill_workflow_hydration_observations(session, contract)


def test_controller_skill_hydration_tool_message_is_sent_to_llm_as_user_context():
    msg = Message(
        id="m-hydrated-tool",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id="s1",
        type=MessageType.TEXT,
        role=MessageRole.TOOL,
        content="OBSERVATION from read_file:\nFile: /tmp/SKILL.md\n...",
        metadata={"tool_name": "read_file", "tool_call_id": "synthetic", "controller_skill_hydration": True},
    )

    formatted = msg.to_llm_format()

    assert formatted == {"role": "user", "content": msg.content}


@pytest.mark.asyncio
async def test_explicit_skill_ppt_loop_starts_with_hydrated_skill_docs(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    model = AsyncMock()
    model.chat.return_value = {"content": "我已读取 skill 工作流，继续生成 deck.html 和 PPTX。", "__tool_calls__": False}
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    tools.get_all_specs.return_value = []
    tools.execute_tool_calls = AsyncMock(return_value=[])
    sessions = AsyncMock()
    session = Session(session_id="s-loop-hydrated", user_id="u1", channel="feishu")

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    sessions.get_or_create.return_value = session
    agent = Agent(model, tools, sessions, max_iterations=1)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    user_msg = Message(
        id="m-loop-hydrated",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的 12 页幻灯片",
    )

    await agent.process_message(user_msg)

    assert model.chat.await_count >= 1
    messages = model.chat.await_args.kwargs["messages"]
    transcript = "\n".join(str(item.get("content", "")) for item in messages)
    assert "OBSERVATION from read_file" in transcript
    assert "built-in-skills/make-a-deck.md" in transcript
    assert "built-in-skills/export-as-pptx-editable.md" in transcript
    assert "<deck-stage" in transcript
    assert not any(item.get("role") == "tool" for item in messages if "skill-hydration" in str(item))


@pytest.mark.asyncio
async def test_explicit_skill_ppt_loop_controller_completes_when_model_never_uses_skill(tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    model = AsyncMock()
    model.chat.return_value = {"content": "稍后重试，我先准备一下。", "__tool_calls__": False}
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    tools.get_all_specs.return_value = []
    tools.execute_tool_calls = AsyncMock(return_value=[])
    sessions = AsyncMock()
    session = Session(session_id="s-loop-skill-controller-complete", user_id="u1", channel="feishu")

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    sessions.get_or_create.return_value = session
    agent = Agent(model, tools, sessions, max_iterations=1)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    user_msg = Message(
        id="m-loop-skill-controller-complete",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的 12 页幻灯片",
    )

    response = await agent.process_message(user_msg)

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    pptx_path = pending_files[0]["file_path"]
    assert pptx_path.endswith(".pptx")
    assert os.path.exists(pptx_path)
    assert "已生成并发送文件" in response.content
    assert "任务未完成" not in response.content
    assert "稍后" not in response.content
    artifact_dir = os.path.dirname(pptx_path)
    assert os.path.exists(os.path.join(artifact_dir, "deck.html"))
    assert os.path.exists(os.path.join(artifact_dir, "skill-workflow-evidence.md"))
    deck_html = open(os.path.join(artifact_dir, "deck.html"), encoding="utf-8").read()
    assert "<deck-stage" in deck_html
    assert "section data-label" in deck_html
    text = _pptx_text(pptx_path)
    assert "RAG" in text
    assert "企业知识库" in text
    assert "当前进展" not in text
    hydrated = [msg for msg in session.messages if msg.metadata.get("controller_skill_hydration")]
    assert hydrated


@pytest.mark.asyncio
async def test_exact_baoyu_rag_prompt_controller_completes_after_setup_probe(tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    setup_command = "mkdir -p /tmp/pyclaw-baoyu-rag && ls -la /tmp/pyclaw-baoyu-rag"
    model = AsyncMock()
    model.chat.side_effect = [
        {
            "content": "我先准备 baoyu design 的工作目录。",
            "__tool_calls__": True,
            "tool_calls": [{
                "id": "setup-baoyu-rag",
                "function": {"name": "terminal", "arguments": json.dumps({"command": setup_command})},
            }],
        },
        {
            "content": "还没生成完，稍后重试。",
            "__tool_calls__": False,
        },
    ]
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    tools.get_all_specs.return_value = [{"name": "terminal"}, {"name": "send_file_to_user"}]
    tools.execute_tool_calls = AsyncMock(return_value=[{
        "role": "tool",
        "tool_call_id": "setup-baoyu-rag",
        "name": "terminal",
        "content": "Exit code: 0\ntotal 0\n",
        "success": True,
        "metadata": {},
    }])
    sessions = AsyncMock()
    session = Session(session_id="s-exact-baoyu-probe", user_id="u1", channel="feishu")

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    sessions.get_or_create.return_value = session
    agent = Agent(model, tools, sessions, max_iterations=6)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))

    response = await agent.process_message(Message(
        id="m-exact-baoyu-probe",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
    ))

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    pptx_path = pending_files[0]["file_path"]
    assert pptx_path.endswith(".pptx")
    assert os.path.exists(pptx_path)
    artifact_dir = os.path.dirname(pptx_path)
    assert os.path.exists(os.path.join(artifact_dir, "deck.html"))
    assert os.path.exists(os.path.join(artifact_dir, "skill-workflow-evidence.md"))
    assert tools.execute_tool_calls.await_count == 1
    assert model.chat.await_count == 1

    assert "已生成并发送文件" in response.content
    forbidden = ("任务未完成", "稍后", "继续生成", "锁屏", "skill 工作流验收", "当前进展")
    assert not any(marker in response.content for marker in forbidden)

    from pptx import Presentation

    presentation = Presentation(pptx_path)
    assert len(presentation.slides) >= 12
    text = _pptx_text(pptx_path)
    for marker in ("RAG", "企业知识库", "检索", "向量", "权限", "评测"):
        assert marker in text
    assert not any(marker in text for marker in forbidden)



@pytest.mark.asyncio
async def test_repeated_baoyu_rag_prompt_rehydrates_docs_with_contract_scoped_ids(tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    session = Session(session_id="s-repeated-baoyu-real-history", user_id="u1", channel="feishu")

    # Real-history regression: an earlier turn in the same Feishu session had
    # already hydrated the same baoyu files.  Old PyClaw used ids based only on
    # abs_path + session_id, so the new turn's controller hydration was skipped
    # as duplicate and the evidence gate later complained that SKILL.md had not
    # been read for the current task.
    old_source_id = "om_old_baoyu"
    old_fingerprint = Agent(AsyncMock(), MagicMock(), MagicMock())._task_fingerprint(task)
    for rel in (
        "SKILL.md",
        "system-prompt.md",
        "references/codex.md",
        "built-in-skills/make-a-deck.md",
        "built-in-skills/export-as-pptx-editable.md",
    ):
        abs_path = str(skill_dir / rel)
        old_digest = hashlib.sha256(abs_path.encode("utf-8")).hexdigest()[:12]
        session.messages.append(Message(
            id=f"skill-hydration-{old_digest}-{session.session_id}",
            channel="feishu",
            channel_user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.TOOL,
            content=(
                "OBSERVATION from read_file:\n"
                f"File: {abs_path}\n"
                "Skill: baoyu-design\n"
                f"Required path: {rel}\n\n"
                f"old hydration for {rel}"
            ),
            metadata={
                "tool_name": "read_file",
                "controller_skill_hydration": True,
                "source_message_id": old_source_id,
                "task_fingerprint": old_fingerprint,
                "skill": "baoyu-design",
                "required_path": rel,
            },
        ))

    model = AsyncMock()
    model.chat.return_value = {
        "content": "已生成的文件未通过 skill 工作流验收，任务未完成。",
        "__tool_calls__": False,
    }
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "terminal"}, {"name": "send_file_to_user"}]
    tools.execute_tool_calls = AsyncMock(return_value=[])
    sessions = AsyncMock()

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    sessions.get_or_create.return_value = session
    agent = Agent(model, tools, sessions, max_iterations=3)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))

    response = await agent.process_message(Message(
        id="om_new_baoyu",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    ))

    current_hydration = [
        msg for msg in session.messages
        if msg.metadata.get("controller_skill_hydration")
        and msg.metadata.get("source_message_id") == "om_new_baoyu"
    ]
    assert any(msg.metadata.get("required_path") == "SKILL.md" for msg in current_hydration)
    assert len(current_hydration) >= 5

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    pptx_path = pending_files[0]["file_path"]
    assert os.path.exists(pptx_path)
    assert os.path.exists(os.path.join(os.path.dirname(pptx_path), "deck.html"))
    assert "已生成并发送文件" in response.content
    assert "skill 工作流验收" not in response.content
    assert "任务未完成" not in response.content
    text = _pptx_text(pptx_path)
    for marker in ("RAG", "企业知识库", "检索", "向量", "权限", "评测"):
        assert marker in text


@pytest.mark.asyncio
async def test_exact_baoyu_rag_prompt_duplicate_side_effect_branch_controller_completes(tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    model = AsyncMock()
    command = "python3 - <<'PY'\nprint('would create deck')\nPY"
    model.chat.return_value = {
        "content": "我继续执行生成命令。",
        "__tool_calls__": True,
        "tool_calls": [{
            "id": "repeat-generate",
            "function": {"name": "terminal", "arguments": json.dumps({"command": command})},
        }],
    }
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    tools.get_all_specs.return_value = [{"name": "terminal"}, {"name": "send_file_to_user"}]
    tools.execute_tool_calls = AsyncMock(return_value=[])
    sessions = AsyncMock()
    session = Session(session_id="s-exact-baoyu-duplicate", user_id="u1", channel="feishu")

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    agent = Agent(model, tools, sessions, max_iterations=6)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    user_msg = Message(
        id="m-exact-baoyu-duplicate",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=task,
    )
    await save_message(session, user_msg)
    contract = agent._infer_completion_contract(session)
    assert contract is not None
    agent.skill_workspace.ensure_required_contexts(session, contract)
    await agent._ensure_skill_workflow_hydration_observations(session, contract)
    await save_message(session, Message(
        id="assistant-previous-repeat",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.ASSISTANT,
        content="执行生成命令。",
        metadata={"tool_calls": [{"id": "previous-generate", "function": {"name": "terminal", "arguments": json.dumps({"command": command})}}]},
    ))
    await save_message(session, Message(
        id="tool-previous-repeat",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.TOOL,
        content="OBSERVATION from terminal:\nExit code: 0\nwould create deck\n",
        metadata={"tool_name": "terminal", "tool_call_id": "previous-generate"},
    ))

    content, pending_files = await agent._agent_loop(session)

    assert tools.execute_tool_calls.await_count == 0
    assert len(pending_files) == 1
    pptx_path = pending_files[0]["file_path"]
    assert pptx_path.endswith(".pptx")
    assert os.path.exists(pptx_path)
    assert os.path.exists(os.path.join(os.path.dirname(pptx_path), "deck.html"))
    assert "已生成并发送文件" in content
    assert "任务未完成" not in content
    assert "副作用工具重复调用" not in content


@pytest.mark.asyncio
async def test_exact_baoyu_rag_prompt_discovers_default_skill_root_and_overrides_process_reply(tmp_path, monkeypatch):
    import re
    import zipfile

    pytest.importorskip("pptx")
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "default-skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    monkeypatch.setattr("pyclaw.core.agent._available_skills_dirs", lambda: [str(skills_root)])

    model = AsyncMock()
    model.chat.return_value = {
        "content": "当前进展：工作目录已创建，但任务未完成，请稍后重试。我可能还需要先锁屏。",
        "__tool_calls__": False,
    }
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    # This mirrors the production/default ToolRegistry path that originally
    # regressed: activate_skill could find ~/.pyclaw/skills, but the agent's
    # controller/index used an empty configured skills_dirs list.
    tools.skills_dirs = []
    tools.get_all_specs.return_value = [{"name": "terminal"}, {"name": "send_file_to_user"}]
    tools.execute_tool_calls = AsyncMock(return_value=[])
    sessions = AsyncMock()
    session = Session(session_id="s-exact-baoyu-default-root", user_id="u1", channel="feishu")

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    sessions.get_or_create.return_value = session
    agent = Agent(model, tools, sessions, max_iterations=4)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))

    assert "baoyu-design" in agent._get_skills_index()
    assert agent._required_skills_for_task("走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片") == ("baoyu-design",)

    response = await agent.process_message(Message(
        id="m-exact-baoyu-default-root",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片",
    ))

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    pptx_path = pending_files[0]["file_path"]
    assert pptx_path.endswith(".pptx")
    assert os.path.exists(pptx_path)
    assert os.path.getsize(pptx_path) > 30_000
    artifact_dir = os.path.dirname(pptx_path)
    assert os.path.exists(os.path.join(artifact_dir, "deck.html"))
    assert os.path.exists(os.path.join(artifact_dir, "skill-workflow-evidence.md"))

    forbidden = ("任务未完成", "稍后", "继续生成", "锁屏", "当前进展", "副作用工具重复调用", "工作流验收")
    assert "已生成并发送文件" in response.content
    assert not any(marker in response.content for marker in forbidden)

    text = _pptx_text(pptx_path)
    for keyword in ("RAG", "企业知识库", "检索", "向量", "权限", "评测"):
        assert keyword in text
    assert "AI Agent 架构蓝图" not in text
    assert not any(marker in text for marker in forbidden)

    with zipfile.ZipFile(pptx_path) as zf:
        slide_names = [name for name in zf.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)]
        slide_xml = "".join(zf.read(name).decode("utf-8", errors="ignore") for name in slide_names)
    assert len(slide_names) >= 12
    assert slide_xml.count("solidFill") >= 20
    assert len(set(re.findall(r"srgbClr val=\"([0-9A-Fa-f]+)\"", slide_xml))) >= 6

@pytest.mark.asyncio
async def test_polluted_continue_contract_recovers_original_skill_task_after_reload(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.session import SessionManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    pptx_dir = skills_root / "pptx"
    (pptx_dir / "scripts").mkdir(parents=True)
    (pptx_dir / "SKILL.md").write_text("---\nname: pptx\n---\n[html2pptx](html2pptx.md)\n", encoding="utf-8")
    (pptx_dir / "html2pptx.md").write_text("# html2pptx workflow\n", encoding="utf-8")
    (pptx_dir / "scripts" / "html2pptx.js").write_text("// converter\n", encoding="utf-8")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    db_path = tmp_path / "pyclaw.db"
    manager = SessionManager(str(db_path))
    await manager.init_db()
    session = await manager.get_or_create("feishu", "u1")
    original_task = "走完整的baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    polluted_task = "继续生成 deck"
    original = Message(
        id="m-original-baoyu",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=original_task,
    )
    failed = Message(
        id="a-failed-baoyu",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.ASSISTANT,
        content="这次没能交付，回复「继续生成 deck」。",
    )
    await manager.save_message(session, original)
    await manager.save_message(session, failed)

    fingerprint_agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    polluted_contract = {
        "kind": "file_deliverable",
        "task_text": polluted_task,
        "artifact_dir": str(tmp_path / "artifacts" / "继续生成_deck_b1203608-4e1_om_x100b6be8"),
        "required_evidence": ["file_created", "send_file_to_user"],
        "max_repair_attempts": 2,
        "source_message_id": "m-continue-deck",
        "task_fingerprint": fingerprint_agent._task_fingerprint(polluted_task),
        "required_skills": ["pptx"],
    }
    session.metadata["current_completion_contract"] = polluted_contract.copy()
    session.metadata["last_incomplete_completion_contract"] = polluted_contract.copy()
    await manager.save_message(session, failed)

    manager2 = SessionManager(str(db_path))
    session2 = await manager2.get_or_create("feishu", "u1")
    model = AsyncMock()
    model.chat.return_value = {"content": "稍后重试，我先准备一下。", "__tool_calls__": False}
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    tools.get_all_specs.return_value = []
    tools.execute_tool_calls = AsyncMock(return_value=[])
    agent = Agent(model, tools, manager2, max_iterations=1)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))

    user2 = Message(
        id="m-continue-deck",
        channel="feishu",
        channel_user_id="u1",
        session_id=session2.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=polluted_task,
    )
    response = await agent.process_message(user2)

    reloaded = await SessionManager(str(db_path)).get_or_create("feishu", "u1")
    active = reloaded.metadata.get("current_completion_contract") or reloaded.metadata.get("completed_completion_contract")
    assert active is not None
    if "task_text" in active:
        assert active["task_text"] == original_task
        assert active["required_skills"] == ["baoyu-design"]
        assert "继续生成_deck" not in active["artifact_dir"]
    else:
        assert active["source_message_id"] == "m-original-baoyu"
        assert active["task_fingerprint"] == agent._task_fingerprint(original_task)
    combined = "\n".join(msg.content for msg in reloaded.messages)
    assert "required_skills: pptx" not in combined
    assert "original_user_task: 继续生成 deck" not in combined
    assert response.content
    assert "任务未完成" not in response.content
    assert "继续生成" not in response.content
    assert "稍后" not in response.content
    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    pptx_path = pending_files[0]["file_path"]
    assert pptx_path.endswith(".pptx")
    assert os.path.exists(pptx_path)
    assert "继续生成_deck" not in pptx_path
    artifact_dir = os.path.dirname(pptx_path)
    assert os.path.exists(os.path.join(artifact_dir, "deck.html"))
    text = _pptx_text(pptx_path)
    assert "RAG" in text
    assert "企业知识库" in text
    assert "当前进展" not in text
    assert "下一步建议" not in text


@pytest.mark.asyncio
async def test_completed_skill_deliverable_continuation_does_not_restart_workflow(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.session import SessionManager

    skills_root = tmp_path / "skills"
    skill_dir = skills_root / "baoyu-design" / "skills" / "baoyu-design"
    _make_baoyu_skill_fixture(skill_dir)
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])

    db_path = tmp_path / "pyclaw.db"
    manager = SessionManager(str(db_path))
    await manager.init_db()
    session = await manager.get_or_create("feishu", "u1")
    original_task = "走完整的 baoyu design skill 做一个关于 RAG 企业知识库的幻灯片"
    fingerprint_agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    await manager.save_message(session, Message(
        id="m-original-completed-baoyu",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=original_task,
    ))
    await manager.save_message(session, Message(
        id="a-original-completed-baoyu",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.ASSISTANT,
        content="已生成并发送文件：RAG_企业知识库.pptx",
    ))
    session.metadata["completed_completion_contract"] = {
        "source_message_id": "m-original-completed-baoyu",
        "task_fingerprint": fingerprint_agent._task_fingerprint(original_task),
    }
    await manager.save_message(session, session.messages[-1])

    manager2 = SessionManager(str(db_path))
    session2 = await manager2.get_or_create("feishu", "u1")
    model = AsyncMock()
    model.chat.return_value = {"content": "之前已经完成并发送了。", "__tool_calls__": False}
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    tools.get_all_specs.return_value = []
    tools.execute_tool_calls = AsyncMock(return_value=[])
    agent = Agent(model, tools, manager2, max_iterations=1)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))

    response = await agent.process_message(Message(
        id="m-after-complete-continue",
        channel="feishu",
        channel_user_id="u1",
        session_id=session2.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="继续生成 deck",
    ))

    assert response.metadata.get("pending_files", []) == []
    reloaded = await SessionManager(str(db_path)).get_or_create("feishu", "u1")
    assert "current_completion_contract" not in reloaded.metadata
    assert not (tmp_path / "artifacts").exists()
    assert model.chat.await_count == 1

def _make_webpage_skill_fixture(skill_dir: Path) -> None:
    skill_dir.mkdir(parents=True, exist_ok=True)
    (skill_dir / "SKILL.md").write_text(
        "---\nname: webpage-coding\ndescription: polished webpage coding\n---\n"
        "# Webpage Coding\n\n"
        "Use this skill for high-polish webpages.\n"
        "Create `reference/brief.md` before coding and then generate `index.html`.\n"
        "The final page must be visually polished, responsive, interactive, and delivered as a file.\n",
        encoding="utf-8",
    )


def test_premium_html_auto_routes_to_webpage_coding_skill(tmp_path, monkeypatch):
    skills_root = tmp_path / "skills"
    _make_webpage_skill_fixture(skills_root / "webpage-coding")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    monkeypatch.setattr("pyclaw.core.agent._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    agent = Agent(AsyncMock(), tools, AsyncMock())

    required = agent._default_required_skills_for_file_deliverable(
        "做一个精美、高颜值的可视化教学网页，图解 GPT-3 结构和训练流程"
    )

    assert required == ("webpage-coding",)


@pytest.mark.asyncio
async def test_premium_html_webpage_skill_orchestration_mentions_brief_and_index(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager

    skills_root = tmp_path / "skills"
    _make_webpage_skill_fixture(skills_root / "webpage-coding")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    monkeypatch.setattr("pyclaw.core.agent._available_skills_dirs", lambda: [str(skills_root)])
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    sessions = AsyncMock()

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    agent = Agent(AsyncMock(), tools, sessions)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    session = Session(session_id="s-webpage-notice", user_id="u1", channel="feishu")
    session.messages.append(Message(
        id="m-webpage-notice",
        channel="feishu",
        channel_user_id="u1",
        user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个精美、高颜值的可视化教学网页，图解 RAG 结构和训练流程",
    ))

    contract = agent._infer_completion_contract(session)
    assert contract is not None
    assert contract.required_skills == ("webpage-coding",)
    await agent._ensure_skill_workflow_hydration_observations(session, contract)
    await agent._ensure_skill_workflow_orchestration_notice(session, contract)

    notice = "\n".join(
        msg.content for msg in session.messages
        if msg.role == MessageRole.USER and msg.metadata.get("skill_workflow_orchestration")
    )
    assert "reference/brief.md" in notice
    assert "index.html" in notice
    assert "webpage-coding" in notice
    assert "bounded_artifact_dir" in notice

    hydrated = [msg for msg in session.messages if msg.metadata.get("controller_skill_hydration")]
    assert any(msg.metadata.get("required_path") == "SKILL.md" for msg in hydrated)


@pytest.mark.asyncio
async def test_premium_html_webpage_skill_controller_completes_when_model_only_delegates(tmp_path, monkeypatch):
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.artifact_acceptance import ArtifactAcceptanceService

    skills_root = tmp_path / "skills"
    _make_webpage_skill_fixture(skills_root / "webpage-coding")
    monkeypatch.setattr("pyclaw.tools.skill_activation._available_skills_dirs", lambda: [str(skills_root)])
    monkeypatch.setattr("pyclaw.core.agent._available_skills_dirs", lambda: [str(skills_root)])

    model = AsyncMock()
    model.chat.return_value = {"content": "交给外部代码生成器了，正在做，稍后发你。", "__tool_calls__": False}
    tools = MagicMock()
    tools._tools = {}
    tools._static_tools = set()
    tools.skills_dirs = [skills_root]
    tools.get_all_specs.return_value = []
    tools.execute_tool_calls = AsyncMock(return_value=[])
    sessions = AsyncMock()
    session = Session(session_id="s-webpage-controller", user_id="u1", channel="feishu")

    async def save_message(sess, msg):
        if not any(existing.id == msg.id for existing in sess.messages):
            sess.messages.append(msg)

    sessions.save_message.side_effect = save_message
    sessions.get_or_create.return_value = session
    agent = Agent(model, tools, sessions, max_iterations=1)
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))

    response = await agent.process_message(Message(
        id="m-webpage-controller",
        channel="feishu",
        channel_user_id="u1",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="做一个精美、高颜值的可视化教学网页，图解 GPT-3 结构和训练流程",
    ))

    pending_files = response.metadata.get("pending_files", [])
    assert len(pending_files) == 1
    html_path = Path(pending_files[0]["file_path"])
    assert html_path.name == "index.html"
    assert html_path.exists()
    assert (html_path.parent / "reference" / "brief.md").exists()
    assert (html_path.parent / "skill-workflow-evidence.md").exists()
    assert "已生成并发送文件" in response.content
    assert "任务未完成" not in response.content
    assert "稍后" not in response.content
    assert "交给外部代码生成器" not in response.content

    # Completed contracts are cleared; reconstruct through acceptance on file content.
    html_text = html_path.read_text(encoding="utf-8")
    assert "GPT-3" in html_text or "gpt-3" in html_text.lower()
    assert "cdn.tailwindcss.com" in html_text
    assert "cdn.jsdelivr.net/npm/chart.js" in html_text
    assert "addEventListener" in html_text
    assert len(html_text) >= 14_000

    # Direct acceptance check with a fresh contract keeps the assertion focused
    # on artifact quality rather than session metadata lifecycle.
    from pyclaw.core.completion_contract import CompletionContract
    check_contract = CompletionContract(
        kind="file_deliverable",
        task_text="做一个精美、高颜值的可视化教学网页，图解 GPT-3 结构和训练流程",
        artifact_dir=str(html_path.parent),
    )
    result = ArtifactAcceptanceService().evaluate(check_contract, [{"file_path": str(html_path)}])
    assert result.accepted, result.reasons


def test_cron_session_disables_completion_contracts_for_text_notification(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    session = Session(
        session_id="cron_808c118e_1",
        user_id="job_808c118e",
        channel="cron",
        metadata={"cron_job_id": "808c118e", "disable_completion_contracts": True},
    )
    cron_prompt = (
        "【定时任务执行 - 请只执行以下任务，不要创建新任务，不要回复关于任务本身的说明】\n"
        "当前执行时间：2026-06-13 21:00:00 CST。"
        "硬性限制：优先使用少量高可信来源；多页面读取优先用 web_extract 一次读取；"
        "最终回复不得提及工具调用。\n\n"
        "首先执行osascript -e 'display notification \"现在是21点\"'，"
        "然后在当前会话推送提醒：「现在是21点啦，建议错峰出行，避开下班高峰路段哦~」，"
        "执行完成后记录日志到~/.pyclaw/cron_history/21_travel_reminder.log"
    )
    session.messages.append(Message(
        id="cron-808c118e",
        channel="cron",
        channel_user_id="job_808c118e",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content=cron_prompt,
        metadata={"disable_completion_contracts": True},
    ))

    assert agent._infer_completion_contract(session) is None
    assert not (tmp_path / "artifacts").exists()
    assert "current_completion_contract" not in session.metadata


def test_cron_session_clears_stale_completion_contract_metadata(tmp_path):
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.completion_contract import CompletionContract

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    stale = CompletionContract(
        kind="file_deliverable",
        task_text="做一个网页",
        artifact_dir=str(tmp_path / "old"),
        source_message_id="old",
        task_fingerprint="old",
    ).to_metadata()
    session = Session(
        session_id="cron_808c118e_2",
        user_id="job_808c118e",
        channel="cron",
        metadata={
            "cron_job_id": "808c118e",
            "disable_completion_contracts": True,
            "current_completion_contract": stale,
            "last_incomplete_completion_contract": stale,
            "last_explicit_skill_completion_contract": stale,
        },
    )
    session.messages.append(Message(
        id="cron-808c118e",
        channel="cron",
        channel_user_id="job_808c118e",
        session_id=session.session_id,
        type=MessageType.TEXT,
        role=MessageRole.USER,
        content="【定时任务执行】当前会话推送提醒，记录日志到~/.pyclaw/cron_history/21_travel_reminder.log",
    ))

    assert agent._infer_completion_contract(session) is None
    assert "current_completion_contract" not in session.metadata
    assert "last_incomplete_completion_contract" not in session.metadata
    assert "last_explicit_skill_completion_contract" not in session.metadata


def test_operational_confirmation_does_not_resume_stale_skill_completion_contract(tmp_path):
    """A generic approval for a Pod/image update must not revive old file gates."""
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.completion_contract import CompletionContract

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    stale = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 的幻灯片",
        artifact_dir=str(tmp_path / "old-deck"),
        required_evidence=("file_created", "send_file_to_user"),
        source_message_id="old-skill-task",
        task_fingerprint="old-skill-fingerprint",
        required_skills=("baoyu-design",),
    ).to_metadata()
    session = Session(
        session_id="s-pod-update-confirmation",
        user_id="u1",
        channel="feishu",
        metadata={
            "current_completion_contract": stale,
            "last_incomplete_completion_contract": stale,
            "last_explicit_skill_completion_contract": stale,
        },
    )
    session.messages.extend([
        Message(
            id="m-pod-update",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="这个pod 7652273671583177522镜像给我升级成cr-aic-cn-beijing.cr.volces.com/hhl/aosp13:xr20260721",
        ),
        Message(
            id="m-confirm-pod-update",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="确认",
        ),
    ])

    assert agent._infer_completion_contract(session) is None
    assert "current_completion_contract" not in session.metadata
    assert "last_incomplete_completion_contract" not in session.metadata
    assert session.metadata["last_explicit_skill_completion_contract"] == stale


def test_operational_success_final_not_rewritten_by_stale_skill_gate(tmp_path):
    """Successful non-artifact tool work should be reported as done, not skill-failed."""
    from pyclaw.core.artifacts import ArtifactManager
    from pyclaw.core.completion_contract import CompletionContract

    agent = Agent(AsyncMock(), MagicMock(), MagicMock())
    agent.artifacts = ArtifactManager(root=str(tmp_path / "artifacts"))
    stale = CompletionContract(
        kind="file_deliverable",
        task_text="走完整的 baoyu design skill 做一个关于 RAG 的幻灯片",
        artifact_dir=str(tmp_path / "old-deck"),
        required_evidence=("file_created", "send_file_to_user"),
        source_message_id="old-skill-task",
        task_fingerprint="old-skill-fingerprint",
        required_skills=("baoyu-design",),
    ).to_metadata()
    session = Session(
        session_id="s-pod-update-final",
        user_id="u1",
        channel="feishu",
        metadata={
            "current_completion_contract": stale,
            "last_incomplete_completion_contract": stale,
            "last_explicit_skill_completion_contract": stale,
        },
    )
    session.messages.extend([
        Message(
            id="m-pod-update",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="这个pod 7652273671583177522镜像给我升级成cr-aic-cn-beijing.cr.volces.com/hhl/aosp13:xr20260721",
        ),
        Message(
            id="m-confirm-pod-update",
            channel="feishu",
            channel_user_id="u1",
            user_id="u1",
            session_id=session.session_id,
            type=MessageType.TEXT,
            role=MessageRole.USER,
            content="确认",
        ),
    ])
    pending_files: list[dict[str, str]] = []

    final = agent._prepare_completion_contract_final_content(
        session=session,
        content="Pod 镜像已更新成功。目标镜像：cr-aic-cn-beijing.cr.volces.com/hhl/aosp13:xr20260721。",
        pending_files=pending_files,
    )

    assert "已更新成功" in final
    assert "skill 工作流验收" not in final
    assert pending_files == []
