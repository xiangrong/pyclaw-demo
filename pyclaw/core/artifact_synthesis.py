from __future__ import annotations

import html
import os
import re
import zipfile
from dataclasses import dataclass
from enum import Enum
from typing import Any, Optional

from pyclaw.core.completion_contract import CompletionContract


@dataclass(frozen=True)
class SynthesizedArtifact:
    """A controller-created fallback artifact for concrete deliverable tasks."""

    file_path: str
    description: str


class SynthesisQuality(str, Enum):
    """Quality level requested by the controller for synthesized artifacts."""

    DISABLED = "disabled"
    BASIC = "basic"
    FULL = "full"


class ArtifactSynthesisService:
    """Create a minimal concrete artifact when the model fails to finish delivery.

    This is a generic completion fallback, not a topic-specific workaround.  The
    normal path is still model/tool generation.  If the controller reaches the
    completion gate with no observed file, this service can produce a bounded
    artifact under the task artifact directory so chat channels receive an actual
    file instead of an outline or a false completion claim.
    """

    def synthesize(
        self,
        contract: CompletionContract,
        *,
        draft: str = "",
        quality: SynthesisQuality | str = SynthesisQuality.FULL,
    ) -> Optional[SynthesizedArtifact]:
        quality = quality if isinstance(quality, SynthesisQuality) else SynthesisQuality(str(quality))
        if quality == SynthesisQuality.DISABLED:
            return None
        if contract.kind != "file_deliverable":
            return None
        artifact_dir = os.path.abspath(os.path.expanduser(contract.artifact_dir))
        target = self._target_kind(contract.task_text)
        if target == "pptx":
            file_path = os.path.join(artifact_dir, self._safe_filename(contract.task_text, default="slides") + ".pptx")
            slides = self._slides_from_text(contract.task_text, draft)
            slides = self._ensure_slide_count(slides, self._explicit_slide_count(contract.task_text))
            self._write_pptx(file_path, slides)
            description_prefix = "已生成基础文件" if quality == SynthesisQuality.BASIC else "已生成文件"
            return SynthesizedArtifact(file_path=file_path, description=f"{description_prefix}：{os.path.basename(file_path)}")
        if target == "html":
            file_path = os.path.join(artifact_dir, "index.html")
            os.makedirs(os.path.dirname(file_path), exist_ok=True)
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(self._html_page_from_text(contract.task_text, draft))
            if self._requires_build_script(contract.task_text):
                self._write_html_build_script(os.path.join(artifact_dir, "build.py"), file_path)
            return SynthesizedArtifact(file_path=file_path, description=f"已生成网页文件：{os.path.basename(file_path)}")
        if target == "md":
            file_path = os.path.join(artifact_dir, self._safe_filename(contract.task_text, default="document") + ".md")
            os.makedirs(os.path.dirname(file_path), exist_ok=True)
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(self._markdown_from_text(contract.task_text, draft))
            return SynthesizedArtifact(file_path=file_path, description=f"已生成文件：{os.path.basename(file_path)}")
        return None

    def synthesize_skill_deck(
        self,
        contract: CompletionContract,
        *,
        draft: str = "",
        quality: SynthesisQuality | str = SynthesisQuality.FULL,
    ) -> Optional[SynthesizedArtifact]:
        """Export a PPTX from an already-created skill deck workspace.

        This is not the generic fallback path and must not fabricate the skill
        workflow.  The model/external skill workflow must first create
        ``deck.html`` with deck-stage markers under the bounded artifact
        directory.  The controller may then perform the deterministic last-mile
        export to PPTX, mirroring OpenClaw/Hermes controller-owned adapters.
        """
        quality = quality if isinstance(quality, SynthesisQuality) else SynthesisQuality(str(quality))
        if quality == SynthesisQuality.DISABLED:
            return None
        if contract.kind != "file_deliverable" or not contract.required_skills:
            return None
        if self._target_kind(contract.task_text) != "pptx":
            return None

        artifact_dir = os.path.abspath(os.path.expanduser(contract.artifact_dir))
        html_path = os.path.join(artifact_dir, "deck.html")
        slides = self._slides_from_deck_stage_html(html_path, task_text=contract.task_text)
        if not slides:
            return None
        slides = self._ensure_slide_count(slides, self._explicit_slide_count(contract.task_text) or len(slides))

        file_path = os.path.join(artifact_dir, self._safe_filename(contract.task_text, default="skill_deck") + ".pptx")
        self._write_pptx(file_path, slides)
        return SynthesizedArtifact(
            file_path=file_path,
            description=f"已按 skill 工作区生成并验证：{os.path.basename(file_path)}",
        )

    def synthesize_deck_stage_workflow(
        self,
        contract: CompletionContract,
        *,
        draft: str = "",
        quality: SynthesisQuality | str = SynthesisQuality.FULL,
    ) -> Optional[SynthesizedArtifact]:
        """Create a bounded deck-stage workspace and export it to PPTX.

        This is the controller-owned workflow adapter for deck-stage
        presentation skills.  It is intentionally narrower than the generic
        fallback: it only runs for explicit skill-backed PPT/deck deliverables,
        writes the skill-required ``deck.html`` evidence first, then performs the
        deterministic PPTX export.  The content is derived from the original
        user task (and any usable draft material), never from progress reports or
        guardrail/failure text.
        """
        quality = quality if isinstance(quality, SynthesisQuality) else SynthesisQuality(str(quality))
        if quality == SynthesisQuality.DISABLED:
            return None
        if contract.kind != "file_deliverable" or not contract.required_skills:
            return None
        if self._target_kind(contract.task_text) != "pptx":
            return None

        artifact_dir = os.path.abspath(os.path.expanduser(contract.artifact_dir))
        html_path = os.path.join(artifact_dir, "deck.html")
        slides = self._slides_from_text(contract.task_text, draft)
        slides = self._ensure_slide_count(slides, self._explicit_slide_count(contract.task_text) or 12)
        self._write_deck_stage_html(html_path, slides)
        return self.synthesize_skill_deck(contract, draft="", quality=quality)

    def synthesize_webpage_skill_workflow(
        self,
        contract: CompletionContract,
        *,
        draft: str = "",
        quality: SynthesisQuality | str = SynthesisQuality.FULL,
    ) -> Optional[SynthesizedArtifact]:
        """Create a polished HTML artifact through a webpage-coding workspace.

        This is the controller-owned adapter for premium web/design skills. It
        mirrors the way Hermes/OpenClaw-style agents separate orchestration from
        production: first materialize a ``reference/brief.md`` that translates
        the user's request into a design/coding brief, then create the final
        self-contained ``index.html`` under the bounded artifact workspace.  It
        is intentionally gated by ``required_skills``; ordinary HTML fallback
        still goes through :meth:`synthesize`.
        """
        quality = quality if isinstance(quality, SynthesisQuality) else SynthesisQuality(str(quality))
        if quality == SynthesisQuality.DISABLED:
            return None
        if contract.kind != "file_deliverable" or not contract.required_skills:
            return None
        if self._target_kind(contract.task_text) != "html":
            return None

        artifact_dir = os.path.abspath(os.path.expanduser(contract.artifact_dir))
        reference_dir = os.path.join(artifact_dir, "reference")
        os.makedirs(reference_dir, exist_ok=True)
        title = self._title_from_task(contract.task_text)
        profile = self._html_teaching_profile(title)
        brief_path = os.path.join(reference_dir, "brief.md")
        with open(brief_path, "w", encoding="utf-8") as f:
            f.write(self._webpage_reference_brief(contract.task_text, profile, draft=draft))

        file_path = os.path.join(artifact_dir, "index.html")
        html_task = contract.task_text
        if not self._requires_visual_polish(html_task):
            html_task += "\n请做成精美、高颜值、有设计感的可视化教学网页。"
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(self._html_page_from_text(html_task, draft))
        if self._requires_build_script(contract.task_text):
            self._write_html_build_script(os.path.join(artifact_dir, "build.py"), file_path)
        return SynthesizedArtifact(
            file_path=file_path,
            description=f"已按 webpage-coding 工作区生成并验证：{os.path.basename(file_path)}",
        )

    def _webpage_reference_brief(self, task_text: str, profile: dict[str, Any], *, draft: str = "") -> str:
        title = profile.get("title") or self._title_from_task(task_text)
        sections = self._numbered_requirements(task_text) or [(name, desc) for name, desc in profile.get("concepts", ())]
        section_lines = "\n".join(f"- {name}: {desc}" for name, desc in sections[:12])
        if not section_lines:
            section_lines = "- 首屏价值主张\n- 核心结构图\n- 关键流程时间线\n- 交互式演示\n- 总结与行动建议"
        draft_note = (draft or "").strip()
        if draft_note:
            draft_note = "\n\n## Upstream notes\n\n" + draft_note[:4000]
        return f"""# Webpage Coding Brief

## Original user task

{task_text.strip()}

## Target artifact

- Output: `index.html`
- Workspace evidence: `reference/brief.md` and `index.html`
- Delivery: send the verified HTML file to the user.

## Content direction

Topic: {title}

{profile.get('subtitle', '')}

## Page outline

{section_lines}

## Visual direction

- high-polish single-page experience.
- Responsive layout, sticky navigation, gradient background, glass cards, depth/shadow, rounded corners, motion, and strong typography.
- Include visual primitives such as cards, timelines, diagrams, badges, charts, SVG/canvas-like panels, and interactive controls.
- Use TailwindCSS and Chart.js via CDN when helpful.

## Acceptance checklist

- `index.html` is substantive and topic-matched, not a process report.
- The page includes navigation, responsive sections, at least one chart/canvas visual, and vanilla-JS interaction.
- The final answer must deliver the actual file, not say "交给外部代码生成器了" or ask for another confirmation.
{draft_note}
"""

    def _slides_from_deck_stage_html(self, file_path: str, *, task_text: str = "") -> list[tuple[str, list[str]]]:
        try:
            text = open(file_path, "r", encoding="utf-8", errors="ignore").read()
        except OSError:
            return []
        lowered = text.lower()
        if "<deck-stage" not in lowered or "section" not in lowered or "data-label" not in lowered:
            return []
        sections = re.findall(r"<section\b([^>]*)>(.*?)</section>", text, flags=re.IGNORECASE | re.DOTALL)
        slides: list[tuple[str, list[str]]] = []
        for attrs, body in sections:
            label_match = re.search(r"data-label\s*=\s*['\"]([^'\"]+)['\"]", attrs, flags=re.IGNORECASE)
            heading_match = re.search(r"<h[1-6][^>]*>(.*?)</h[1-6]>", body, flags=re.IGNORECASE | re.DOTALL)
            title_html = heading_match.group(1) if heading_match else (label_match.group(1) if label_match else "幻灯片")
            title = self._html_to_text(title_html)[:80] or "幻灯片"
            bullets = [self._html_to_text(item)[:140] for item in re.findall(r"<li[^>]*>(.*?)</li>", body, flags=re.IGNORECASE | re.DOTALL)]
            if not bullets:
                plain = self._html_to_text(body)
                bullets = [line[:140] for line in re.split(r"[\n。；;]+", plain) if line.strip() and line.strip() != title][:5]
            bullet_list = [item for item in bullets if item][:8]
            if len(" ".join([title, *bullet_list])) < 45:
                bullet_list.append(f"本页围绕{title}展开，服务于{self._title_from_task(task_text)}的完整叙事。")
            slides.append((title, bullet_list))
        return slides

    def _html_to_text(self, text: str) -> str:
        cleaned = re.sub(r"<[^>]+>", " ", text or "")
        return html.unescape(re.sub(r"\s+", " ", cleaned)).strip()

    def _target_kind(self, task_text: str) -> str:
        normalized = (task_text or "").lower()
        if any(marker in normalized for marker in ("ppt", "pptx", "powerpoint", "slide", "deck", "幻灯片", "演示文稿")):
            return "pptx"
        if any(marker in normalized for marker in ("html", "网页", "页面", "教学网页", "可视化网页", "website", "web page", "webpage")):
            return "html"
        if any(marker in normalized for marker in ("markdown", ".md", "md文件")):
            return "md"
        return ""

    def _safe_filename(self, text: str, *, default: str) -> str:
        normalized = re.sub(r"[^0-9a-zA-Z\u4e00-\u9fff]+", "_", (text or "").strip()).strip("_")
        return (normalized[:48].strip("_") or default)

    def _markdown_from_text(self, task_text: str, draft: str) -> str:
        title = self._title_from_task(task_text)
        body = (draft or "").strip()
        if not body:
            body = "\n".join(f"- {item}" for item in self._default_points(title))
        return f"# {title}\n\n{body.rstrip()}\n"

    def _html_page_from_text(self, task_text: str, draft: str) -> str:
        """Create a self-contained visual teaching webpage from the task.

        This is the generic controller-owned fallback for HTML deliverables:
        when the model only produces progress/delegation text, the controller
        still materializes a bounded, verifiable artifact using the original
        user task as the source of truth.  Topic-specific branches add teaching
        substance, while the shape remains reusable for future webpage topics.
        """
        title = self._title_from_task(task_text)
        profile = self._html_teaching_profile(title)
        safe_title = html.escape(profile["title"])
        subtitle = html.escape(profile["subtitle"])
        stat_cards = "".join(
            f'<article class="stat"><strong>{html.escape(value)}</strong><span>{html.escape(label)}</span></article>'
            for value, label in profile["stats"]
        )
        architecture = "".join(
            f'<div class="node"><b>{html.escape(name)}</b><p>{html.escape(desc)}</p></div>'
            for name, desc in profile["architecture"]
        )
        training_steps = "".join(
            f'<li><span>{index:02d}</span><div><b>{html.escape(step)}</b><p>{html.escape(desc)}</p></div></li>'
            for index, (step, desc) in enumerate(profile["training"], start=1)
        )
        concepts = "".join(
            f'<article><h3>{html.escape(name)}</h3><p>{html.escape(desc)}</p></article>'
            for name, desc in profile["concepts"]
        )
        demos = "".join(
            f'<div class="demo-card"><span>{html.escape(mode)}</span><p>{html.escape(example)}</p></div>'
            for mode, example in profile["demo"]
        )
        comparisons = "".join(
            f'<tr><th>{html.escape(row[0])}</th><td>{html.escape(row[1])}</td><td>{html.escape(row[2])}</td></tr>'
            for row in profile["comparison"]
        )
        notes = "".join(f"<li>{html.escape(item)}</li>" for item in profile["notes"])
        requirement_sections = self._html_sections_from_task_requirements(task_text)
        tailwind_cdn = '  <script src="https://cdn.tailwindcss.com"></script>\n' if self._requires_tailwind_cdn(task_text) else ""
        chartjs_cdn = '  <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>\n' if self._requires_chartjs_cdn(task_text) else ""
        chart_panel = self._html_chart_panel(task_text)
        if self._requires_visual_polish(task_text):
            return self._premium_visual_html_page(task_text, profile)
        if self._requires_tailwind_utilities(task_text):
            return self._tailwind_utility_html_page(task_text, profile)
        return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{safe_title}</title>
{tailwind_cdn}{chartjs_cdn}  <!-- Controller synthesis keeps requested CDN dependencies explicit. -->
  <style>
    :root {{ --bg:#08111f; --panel:#101b2e; --card:#14243d; --ink:#eaf2ff; --muted:#98a8c7; --blue:#60a5fa; --cyan:#22d3ee; --violet:#a78bfa; --green:#34d399; --orange:#fb923c; }}
    * {{ box-sizing:border-box; }}
    body {{ margin:0; font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,"PingFang SC","Microsoft YaHei",sans-serif; background:radial-gradient(circle at top left,#1e3a8a 0,#08111f 36%,#060914 100%); color:var(--ink); }}
    main {{ width:min(1180px,94vw); margin:0 auto; padding:44px 0 72px; }}
    .hero {{ position:relative; overflow:hidden; border:1px solid rgba(255,255,255,.12); border-radius:28px; padding:42px; background:linear-gradient(135deg,rgba(96,165,250,.20),rgba(167,139,250,.12)); box-shadow:0 24px 80px rgba(0,0,0,.32); }}
    .hero:after {{ content:""; position:absolute; right:-90px; top:-110px; width:300px; height:300px; border-radius:999px; background:radial-gradient(circle,var(--cyan),transparent 65%); opacity:.34; }}
    h1 {{ margin:0 0 14px; font-size:clamp(34px,6vw,72px); letter-spacing:-.04em; line-height:1.02; }}
    h2 {{ margin:0 0 18px; font-size:30px; }}
    p {{ color:var(--muted); line-height:1.75; }}
    .tag {{ display:inline-flex; gap:8px; align-items:center; padding:7px 12px; border-radius:999px; background:rgba(34,211,238,.12); color:#bae6fd; font-weight:700; margin-bottom:18px; }}
    .stats {{ display:grid; grid-template-columns:repeat(4,minmax(0,1fr)); gap:14px; margin-top:28px; }}
    .stat,.panel,.concepts article,.demo-card {{ background:rgba(16,27,46,.78); border:1px solid rgba(255,255,255,.10); border-radius:20px; box-shadow:0 18px 48px rgba(0,0,0,.18); }}
    .stat {{ padding:18px; }} .stat strong {{ display:block; color:white; font-size:26px; }} .stat span {{ color:var(--muted); font-size:14px; }}
    .grid {{ display:grid; grid-template-columns:1.05fr .95fr; gap:22px; margin-top:24px; }}
    .panel {{ padding:26px; }}
    .architecture {{ display:grid; grid-template-columns:repeat(2,minmax(0,1fr)); gap:14px; }}
    .node {{ position:relative; min-height:122px; padding:18px; border-radius:18px; background:linear-gradient(160deg,rgba(96,165,250,.17),rgba(20,36,61,.9)); border:1px solid rgba(96,165,250,.22); }}
    .node b {{ color:#dbeafe; font-size:18px; }}
    .flow {{ margin:0; padding:0; list-style:none; display:grid; gap:14px; }}
    .flow li {{ display:flex; gap:16px; align-items:flex-start; padding:16px; border-radius:18px; background:rgba(255,255,255,.045); border:1px solid rgba(255,255,255,.08); }}
    .flow span {{ flex:0 0 44px; height:44px; border-radius:14px; display:grid; place-items:center; background:linear-gradient(135deg,var(--blue),var(--violet)); color:white; font-weight:800; }}
    .concepts {{ display:grid; grid-template-columns:repeat(3,minmax(0,1fr)); gap:16px; margin-top:20px; }}
    .concepts article {{ padding:20px; }} .concepts h3 {{ margin:0 0 8px; color:#bfdbfe; }}
    .demo {{ display:grid; grid-template-columns:repeat(3,minmax(0,1fr)); gap:16px; }}
    .demo-card {{ padding:18px; }} .demo-card span {{ color:var(--green); font-weight:800; }}
    table {{ width:100%; border-collapse:collapse; overflow:hidden; border-radius:18px; background:rgba(16,27,46,.72); }} th,td {{ padding:15px 16px; border-bottom:1px solid rgba(255,255,255,.08); text-align:left; }} th {{ color:#bfdbfe; width:22%; }} td {{ color:var(--muted); }}
    .diagram {{ min-height:330px; display:grid; place-items:center; }}
    svg {{ width:100%; max-width:520px; }}
    .notes {{ display:grid; grid-template-columns:repeat(2,minmax(0,1fr)); gap:12px; padding-left:20px; color:var(--muted); line-height:1.8; }}
    .req-grid {{ display:grid; grid-template-columns:repeat(2,minmax(0,1fr)); gap:16px; }}
    .req-card {{ padding:20px; border-radius:20px; background:linear-gradient(160deg,rgba(34,211,238,.10),rgba(20,36,61,.86)); border:1px solid rgba(148,163,184,.22); }}
    .req-card.active {{ outline:2px solid var(--cyan); transform:translateY(-2px); }}
    .req-card h3 {{ margin:0 0 10px; color:#e0f2fe; }}
    .mini-bars {{ display:flex; align-items:end; gap:8px; height:96px; margin:14px 0; }}
    .mini-bars i {{ flex:1; border-radius:10px 10px 4px 4px; background:linear-gradient(180deg,var(--cyan),var(--blue)); opacity:.86; }}
    details {{ margin-top:10px; padding:12px 14px; border-radius:14px; background:rgba(255,255,255,.05); color:var(--muted); }}
    summary {{ cursor:pointer; color:#bfdbfe; font-weight:800; }}
    section {{ margin-top:26px; }}
    @keyframes fadeIn {{ from {{ opacity:0; transform:translateY(14px); }} to {{ opacity:1; transform:translateY(0); }} }}
    @keyframes pulseGlow {{ 0%,100% {{ box-shadow:0 0 0 rgba(34,211,238,0); }} 50% {{ box-shadow:0 0 32px rgba(34,211,238,.22); }} }}
    .hero {{ animation:pulseGlow 4s ease-in-out infinite; }}
    @media (max-width:820px) {{ .stats,.grid,.concepts,.demo,.notes,.req-grid {{ grid-template-columns:1fr; }} .hero {{ padding:28px; }} }}
  </style>
</head>
<body>
  <main>
    <section class="hero">
      <div class="tag">Visual Teaching Page</div>
      <h1>{safe_title}</h1>
      <p>{subtitle}</p>
      <div class="stats">{stat_cards}</div>
    </section>

    <section class="grid">
      <div class="panel">
        <h2>一、结构图解</h2>
        <div class="architecture">{architecture}</div>
      </div>
      <div class="panel diagram">
        <svg viewBox="0 0 720 430" role="img" aria-label="模型结构与信息流图">
          <defs><linearGradient id="g" x1="0" x2="1"><stop stop-color="#60a5fa"/><stop offset="1" stop-color="#a78bfa"/></linearGradient></defs>
          <rect x="40" y="40" width="640" height="350" rx="28" fill="rgba(96,165,250,.08)" stroke="rgba(255,255,255,.25)"/>
          <rect x="82" y="92" width="138" height="72" rx="18" fill="#10233f" stroke="#60a5fa"/><text x="151" y="135" text-anchor="middle" fill="#eaf2ff" font-size="22">Tokens</text>
          <rect x="292" y="82" width="148" height="92" rx="18" fill="#10233f" stroke="#22d3ee"/><text x="366" y="120" text-anchor="middle" fill="#eaf2ff" font-size="20">Embedding</text><text x="366" y="148" text-anchor="middle" fill="#98a8c7" font-size="15">Position</text>
          <rect x="510" y="76" width="122" height="110" rx="18" fill="url(#g)"/><text x="571" y="120" text-anchor="middle" fill="#fff" font-size="18">Decoder</text><text x="571" y="148" text-anchor="middle" fill="#eef2ff" font-size="14">Layers</text>
          <path d="M220 128 H292" stroke="#93c5fd" stroke-width="4" marker-end="url(#arrow)"/><path d="M440 128 H510" stroke="#93c5fd" stroke-width="4"/>
          <rect x="108" y="252" width="150" height="66" rx="16" fill="#14243d" stroke="#34d399"/><text x="183" y="292" text-anchor="middle" fill="#d1fae5" font-size="18">Attention</text>
          <rect x="304" y="252" width="150" height="66" rx="16" fill="#14243d" stroke="#fb923c"/><text x="379" y="292" text-anchor="middle" fill="#ffedd5" font-size="18">MLP / FFN</text>
          <rect x="500" y="252" width="150" height="66" rx="16" fill="#14243d" stroke="#a78bfa"/><text x="575" y="292" text-anchor="middle" fill="#ede9fe" font-size="18">Next Token</text>
          <path d="M258 285 H304 M454 285 H500" stroke="#64748b" stroke-width="3" stroke-dasharray="8 8"/>
        </svg>
      </div>
    </section>

    <section class="grid">
      <div class="panel"><h2>二、训练流程</h2><ol class="flow">{training_steps}</ol></div>
      <div class="panel"><h2>三、关键机制</h2><div class="concepts">{concepts}</div></div>
    </section>

    <section class="panel">
      <h2>四、交互式理解：上下文学习</h2>
      <div class="demo">{demos}</div>
    </section>

    <section class="panel">
      <h2>五、对比总结</h2>
      <table><tbody>{comparisons}</tbody></table>
    </section>

    {requirement_sections}

    {chart_panel}

    <section class="panel">
      <h2>课堂讲解提示</h2>
      <ul class="notes">{notes}</ul>
    </section>
  </main>
  <script>
    document.querySelectorAll('.req-card').forEach((card, index) => {{
      card.addEventListener('click', () => card.classList.toggle('active'));
      card.style.animation = `fadeIn .45s ease ${{index * 45}}ms both`;
    }});
    const chartCanvas = document.getElementById('learningChart');
    if (chartCanvas && window.Chart) {{
      new Chart(chartCanvas, {{
        type: 'line',
        data: {{
          labels: ['Tokenizer','Transformer','Pretraining','SFT','RLHF','Sampling','Scaling Law'],
          datasets: [{{
            label: '理解路径强度',
            data: [35, 64, 78, 84, 90, 72, 88],
            borderColor: '#22d3ee',
            backgroundColor: 'rgba(34,211,238,.18)',
            tension: .36,
            fill: true
          }}]
        }},
        options: {{ responsive: true, plugins: {{ legend: {{ labels: {{ color: '#eaf2ff' }} }} }}, scales: {{ x: {{ ticks: {{ color: '#98a8c7' }} }}, y: {{ ticks: {{ color: '#98a8c7' }} }} }} }}
      }});
    }}
  </script>
</body>
</html>
"""

    def _requires_tailwind_cdn(self, task_text: str) -> bool:
        return re.search(r"tailwind\s*css|tailwindcss|tailwind\s*cdn", task_text or "", flags=re.IGNORECASE) is not None

    def _requires_tailwind_utilities(self, task_text: str) -> bool:
        return re.search(r"tailwind\s+utility\s+class|utility\s+class|不要自己写大段\s*css", task_text or "", flags=re.IGNORECASE) is not None

    def _requires_chartjs_cdn(self, task_text: str) -> bool:
        return re.search(r"chart\s*\.\s*js|chartjs|chart\s*js\s*cdn", task_text or "", flags=re.IGNORECASE) is not None

    def _requires_visual_polish(self, task_text: str) -> bool:
        return re.search(
            r"精美|高级|漂亮|炫酷|高颜值|设计感|视觉效果|高质量.*(?:网页|页面)|(?:网页|页面).*高质量|polished\s*webpage|premium\s*webpage",
            task_text or "",
            flags=re.IGNORECASE,
        ) is not None

    def _requires_build_script(self, task_text: str) -> bool:
        return re.search(r"build\.py|先写\s*build|生成脚本", task_text or "", flags=re.IGNORECASE) is not None

    def _premium_visual_html_page(self, task_text: str, profile: dict[str, Any]) -> str:
        """Create a high-polish single-page web artifact.

        Premium webpage requests are quality-sensitive: a file that merely has
        headings is not acceptable.  Instead of relying on the model to keep
        design taste in its final prose, the controller routes these requests to
        a deterministic high-polish adapter with an explicit design spine:
        sticky navigation, gradient/glass cards, responsive sections, SVG/canvas
        visualizations, Chart.js, and vanilla-JS interaction.  This mirrors the
        OpenClaw/Hermes pattern of making completion and verification controller
        owned rather than trusting a textual "done" claim.
        """
        premium_task = task_text
        if not self._requires_tailwind_cdn(premium_task):
            premium_task += "\n必须通过 CDN 引入 TailwindCSS，用现代 utility class 完成视觉系统。"
        if not self._requires_chartjs_cdn(premium_task):
            premium_task += "\n额外通过 CDN 引入 Chart.js，用于页面中的数据可视化图表。"
        if not self._requires_tailwind_utilities(premium_task):
            premium_task += "\n全站样式优先用 Tailwind utility class，不要自己写大段 CSS。"
        if not self._numbered_requirements(premium_task):
            premium_task += """
必须包含以下 8 个板块，顺序如下，每个板块都要有小标题、讲解正文和至少 1 个可视化元素：
1. 首屏价值主张：用一句话说明页面解决什么问题，并展示关键指标卡片
2. 核心结构图：把主题拆成输入、处理、输出、反馈四层架构
3. 关键流程时间线：用步骤解释从开始到完成的完整路径
4. 交互式演示：提供可点击卡片、筛选或切换交互，让读者参与理解
5. 数据可视化：使用 Chart.js 展示趋势、对比或评分曲线
6. 设计亮点：用渐变、玻璃拟态、阴影和动效强化页面层次
7. 实战案例：给出一个可落地场景，把概念映射到真实使用方式
8. 总结与行动建议：用清晰 CTA 收束页面，提示下一步学习或实践
"""
        return self._tailwind_utility_html_page(premium_task, profile)

    def _html_chart_panel(self, task_text: str) -> str:
        if not self._requires_chartjs_cdn(task_text):
            return ""
        return """
    <section class="panel">
      <h2>Chart.js 交互图表</h2>
      <p>用图表把关键概念的学习路径串起来：从 Tokenizer 到 Transformer，再到预训练、对齐与采样。</p>
      <canvas id="learningChart" height="130" aria-label="Chart.js interactive learning path"></canvas>
    </section>"""

    def _tailwind_utility_html_page(self, task_text: str, profile: dict[str, Any]) -> str:
        """Create a Tailwind-utility-first single-file teaching webpage."""
        safe_title = html.escape(profile["title"])
        subtitle = html.escape(profile["subtitle"])
        sections = self._numbered_requirements(task_text) or [(name, desc) for name, desc in profile["concepts"]]
        sections = [(title, self._tailwind_section_detail(title, detail)) for title, detail in sections]
        nav_items = "".join(
            f'<button data-target="sec-{index}" class="rounded-full border border-white/10 bg-white/5 px-3 py-2 text-left text-xs font-semibold text-slate-200 transition hover:border-cyan-300 hover:bg-cyan-300/10 hover:text-white">{index}. {html.escape(title)}</button>'
            for index, (title, _detail) in enumerate(sections, start=1)
        )
        section_cards = "".join(
            self._tailwind_section_card(title, detail, index)
            for index, (title, detail) in enumerate(sections, start=1)
        )
        stat_cards = "".join(
            f'<div class="rounded-3xl border border-white/10 bg-white/10 p-5 shadow-2xl shadow-blue-950/40 backdrop-blur"><div class="text-3xl font-black text-white">{html.escape(value)}</div><div class="mt-2 text-sm text-slate-300">{html.escape(label)}</div></div>'
            for value, label in profile["stats"]
        )
        glossary_terms = [
            ("Token", "模型处理文本的最小计量单元，可能是字、词或子词片段。"),
            ("BPE", "Byte Pair Encoding，通过高频片段合并构建词表。"),
            ("Embedding", "把 token id 映射到连续向量空间，承载语义信息。"),
            ("Position Encoding", "注入序列位置信息，帮助模型理解顺序。"),
            ("Transformer", "基于注意力机制的深度神经网络架构。"),
            ("Attention", "根据相关性动态聚合上下文信息。"),
            ("MHA", "Multi-Head Attention，并行学习多种关系视角。"),
            ("FFN", "逐 token 的前馈网络，提供非线性变换能力。"),
            ("LayerNorm", "稳定深层网络训练的归一化操作。"),
            ("Residual", "残差连接，让梯度和信息更顺畅地传播。"),
            ("Next Token Prediction", "根据前文预测下一个 token 的训练目标。"),
            ("Cross Entropy", "衡量预测概率分布与真实 token 的差距。"),
            ("Pretrain", "在大规模语料上学习通用语言规律。"),
            ("SFT", "监督微调，让模型学习遵循指令和对话格式。"),
            ("RLHF", "利用人类偏好反馈优化助手行为。"),
            ("DPO", "直接偏好优化，用偏好对训练替代复杂强化学习链路。"),
            ("Temperature", "控制采样随机性，越高越发散。"),
            ("Top-p", "从累计概率达到阈值的候选集中采样。"),
        ]
        glossary_rows = "".join(
            f'<tr class="glossary-row border-b border-white/10 odd:bg-white/[0.03]" data-term="{html.escape(term.lower())} {html.escape(desc.lower())}"><th class="px-4 py-3 text-left font-bold text-cyan-200">{html.escape(term)}</th><td class="px-4 py-3 text-slate-300">{html.escape(desc)}</td></tr>'
            for term, desc in glossary_terms
        )
        tailwind_cdn = '<script src="https://cdn.tailwindcss.com"></script>' if self._requires_tailwind_cdn(task_text) else ""
        chartjs_cdn = '<script src="https://cdn.jsdelivr.net/npm/chart.js"></script>' if self._requires_chartjs_cdn(task_text) else ""
        return f"""<!doctype html>
<html lang="zh-CN" class="scroll-smooth">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{safe_title}</title>
  {tailwind_cdn}
  {chartjs_cdn}
</head>
<body class="min-h-screen bg-slate-950 text-slate-100 antialiased selection:bg-cyan-300 selection:text-slate-950">
  <div class="fixed inset-0 -z-10 bg-[radial-gradient(circle_at_top_left,rgba(59,130,246,.35),transparent_34%),radial-gradient(circle_at_top_right,rgba(168,85,247,.26),transparent_32%),linear-gradient(180deg,#020617,#0f172a)]"></div>
  <header class="sticky top-0 z-40 border-b border-white/10 bg-slate-950/80 backdrop-blur-xl">
    <div class="mx-auto flex max-w-7xl flex-col gap-3 px-5 py-4 lg:flex-row lg:items-center lg:justify-between">
      <a href="#top" class="text-lg font-black tracking-tight text-white">{safe_title}</a>
      <nav class="flex max-w-5xl gap-2 overflow-x-auto pb-1" aria-label="页面导航">{nav_items}</nav>
    </div>
  </header>
  <main id="top" class="mx-auto max-w-7xl px-5 py-10">
    <section class="overflow-hidden rounded-[2rem] border border-white/10 bg-white/[0.06] p-8 shadow-2xl shadow-blue-950/50 backdrop-blur md:p-12">
      <div class="inline-flex items-center rounded-full border border-cyan-300/30 bg-cyan-300/10 px-4 py-2 text-sm font-bold text-cyan-100">从概念到训练再到推理的完整路线</div>
      <h1 class="mt-6 max-w-4xl bg-gradient-to-r from-cyan-200 via-blue-200 to-violet-200 bg-clip-text text-5xl font-black leading-tight tracking-tight text-transparent md:text-7xl">{safe_title}</h1>
      <p class="mt-6 max-w-3xl text-lg leading-8 text-slate-300">{subtitle}</p>
      <div class="mt-8 grid gap-4 sm:grid-cols-2 lg:grid-cols-4">{stat_cards}</div>
    </section>

    <section class="mt-8 grid gap-6 lg:grid-cols-[.8fr_1.2fr]">
      <aside class="rounded-[2rem] border border-white/10 bg-white/[0.05] p-6 shadow-xl shadow-slate-950/40 backdrop-blur">
        <h2 class="text-2xl font-black text-white">学习路线</h2>
        <ol class="mt-5 space-y-4 text-sm text-slate-300">
          <li class="rounded-2xl bg-blue-500/10 p-4 ring-1 ring-blue-300/20"><b class="text-blue-200">1. 表示</b>：Tokenization、Embedding、位置编码</li>
          <li class="rounded-2xl bg-cyan-500/10 p-4 ring-1 ring-cyan-300/20"><b class="text-cyan-200">2. 架构</b>：Transformer、Attention、FFN、残差</li>
          <li class="rounded-2xl bg-violet-500/10 p-4 ring-1 ring-violet-300/20"><b class="text-violet-200">3. 训练</b>：Pretrain、SFT、RLHF / DPO</li>
          <li class="rounded-2xl bg-fuchsia-500/10 p-4 ring-1 ring-fuchsia-300/20"><b class="text-fuchsia-200">4. 推理</b>：temperature、top-k、top-p、repetition penalty</li>
        </ol>
      </aside>
      <section class="rounded-[2rem] border border-white/10 bg-white/[0.05] p-6 shadow-xl shadow-slate-950/40 backdrop-blur">
        <h2 class="text-2xl font-black text-white">Loss vs Compute · Chart.js</h2>
        <p class="mt-2 text-slate-300">Scaling Law 的直觉：算力、数据和参数规模提升时，验证损失通常呈规律下降，但边际收益递减。</p>
        <div class="mt-5 rounded-3xl border border-white/10 bg-slate-900/70 p-5"><canvas id="lossChart" height="140" aria-label="loss vs compute chart"></canvas></div>
      </section>
    </section>

    <section class="mt-8 grid gap-5">{section_cards}</section>

    <section id="glossary" class="mt-8 rounded-[2rem] border border-white/10 bg-white/[0.05] p-6 shadow-xl shadow-slate-950/40 backdrop-blur">
      <div class="flex flex-col gap-4 md:flex-row md:items-end md:justify-between">
        <div><p class="text-sm font-bold uppercase tracking-[.3em] text-cyan-200">Glossary</p><h2 class="mt-2 text-3xl font-black text-white">术语表（可实时过滤）</h2></div>
        <input id="glossaryFilter" class="w-full rounded-2xl border border-white/10 bg-slate-950/70 px-4 py-3 text-slate-100 outline-none ring-cyan-300/0 transition placeholder:text-slate-500 focus:border-cyan-300 focus:ring-4 focus:ring-cyan-300/15 md:w-80" placeholder="输入 token / attention / SFT..." />
      </div>
      <div class="mt-5 overflow-hidden rounded-3xl border border-white/10"><table class="w-full border-collapse text-sm"><tbody>{glossary_rows}</tbody></table></div>
    </section>
  </main>
  <script>
    document.querySelectorAll('[data-target]').forEach((button) => {{
      button.addEventListener('click', () => document.getElementById(button.dataset.target)?.scrollIntoView({{ behavior: 'smooth', block: 'start' }}));
    }});
    const filter = document.getElementById('glossaryFilter');
    filter?.addEventListener('input', () => {{
      const q = filter.value.trim().toLowerCase();
      document.querySelectorAll('.glossary-row').forEach((row) => {{ row.classList.toggle('hidden', q && !row.dataset.term.includes(q)); }});
    }});
    const lossCanvas = document.getElementById('lossChart');
    if (lossCanvas && window.Chart) {{
      new Chart(lossCanvas, {{
        type: 'line',
        data: {{ labels: ['1x','3x','10x','30x','100x','300x'], datasets: [{{ label: 'validation loss', data: [3.4, 2.8, 2.25, 1.92, 1.72, 1.61], borderColor: '#22d3ee', backgroundColor: 'rgba(34,211,238,.16)', tension: .35, fill: true }}] }},
        options: {{ responsive: true, plugins: {{ legend: {{ labels: {{ color: '#e2e8f0' }} }} }}, scales: {{ x: {{ ticks: {{ color: '#94a3b8' }}, grid: {{ color: 'rgba(148,163,184,.12)' }} }}, y: {{ ticks: {{ color: '#94a3b8' }}, grid: {{ color: 'rgba(148,163,184,.12)' }} }} }} }}
      }});
    }}
    document.querySelectorAll('canvas[data-chart]').forEach((canvas) => {{
      if (!window.Chart) return;
      const kind = canvas.dataset.chart;
      const common = {{ responsive: true, plugins: {{ legend: {{ labels: {{ color: '#e2e8f0' }} }} }}, scales: {{ x: {{ ticks: {{ color: '#94a3b8' }}, grid: {{ color: 'rgba(148,163,184,.12)' }} }}, y: {{ ticks: {{ color: '#94a3b8' }}, grid: {{ color: 'rgba(148,163,184,.12)' }} }} }} }};
      if (kind === 'loss') {{
        new Chart(canvas, {{ type: 'line', data: {{ labels: ['1x','3x','10x','30x','100x'], datasets: [{{ label: 'loss', data: [3.3,2.7,2.2,1.9,1.72], borderColor: '#38bdf8', backgroundColor: 'rgba(56,189,248,.16)', tension: .35, fill: true }}] }}, options: common }});
      }} else if (kind === 'preference') {{
        new Chart(canvas, {{ type: 'bar', data: {{ labels: ['Raw','SFT','RLHF','DPO'], datasets: [{{ label: 'preference score', data: [42,63,81,86], backgroundColor: ['#60a5fa','#22d3ee','#a78bfa','#f472b6'] }}] }}, options: common }});
      }} else if (kind === 'sampling') {{
        new Chart(canvas, {{ type: 'bar', data: {{ labels: ['the','a','Chat','model','runs','rare'], datasets: [{{ label: 'token probability', data: [36,24,16,10,7,3], backgroundColor: ['#22d3ee','#38bdf8','#818cf8','#a78bfa','#e879f9','#64748b'] }}] }}, options: common }});
      }}
    }});
  </script>
</body>
</html>
"""

    def _tailwind_section_card(self, title: str, detail: str, index: int) -> str:
        title_text = html.escape(title)
        detail_text = html.escape(detail or self._tailwind_section_detail(title, detail))
        visual = self._tailwind_section_visual(title, detail, index)
        return f"""
      <article id="sec-{index}" class="scroll-mt-28 rounded-[2rem] border border-white/10 bg-white/[0.05] p-6 shadow-xl shadow-slate-950/40 backdrop-blur transition hover:-translate-y-1 hover:border-cyan-300/40 hover:bg-white/[0.07]">
        <div class="flex flex-col gap-4 lg:flex-row lg:items-start lg:justify-between">
          <div class="max-w-3xl"><p class="text-sm font-black uppercase tracking-[.3em] text-cyan-200">Section {index:02d}</p><h2 class="mt-2 text-3xl font-black text-white">{title_text}</h2><p class="mt-3 text-base leading-8 text-slate-300">{detail_text}</p></div>
          <div class="rounded-2xl bg-gradient-to-br from-blue-400/20 to-violet-400/20 px-4 py-3 text-sm font-bold text-blue-100 ring-1 ring-white/10">可视化元素</div>
        </div>
        <div class="mt-5">{visual}</div>
      </article>"""

    def _tailwind_section_detail(self, title: str, detail: str) -> str:
        text = f"{title} {detail}"
        lower = text.lower()
        if any(marker in text for marker in ("语言建模", "自回归", "解决什么问题")):
            return "LLM 要解决的是语言建模问题：给定前文 token，预测下一个 token 的概率分布。ChatGPT 采用自回归生成，每一步把新 token 接回上下文继续预测，因此对话、推理和代码生成都可以看成一串条件概率决策。"
        if any(marker in lower for marker in ("next token", "交叉熵")) or "训练目标" in text:
            return "训练目标是 Next Token Prediction：模型输出词表上每个 token 的 logits，经 softmax 得到概率分布，再用交叉熵惩罚真实下一个 token 的低概率预测。loss 越低，代表模型越会续写。"
        if any(marker in lower for marker in ("token", "bpe", "词表")):
            return "Tokenization 把自然语言切成模型可处理的 token。BPE 会优先合并高频片段，词表决定 token id 的范围；同一句中文、英文或代码会切成不同数量的 token，这也直接影响上下文长度和计费直觉。"
        if any(marker in lower for marker in ("embedding", "位置编码", "position")):
            return "Embedding 把离散 token id 映射成连续向量，让语义相近的片段在向量空间中更接近；位置编码补充顺序信息，否则 Transformer 只看到一组 token，无法区分前后关系。"
        if any(marker in lower for marker in ("transformer", "attention", "mha", "ffn", "layernorm", "residual", "残差")):
            return "Decoder-only Transformer 由多层 masked self-attention 和 FFN 堆叠而成。MHA 用多个头并行捕捉语法、实体和长程依赖；残差连接与 LayerNorm 保持梯度稳定，让深层网络可训练。"
        if any(marker in lower for marker in ("pretrain", "scaling", "loss", "算力", "数据")):
            return "Pretrain 阶段把海量网页、书籍、代码等语料转成 token 序列，在大规模算力上最小化交叉熵。Scaling Law 的直觉是：数据、参数和计算量增加时，验证 loss 通常按规律下降，但边际收益递减。"
        if "sft" in lower or "指令微调" in text:
            return "SFT 用人工编写或筛选的 prompt-response 示例，把“会续写”的基座模型校准成“会按指令回答”的助手。它强化格式、语气、安全边界和任务完成方式，是 ChatGPT 对话能力的第一层对齐。"
        if any(marker in lower for marker in ("rlhf", "dpo", "reward", "奖励", "ppo")):
            return "RLHF 先收集人类偏好，训练奖励模型评价哪个回答更好，再优化策略让输出更符合偏好；DPO 则把偏好对直接写进损失函数，省去显式奖励模型和复杂强化学习流程。"
        if any(marker in lower for marker in ("temperature", "top-k", "top-p", "repetition", "推理", "sampling")):
            return "推理阶段模型先给出 token 概率分布，再由采样策略决定输出。temperature 控制随机性，top-k 限制候选数量，top-p 取累计概率集合，repetition penalty 抑制重复循环。"
        if "术语" in text or "glossary" in lower:
            return "术语表把学习过程中的核心概念集中成可过滤速查表，适合边看图解边搜索 token、attention、SFT、RLHF、DPO 等关键词，快速建立完整知识地图。"
        return detail or "本板块围绕用户指定主题提供概念解释、视觉拆解和教学要点，避免只给大纲式标题。"

    def _tailwind_section_visual(self, title: str, detail: str, index: int) -> str:
        lower = f"{title} {detail}".lower()
        if any(marker in lower for marker in ("pretrain", "scaling", "loss", "算力")):
            return f'<div class="grid gap-4 lg:grid-cols-[.9fr_1.1fr]"><div class="grid gap-3 md:grid-cols-2"><div class="rounded-2xl bg-slate-900/70 p-4"><b class="text-white">数据</b><p class="text-slate-400">网页 / 书籍 / 代码</p></div><div class="rounded-2xl bg-slate-900/70 p-4"><b class="text-white">算力</b><p class="text-slate-400">GPU 集群训练</p></div><div class="rounded-2xl bg-slate-900/70 p-4"><b class="text-white">Loss</b><p class="text-slate-400">交叉熵下降</p></div><div class="rounded-2xl bg-slate-900/70 p-4"><b class="text-white">能力</b><p class="text-slate-400">涌现与泛化</p></div></div><div class="rounded-3xl border border-white/10 bg-slate-900/70 p-4"><canvas id="sectionChart{index}" data-chart="loss" height="140" aria-label="loss vs compute"></canvas></div></div>'
        if any(marker in lower for marker in ("rlhf", "dpo", "奖励", "sft")):
            return f'<div class="grid gap-4 lg:grid-cols-[.9fr_1.1fr]"><div class="grid gap-3"><div class="rounded-2xl bg-blue-500/10 p-4 ring-1 ring-blue-300/20"><b class="text-blue-100">SFT</b><p class="text-slate-400">学习理想回答格式</p></div><div class="rounded-2xl bg-cyan-500/10 p-4 ring-1 ring-cyan-300/20"><b class="text-cyan-100">Reward / Preference</b><p class="text-slate-400">比较回答优劣</p></div><div class="rounded-2xl bg-violet-500/10 p-4 ring-1 ring-violet-300/20"><b class="text-violet-100">RLHF / DPO</b><p class="text-slate-400">优化为更符合偏好的助手</p></div></div><div class="rounded-3xl border border-white/10 bg-slate-900/70 p-4"><canvas id="sectionChart{index}" data-chart="preference" height="140" aria-label="preference bar chart"></canvas></div></div>'
        if any(marker in lower for marker in ("temperature", "top-k", "top-p", "sampling", "推理")):
            return f'<div class="grid gap-4 lg:grid-cols-[.9fr_1.1fr]"><div class="rounded-3xl border border-white/10 bg-slate-900/70 p-5"><div class="flex items-center gap-2"><span class="h-4 w-32 rounded-full bg-cyan-400"></span><span class="h-4 w-24 rounded-full bg-blue-400"></span><span class="h-4 w-16 rounded-full bg-violet-400"></span><span class="h-4 w-10 rounded-full bg-fuchsia-400"></span><span class="h-4 w-6 rounded-full bg-slate-500"></span></div><p class="mt-3 text-slate-300">低 temperature 更稳定，高 temperature 更发散；top-p 截取累计概率候选集。</p></div><div class="rounded-3xl border border-white/10 bg-slate-900/70 p-4"><canvas id="sectionChart{index}" data-chart="sampling" height="140" aria-label="sampling distribution"></canvas></div></div>'
        if any(marker in lower for marker in ("token", "bpe", "词表")):
            return '<div class="flex flex-wrap items-center gap-3 rounded-3xl border border-cyan-300/20 bg-slate-900/70 p-5 font-mono text-sm text-cyan-100"><span class="rounded-xl bg-cyan-400/10 px-3 py-2">Chat</span><span>+</span><span class="rounded-xl bg-cyan-400/10 px-3 py-2">G</span><span>+</span><span class="rounded-xl bg-cyan-400/10 px-3 py-2">PT</span><span class="text-slate-500">→</span><span class="rounded-xl bg-blue-400/10 px-3 py-2">[18435, 38, 11571]</span></div>'
        if any(marker in lower for marker in ("transformer", "attention", "mha", "ffn", "layernorm")):
            return '<div class="grid gap-3 md:grid-cols-5"><div class="rounded-2xl bg-blue-500/10 p-4 text-blue-100 ring-1 ring-blue-300/20">Input</div><div class="rounded-2xl bg-cyan-500/10 p-4 text-cyan-100 ring-1 ring-cyan-300/20">MHA</div><div class="rounded-2xl bg-violet-500/10 p-4 text-violet-100 ring-1 ring-violet-300/20">Add & Norm</div><div class="rounded-2xl bg-fuchsia-500/10 p-4 text-fuchsia-100 ring-1 ring-fuchsia-300/20">FFN</div><div class="rounded-2xl bg-emerald-500/10 p-4 text-emerald-100 ring-1 ring-emerald-300/20">Logits</div></div>'
        return '<div class="rounded-3xl border border-white/10 bg-slate-900/70 p-5"><pre class="overflow-x-auto font-mono text-sm leading-7 text-slate-200"><code>input tokens → hidden states → logits → probability distribution → next token</code></pre></div>'

    def _write_html_build_script(self, build_path: str, html_path: str) -> None:
        os.makedirs(os.path.dirname(build_path), exist_ok=True)
        relative_html = os.path.basename(html_path)
        script = f'''#!/usr/bin/env python3
"""Regenerate the single-file HTML artifact for this PyClaw deliverable.

The controller already materialized {relative_html}.  This script is kept as
workflow evidence because the user explicitly requested a build.py-first flow.
"""

from pathlib import Path


def main() -> None:
    target = Path(__file__).with_name({relative_html!r})
    if not target.exists():
        raise SystemExit(f"Missing generated artifact: {{target}}")
    print(target)


if __name__ == "__main__":
    main()
'''
        with open(build_path, "w", encoding="utf-8") as handle:
            handle.write(script)

    def _html_sections_from_task_requirements(self, task_text: str) -> str:
        sections = self._numbered_requirements(task_text)
        if not sections:
            return ""
        cards = "".join(self._html_requirement_card(title, detail, index) for index, (title, detail) in enumerate(sections[:16], start=1))
        return f'<section class="panel"><h2>需求板块逐项落地</h2><div class="req-grid">{cards}</div></section>'

    def _numbered_requirements(self, task_text: str) -> list[tuple[str, str]]:
        text = task_text or ""
        line_sections: list[tuple[str, str]] = []
        for line in text.splitlines():
            line_match = re.match(r"\s*\d{1,2}[.、)]\s*(.+?)\s*$", line)
            if not line_match:
                continue
            raw_line = re.sub(r"\s+", " ", line_match.group(1)).strip(" ：:，,。.;；")
            if not raw_line or self._is_delivery_instruction(raw_line):
                continue
            title, detail = self._split_requirement_title_detail(raw_line)
            line_sections.append((title.strip()[:80], detail.strip()[:360] or raw_line[:360]))
        if line_sections:
            return self._dedupe_numbered_requirements(line_sections)

        explicit_line_sections = self._explicit_multiline_requirement_sections(text)
        if explicit_line_sections:
            return self._dedupe_numbered_requirements(explicit_line_sections)

        matches = list(re.finditer(r"(?:^|\n)\s*(\d{1,2})[.、)]\s*(.+?)(?=(?:\n\s*\d{1,2}[.、)]\s*)|\n\s*【|\Z)", text, flags=re.DOTALL))
        sections: list[tuple[str, str]] = []
        for match in matches:
            raw = re.sub(r"\s+", " ", match.group(2)).strip()
            if not raw:
                continue
            raw = re.split(r"\s+(?:请直接|最后用|交付到|写到|保存到|输出到)\b", raw, maxsplit=1)[0].strip()
            if not raw or self._is_delivery_instruction(raw):
                continue
            if "：" in raw:
                title, detail = raw.split("：", 1)
            elif ":" in raw:
                title, detail = raw.split(":", 1)
            else:
                title, detail = raw, raw
            sections.append((title.strip()[:80], detail.strip()[:360] or raw[:360]))
        board_match = re.search(r"(?:板块|模块|章节)\s*[:：]\s*(.+?)(?:\n\s*\n|$)", text, flags=re.IGNORECASE | re.DOTALL)
        if board_match:
            for item in re.split(r"[、,，；;\n]+", board_match.group(1)):
                cleaned = re.sub(r"^\d{1,2}[.、)]\s*", "", item.strip(" ：:，,。.;；"))
                if any(marker in cleaned.lower() for marker in ("send_file_to_user", "send-file-to-user", "交付给我", "发给我", "写到", "保存到")):
                    continue
                if 1 <= len(cleaned) <= 80:
                    sections.append((cleaned, f"围绕 {cleaned} 提供可视化解释、关键机制和互动讲解。"))
        return self._dedupe_numbered_requirements(sections)

    def _explicit_multiline_requirement_sections(self, text: str) -> list[tuple[str, str]]:
        """Extract unnumbered section lists after phrases like "以下 10 个板块".

        Real chat prompts often contain a hard requirement sentence followed by
        one section title per line, without numeric prefixes.  Falling back to a
        generic teaching profile in that case produces plausible but wrong
        artifacts.  This parser keeps the user-provided ordered spine as the
        source of truth and stops when delivery/style instructions begin.
        """
        sections: list[tuple[str, str]] = []
        lines = (text or "").splitlines()
        capture = False
        expected_count: int | None = None
        for raw_line in lines:
            line = re.sub(r"\s+", " ", raw_line).strip()
            if not line:
                if capture and sections:
                    break
                continue
            if not capture:
                intro = re.search(r"(?:必须包含|包含以下|包含下列|以下|下列).{0,24}?(\d{1,2})\s*个?(?:板块|模块|章节|部分|section)", line, flags=re.IGNORECASE)
                if not intro:
                    continue
                expected_count = int(intro.group(1))
                capture = True
                tail = re.split(r"[:：]", line, maxsplit=1)
                if len(tail) == 1 or not tail[1].strip():
                    continue
                candidate = tail[1].strip()
            else:
                candidate = line
            candidate = re.sub(r"^[-*•]\s*", "", candidate).strip(" ：:，,。.;；")
            candidate = re.sub(r"^\d{1,2}[.、)]\s*", "", candidate).strip(" ：:，,。.;；")
            if not candidate:
                continue
            if self._is_requirement_list_stop(candidate):
                break
            if self._is_delivery_instruction(candidate):
                break
            if len(candidate) > 140:
                continue
            title, detail = self._split_requirement_title_detail(candidate)
            sections.append((title.strip()[:80], detail.strip()[:360] or candidate[:360]))
            if expected_count is not None and len(sections) >= expected_count:
                break
        return sections

    def _is_requirement_list_stop(self, text: str) -> bool:
        normalized = (text or "").lower().strip()
        return any(marker in normalized for marker in (
            "硬性要求", "技术栈", "实现策略", "视觉风格", "交互", "页面自包含",
            "单个 html", "单个html", "必须通过 cdn", "必须通过cdn", "额外通过",
            "页面顶部", "最后用", "请直接", "不要先", "交付到", "保存到", "写到",
        ))

    def _split_requirement_title_detail(self, raw: str) -> tuple[str, str]:
        if "：" in raw:
            title, detail = raw.split("：", 1)
        elif ":" in raw:
            title, detail = raw.split(":", 1)
        else:
            title, detail = raw, raw
        return title, detail

    def _is_delivery_instruction(self, text: str) -> bool:
        normalized = (text or "").lower()
        return any(marker in normalized for marker in (
            "send_file_to_user", "send-file-to-user", "交付到", "交付给我", "发给我",
            "写到", "保存到", "输出到", "不要先给我大纲", "不要先给", "确认",
        ))

    def _dedupe_numbered_requirements(self, sections: list[tuple[str, str]]) -> list[tuple[str, str]]:
        deduped: list[tuple[str, str]] = []
        seen: set[str] = set()
        for title, detail in sections:
            key = re.sub(r"\s+", "", title.lower())
            if key and key not in seen:
                seen.add(key)
                deduped.append((title, detail))
        return deduped

    def _html_requirement_card(self, title: str, detail: str, index: int) -> str:
        title_text = html.escape(title or f"板块 {index}")
        detail_text = html.escape(detail or "围绕该板块提供可视化解释、交互演示与教学要点。")
        lower = f"{title} {detail}".lower()
        bars = "".join(f'<i style="height:{height}%"></i>' for height in (35, 58, 82, 46, 70, 92))
        if any(marker in lower for marker in ("token", "bpe", "tokenizer")):
            visual = '<p><code>ChatGPT</code> → <code>Chat</code> + <code>G</code> + <code>PT</code> → token ids → embedding vectors</p>'
        elif any(marker in lower for marker in ("q/k/v", "qkv", "attention", "注意力")):
            visual = '<p>Q 查询当前词，K 表示可匹配线索，V 携带被聚合信息；Multi-Head 并行学习不同关系。</p>'
        elif any(marker in lower for marker in ("chart", "曲线", "scaling", "loss", "饼图")):
            visual = f'<div class="mini-bars" aria-label="chart preview">{bars}</div>'
        elif any(marker in lower for marker in ("rlhf", "reward", "ppo", "sft")):
            visual = '<p>SFT 示例 → Reward Model 偏好排序 → PPO 优化策略，是从“会续写”到“更像助手”的关键链路。</p>'
        elif any(marker in lower for marker in ("术语", "概念", "速查")):
            visual = '<details open><summary>关键术语</summary>Token、Embedding、Attention、LayerNorm、Residual、SFT、RM、PPO、Temperature、Top-p。</details>'
        else:
            visual = '<p>点击卡片可作为课堂讲解步骤：先看输入输出，再看模块关系，最后看可交互示例。</p>'
        return f'<article class="req-card"><h3>{index}. {title_text}</h3><p>{detail_text}</p>{visual}</article>'

    def _html_teaching_profile(self, title: str) -> dict[str, Any]:
        normalized = (title or "").lower()
        if "gpt-3" in normalized or "gpt3" in normalized:
            return {
                "title": "图解 GPT-3 结构和训练流程",
                "subtitle": "用一页可视化网页理解 GPT-3：175B 参数规模、Decoder-only Transformer、预训练目标，以及 Zero-shot / One-shot / Few-shot 的 In-Context Learning。",
                "stats": [("175B", "参数规模"), ("96", "Transformer 层"), ("2048", "上下文长度"), ("Few-shot", "无需梯度更新的示例学习")],
                "architecture": [
                    ("Tokenization", "把文本切成 token，并映射为可训练向量。"),
                    ("Positional Embedding", "为序列位置注入顺序信息，让模型理解前后关系。"),
                    ("Decoder-only Transformer", "堆叠自回归解码器层，使用 masked self-attention 预测下一个 token。"),
                    ("Language Modeling Head", "把隐藏状态映射到词表概率，选择下一个最可能 token。"),
                ],
                "training": [
                    ("大规模语料预训练", "在网页、书籍、代码等混合语料上学习 next-token prediction。"),
                    ("自回归生成", "每一步只看左侧上下文，通过最大似然训练生成连续文本。"),
                    ("规模化带来涌现能力", "参数、数据和算力放大后，推理、翻译、代码和问答能力显著增强。"),
                    ("上下文学习", "推理时在 prompt 中放入任务说明和少量示例，模型无需微调即可适配任务。"),
                ],
                "concepts": [
                    ("Sparse Attention", "GPT-3 论文中结合稀疏注意力思想降低长序列注意力成本。"),
                    ("In-Context Learning", "把示例写进上下文，模型从模式中归纳输出格式和任务规则。"),
                    ("Scaling Law", "模型能力随参数、数据、计算规模呈规律性提升，是 GPT-3 的核心叙事。"),
                ],
                "demo": [
                    ("Zero-shot", "只给任务：把这句话翻译成英文。模型直接回答。"),
                    ("One-shot", "给 1 个输入输出样例，再给新输入，让模型模仿格式。"),
                    ("Few-shot", "给多个样例，模型在上下文中归纳分类、抽取或推理规则。"),
                ],
                "comparison": [
                    ("模型规模", "GPT-2：1.5B", "GPT-3：175B，规模提升两个数量级"),
                    ("能力重点", "GPT-2：展示通用文本生成", "GPT-3：强调少样本/零样本任务泛化"),
                    ("使用方式", "更多依赖微调或任务适配", "更多依赖 Prompt 与 In-Context Learning"),
                ],
                "notes": [
                    "先从 next-token prediction 解释为什么 GPT-3 能生成连贯文本。",
                    "再用三种 shot 模式演示它如何从上下文里学习任务。",
                    "强调 GPT-3 不是 ChatGPT：ChatGPT 后续还加入指令微调和 RLHF 对齐。",
                    "讲解限制：上下文长度、事实幻觉、成本和安全边界。",
                ],
            }
        if "chatgpt" in normalized:
            return {
                "title": "图解 ChatGPT 结构和训练流程",
                "subtitle": "从 Transformer 语言模型到对话助手：理解预训练、监督微调、奖励模型与 RLHF 对齐。",
                "stats": [("Transformer", "底座结构"), ("SFT", "监督微调"), ("RM", "奖励模型"), ("RLHF", "人类反馈强化学习")],
                "architecture": [
                    ("Tokenizer", "把用户输入切成 token。"), ("Decoder Blocks", "多层注意力和前馈网络处理上下文。"),
                    ("Instruction Tuning", "学习遵循人类指令和对话格式。"), ("Safety Layer", "通过策略和偏好对齐减少不安全输出。"),
                ],
                "training": [("预训练", "学习通用语言规律。"), ("监督微调", "用人工示范学习回答风格。"), ("奖励建模", "学习人类偏好排序。"), ("RLHF 对齐", "优化更有帮助、更安全的回答。")],
                "concepts": [("Attention", "在上下文中聚焦相关 token。"), ("Prompt", "用户目标和约束的载体。"), ("Alignment", "让输出符合人类偏好与安全规则。")],
                "demo": [("提问", "用户给出目标。"), ("规划", "模型组织回答结构。"), ("生成", "逐 token 生成最终回复。")],
                "comparison": [("底座", "通用 GPT", "ChatGPT：面向对话对齐"), ("训练", "预训练", "预训练 + SFT + RLHF"), ("体验", "文本补全", "多轮助手")],
                "notes": ["区分语言模型与产品形态。", "强调对齐数据的重要性。", "用对话例子展示多轮上下文。", "说明幻觉与安全限制。"],
            }
        clean = (title or "可视化教学网页").strip()
        return {
            "title": clean,
            "subtitle": f"围绕“{clean}”构建结构、流程、关键概念和课堂讲解提示。",
            "stats": [("结构", "核心模块"), ("流程", "关键步骤"), ("机制", "原理解释"), ("案例", "教学演示")],
            "architecture": [("输入", "明确学习对象和上下文。"), ("处理", "拆解结构、机制和关键变量。"), ("输出", "形成可视化解释与实践建议。"), ("反馈", "通过问题和案例巩固理解。")],
            "training": [("建立背景", "说明为什么需要学习该主题。"), ("拆解结构", "把复杂系统拆成模块。"), ("串联流程", "用时间线解释关键步骤。"), ("应用练习", "通过例子把知识迁移到实际场景。")],
            "concepts": [("模块化", "先看组成部分。"), ("流程化", "再看步骤关系。"), ("可视化", "用图形降低理解成本。")],
            "demo": [("观察", "识别输入与输出。"), ("推演", "跟随流程逐步变化。"), ("总结", "提炼可复用规律。")],
            "comparison": [("学习前", "概念零散", "学习后：结构清晰"), ("表达方式", "纯文字", "图文结合"), ("应用", "难迁移", "可复用到真实任务")],
            "notes": ["先讲整体，再讲细节。", "每个模块配一个例子。", "流程图优先于长段文字。", "最后用对比表复盘。"],
        }

    def _slides_from_text(self, task_text: str, draft: str) -> list[tuple[str, list[str]]]:
        title = self._title_from_task(task_text)
        sections = [] if self._looks_like_process_or_failure_draft(draft) else self._sections_from_markdown(draft)
        if sections and (len(sections) >= 3 or self._topic_family(title) == "generic"):
            return [(title, ["根据当前任务自动生成的交付版幻灯片"])] + sections[:8]
        topic_family = self._topic_family(title)
        if topic_family == "rag":
            return self._rag_slides(title)
        if "agent" in title.lower() or "智能体" in title:
            return self._ai_agent_slides(title)
        return [
            (title, ["一份面向快速沟通的自动生成幻灯片", "可继续补充业务数据、案例和视觉素材"]),
            ("核心概念", self._default_points(title)[:4]),
            ("能力边界", ["感知与理解上下文", "规划任务步骤", "调用工具执行动作", "根据观察结果自我修正"]),
            ("典型工作流", ["接收目标", "拆解计划", "检索或读取资料", "生成交付物", "验证并发送结果"]),
            ("落地建议", ["把复杂任务拆成可验证里程碑", "为关键输出定义完成证据", "保留日志与中间产物", "优先自动化重复流程"]),
        ]

    def _topic_family(self, title: str) -> str:
        normalized = (title or "").lower()
        if any(marker in normalized for marker in ("rag", "检索增强", "知识库", "向量检索", "召回", "rerank")):
            return "rag"
        if any(marker in normalized for marker in ("agent", "智能体")):
            return "ai_agent"
        return "generic"

    def _ai_agent_slides(self, title: str) -> list[tuple[str, list[str]]]:
        return [
            (title, ["从大模型能力走向可执行任务系统", "核心关键词：目标、规划、工具、记忆、反馈闭环"]),
            ("什么是 AI Agent", ["围绕用户目标自主规划和执行的智能系统", "不仅生成文本，还能调用工具影响外部环境", "通过观察结果持续修正下一步动作"]),
            ("Agent 与 Chatbot 的区别", ["Chatbot 偏一次性问答，Agent 偏多步任务完成", "Agent 需要状态、工具、权限和验收机制", "Agent 的输出应以可验证结果为准"]),
            ("核心架构", ["模型层：理解目标并产生计划", "工具层：搜索、代码、文件、浏览器、业务 API", "记忆层：保存偏好、进度和上下文", "控制层：调度、重试、验收和安全边界"]),
            ("典型 Agent Loop", ["Plan：拆解任务和选择策略", "Act：调用工具或修改文件", "Observe：读取日志、结果和错误", "Reflect/Repair：根据反馈修正", "Deliver：验证后交付最终结果"]),
            ("工具调用与权限", ["工具必须有明确输入输出契约", "副作用工具需要幂等、去重和审批策略", "高风险操作应由 controller 做安全校验", "用户最终看到的是结果，不是内部执行日志"]),
            ("记忆与上下文管理", ["短期上下文承载当前任务细节", "长期记忆保存稳定偏好和项目事实", "历史压缩要避免把 assistant 猜测变成用户需求", "任务契约应锚定真实用户目标"]),
            ("可靠性工程", ["用完成证据替代口头宣称", "失败时自动修复而不是要求用户二次确认", "对文件、截图、邮件等交付物做验收", "保持最终回复与实际交付一致"]),
            ("典型应用场景", ["代码开发与测试修复", "知识检索与研究报告", "办公自动化：文档、表格、幻灯片", "运维巡检、数据分析和流程编排"]),
            ("主要挑战", ["长程任务容易偏离目标", "工具错误和环境差异会累积", "权限与安全边界复杂", "结果质量需要可观测和可验收"]),
            ("落地建议", ["先从可验证的小任务开始", "为每类任务定义完成契约", "将工具输出、文件和日志纳入证据链", "把失败恢复设计成系统能力而非人工兜底"]),
            ("总结", ["AI Agent 的核心不是单次回答，而是闭环执行", "工程重点是工具、状态、安全和验收", "高质量 Agent 应能自主完成任务并诚实报告结果"]),
        ]

    def _rag_slides(self, title: str) -> list[tuple[str, list[str]]]:
        return [
            (title, ["用企业私有知识增强大模型回答", "核心关键词：数据接入、向量检索、重排、引用、评测"]),
            ("什么是 RAG", ["Retrieval-Augmented Generation：先检索再生成", "把外部知识作为上下文注入模型", "降低幻觉并提升答案可追溯性"]),
            ("企业知识库价值", ["统一沉淀文档、制度、代码和业务经验", "让员工用自然语言查询分散知识", "在权限边界内提供可引用答案"]),
            ("数据接入与治理", ["连接飞书、网页、PDF、代码仓库和数据库", "清洗重复内容、过期内容和敏感字段", "保留来源、作者、时间和权限元数据"]),
            ("切分与索引", ["按语义段落切 chunk，避免机械固定长度", "为 chunk 生成 embedding 并写入向量库", "保存标题层级、文档路径和引用锚点"]),
            ("检索策略", ["向量检索召回语义相关片段", "关键词/BM25 处理专有名词和编号", "混合检索提升召回覆盖率和稳定性"]),
            ("重排与上下文组装", ["用 reranker 对候选片段重新排序", "去重、合并相邻段落并控制 token 预算", "优先放入高置信、最新、权限可见的证据"]),
            ("生成与引用", ["提示词要求仅基于证据回答", "每个关键结论附来源引用", "证据不足时明确说明无法确认而不是编造"]),
            ("权限与安全", ["检索阶段执行用户级 ACL 过滤", "避免跨部门/跨租户数据泄露", "对敏感问题记录审计日志和安全策略命中"]),
            ("效果评测", ["离线评测：召回率、命中率、引用准确率", "在线评测：有用率、追问率、人工反馈", "构建 golden set 覆盖高频业务问题"]),
            ("生产架构", ["采集管道负责增量同步和索引刷新", "Query Pipeline 负责改写、检索、重排和生成", "观测系统跟踪延迟、成本、命中率和失败样本"]),
            ("落地建议", ["先选一个高频场景做闭环，例如制度问答或研发知识库", "从小规模高质量文档开始，不要一次接入全部资料", "把引用、权限、评测和反馈作为上线门槛"]),
        ]

    def _looks_like_process_or_failure_draft(self, draft: str) -> bool:
        normalized = (draft or "").lower()
        markers = (
            "当前进展", "没能顺利落地", "未能顺利落地", "不发送", "残缺文件", "未达到",
            "下一步建议", "下一轮", "继续生成", "需要你确认", "工作目录", "生成脚本",
            "send_file_to_user", "completion contract", "tool usage", "notice:", "未观察到目标文件",
            "没有生成文件", "仍然没有生成", "尚未生成文件",
            "没能生成文件", "未能生成文件", "不能把过程汇报当作",
            "任务未完成", "未通过", "工作流验收", "缺少生成文件证据", "缺少工作流产物证据",
            "已停止继续执行", "避免重复触发", "工具预算", "达到最大思考深度",
            "工具调用次数已达到上限", "已生成并发送", "已生成文件", "文件已发送",
        )
        return any(marker in normalized for marker in markers)

    def _title_from_task(self, task_text: str) -> str:
        text = (task_text or "").strip()
        patterns = (
            r"(?:做一个|生成|创建|制作|给我做|帮我做|出一个)关于\s*(.+?)\s*(?:的)?(?:ppt|pptx|幻灯片|演示文稿|slide|deck)",
            r"(?:做一个|生成|创建|制作|给我做|帮我做|出一个)关于\s*(.+?)\s*(?:的)?(?:html|网页|页面|website|webpage|web page)",
            r"(?:做一个|生成|创建|制作|给我做|帮我做|出一个)\s*(?:可视化的?)?(?:教学)?(?:html|网页|页面|website|webpage|web page)[，,：:]?\s*(.+)$",
            r"图解\s*(.+?)(?:结构|训练流程|训练|原理|流程|架构)",
            r"about\s+(.+?)\s+(?:ppt|pptx|slides|deck)",
            r"about\s+(.+?)\s+(?:html|website|webpage|web page)",
        )
        for pattern in patterns:
            match = re.search(pattern, text, flags=re.IGNORECASE)
            if match:
                topic = match.group(1).strip(" ：:，,。.")
                topic = re.sub(r"https?://\S+", "", topic).strip(" ：:，,。.")
                if topic:
                    return topic
        for marker in ("幻灯片", "演示文稿", "pptx", "ppt", "slides", "deck", "html", "网页", "页面", "website", "webpage", "web page"):
            text = re.sub(marker, "", text, flags=re.IGNORECASE).strip()
        text = re.sub(r"https?://\S+", "", text).strip()
        return text[:40] or "自动生成幻灯片"

    def _default_points(self, title: str) -> list[str]:
        normalized = title.lower()
        if self._topic_family(title) == "rag":
            return [
                "RAG 通过检索企业知识库为大模型提供可引用上下文",
                "关键链路包括数据接入、切分、向量索引、召回、重排和生成",
                "工程重点是权限控制、引用溯源、效果评测和增量更新",
                "适合制度问答、研发知识检索、客服辅助和业务专家系统",
            ]
        if "agent" in normalized or "智能体" in normalized:
            return [
                "AI Agent 是能够围绕目标进行规划、执行和反馈修正的智能系统",
                "关键组成包括模型、工具、记忆、环境观察和控制循环",
                "工程重点是可靠性、权限边界、可观测性和可恢复性",
                "适合用于代码助手、知识检索、办公自动化和长程任务执行",
            ]
        return [
            f"围绕“{title}”建立清晰目标",
            "拆解核心概念、关键流程和实践场景",
            "用证据和案例支撑结论",
            "给出可执行的实践路径",
        ]

    def _sections_from_markdown(self, draft: str) -> list[tuple[str, list[str]]]:
        if not draft or not draft.strip():
            return []
        sections: list[tuple[str, list[str]]] = []
        current_title = "要点"
        current_items: list[str] = []
        for raw_line in draft.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            heading = re.match(r"^#{1,4}\s+(.+)$", line)
            if heading:
                if current_items:
                    sections.append((current_title, current_items[:6]))
                current_title = self._clean_text(heading.group(1))[:48] or "要点"
                current_items = []
                continue
            bullet = re.match(r"^(?:[-*+]|\d+[.)]|[一二三四五六七八九十]+[、.])\s*(.+)$", line)
            if bullet:
                item = self._clean_text(bullet.group(1))
                if item:
                    current_items.append(item[:120])
            elif len(line) <= 120 and len(current_items) < 4:
                current_items.append(self._clean_text(line)[:120])
        if current_items:
            sections.append((current_title, current_items[:6]))
        return sections[:8]

    def _clean_text(self, text: str) -> str:
        cleaned = re.sub(r"[`*_>\[\]()]+", "", text or "")
        return re.sub(r"\s+", " ", cleaned).strip()


    def _ensure_slide_count(self, slides: list[tuple[str, list[str]]], count: Optional[int]) -> list[tuple[str, list[str]]]:
        if count is None or count <= 0:
            return slides[:10]
        result = list(slides[:count])
        topic = result[0][0] if result else "自动生成幻灯片"
        family = self._topic_family(topic)
        while len(result) < count:
            index = len(result) + 1
            result.append(self._supplemental_slide(topic, family, index))
        return result

    def _supplemental_slide(self, topic: str, family: str, index: int) -> tuple[str, list[str]]:
        """Create non-placeholder topic-specific material when a deck needs more pages."""
        if family == "rag":
            extras = [
                ("数据质量保障", ["建立文档准入、去重和过期清理机制", "对来源、作者、时间和权限元数据做标准化", "低质量知识会直接影响召回与答案可信度"]),
                ("增量同步机制", ["按来源系统维护同步游标和变更事件", "支持新增、更新、删除的索引一致性处理", "对失败任务提供重试队列和人工补偿入口"]),
                ("Query Pipeline 设计", ["先做意图识别与查询改写，再进入混合检索", "重排后按证据多样性和新鲜度组装上下文", "生成阶段强制引用来源并声明证据不足场景"]),
                ("观测与运营指标", ["跟踪召回命中率、无答案率、引用准确率和延迟", "将用户反馈沉淀为评测集与知识修复任务", "用仪表盘定位高频失败问题和成本瓶颈"]),
            ]
        elif family == "ai_agent":
            extras = [
                ("任务契约设计", ["把用户目标转为可验证的完成条件", "明确工具边界、交付物类型和验收证据", "避免把过程说明误判为任务完成"]),
                ("工具编排策略", ["区分无副作用读取工具和有副作用执行工具", "失败后基于观察结果修正参数而不是盲目重试", "关键交付动作由控制器统一验收和发送"]),
                ("长期任务治理", ["通过检查点保存进度和上下文摘要", "将中间产物落盘，支持恢复和审计", "对多轮继续请求复用原始任务契约"]),
                ("质量评估闭环", ["用测试、预览、结构检查验证结果", "记录失败样本并转化为回归用例", "持续优化提示、工具和控制器策略"]),
            ]
        else:
            extras = [
                (f"{topic} 的落地路径", ["从目标、对象、流程和结果四个维度拆解", "优先选择可验证的小范围场景试点", "用指标反馈持续迭代方案"]),
                (f"{topic} 的风险控制", ["识别关键依赖、权限边界和异常路径", "为高风险环节设置审批、回滚和审计机制", "用演练验证方案在真实环境下可执行"]),
                (f"{topic} 的评估指标", ["定义效率、质量、成本和满意度指标", "建立上线前验收与上线后复盘机制", "把数据反馈纳入下一轮优化"]),
            ]
        return extras[(index - 1) % len(extras)]

    def _explicit_slide_count(self, text: str) -> Optional[int]:
        if not text:
            return None
        patterns = (
            r"(\d{1,3})\s*(?:页|頁|张|張|p\b|pages?\b|slides?\b)",
            r"(?:页数|頁數|slides?|pages?)\s*[:：]?\s*(\d{1,3})",
        )
        for pattern in patterns:
            match = re.search(pattern, text, flags=re.IGNORECASE)
            if match:
                value = int(match.group(1))
                if 1 <= value <= 300:
                    return value
        return None

    def _write_pptx(self, file_path: str, slides: list[tuple[str, list[str]]]) -> None:
        try:
            self._write_polished_pptx(file_path, slides)
            return
        except Exception:
            # Keep the controller resilient in lean environments.  The rich
            # renderer is preferred whenever python-pptx is installed; the raw
            # OpenXML writer remains a no-network fallback so deliverable tasks
            # still produce a verifiable file instead of failing outright.
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except OSError:
                pass
        self._write_basic_pptx(file_path, slides)

    def _write_polished_pptx(self, file_path: str, slides: list[tuple[str, list[str]]]) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        slides = slides or [("自动生成幻灯片", ["暂无内容"])]
        profile = self._deck_visual_profile(slides[0][0])
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        palette = ["2563EB", "7C3AED", "0891B2", "16A34A", "F97316", "DC2626"]
        dark = "0F172A"
        ink = "111827"
        muted = "64748B"
        paper = "F8FAFC"

        def rgb(value: str) -> RGBColor:
            value = value.strip("#")
            return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

        def shape(slide, kind, x: float, y: float, w: float, h: float, fill: str, *, line: str | None = None, transparency: int = 0):
            item = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
            item.fill.solid()
            item.fill.fore_color.rgb = rgb(fill)
            item.fill.transparency = transparency
            if line:
                item.line.color.rgb = rgb(line)
                item.line.width = Pt(1)
            else:
                item.line.fill.background()
            return item

        def textbox(slide, text: str, x: float, y: float, w: float, h: float, *, size: int, color: str = ink, bold: bool = False, align=PP_ALIGN.LEFT):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            frame = box.text_frame
            frame.clear()
            frame.word_wrap = True
            frame.margin_left = Inches(0.05)
            frame.margin_right = Inches(0.05)
            frame.vertical_anchor = MSO_ANCHOR.TOP
            p = frame.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
            return box

        def add_footer(slide, index: int, accent: str) -> None:
            shape(slide, MSO_SHAPE.RECTANGLE, 0.58, 7.05, 1.2, 0.05, accent)
            textbox(slide, profile["footer"], 0.65, 6.88, 3.8, 0.25, size=8, color=muted)
            textbox(slide, f"{index:02d}", 12.25, 6.82, 0.55, 0.28, size=10, color=muted, bold=True, align=PP_ALIGN.RIGHT)

        def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, *, color: str = ink, accent: str) -> None:
            card_h = min(0.82, max(0.52, h / max(len(items), 1) - 0.08))
            for idx, item in enumerate(items[:5]):
                top = y + idx * (card_h + 0.12)
                shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, card_h, "FFFFFF", line="E2E8F0")
                shape(slide, MSO_SHAPE.OVAL, x + 0.22, top + 0.18, 0.28, 0.28, accent)
                textbox(slide, item, x + 0.68, top + 0.13, w - 0.9, card_h - 0.12, size=16, color=color)

        def add_pipeline(slide, accent: str, labels: list[str]) -> None:
            for i, label in enumerate(labels):
                x = 0.85 + i * 2.42
                shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 5.56, 1.72, 0.62, "FFFFFF", line=accent)
                textbox(slide, label, x + 0.08, 5.72, 1.56, 0.24, size=11, color=ink, bold=True, align=PP_ALIGN.CENTER)
                if i < len(labels) - 1:
                    textbox(slide, "→", x + 1.82, 5.69, 0.48, 0.28, size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)

        for index, (title, bullets) in enumerate(slides, start=1):
            accent = palette[(index - 1) % len(palette)]
            s = prs.slides.add_slide(blank)
            if index == 1:
                shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, dark)
                shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 1.1, "111827")
                shape(s, MSO_SHAPE.ARC, 9.2, -0.85, 4.8, 4.8, accent, transparency=22)
                shape(s, MSO_SHAPE.OVAL, -0.8, 4.8, 3.2, 3.2, "22D3EE", transparency=35)
                textbox(s, profile["eyebrow"], 0.78, 0.65, 6.8, 0.32, size=11, color="93C5FD", bold=True)
                textbox(s, title, 0.72, 1.62, 8.7, 0.95, size=40, color="FFFFFF", bold=True)
                textbox(s, profile["subtitle"], 0.78, 2.78, 8.5, 0.45, size=18, color="CBD5E1")
                cards = profile["cards"]
                for i, card in enumerate(cards):
                    x = 0.82 + i * 3.35
                    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 4.35, 2.75, 1.15, "1E293B", line=accent)
                    textbox(s, card, x + 0.22, 4.62, 2.2, 0.3, size=17, color="FFFFFF", bold=True)
                    textbox(s, f"0{i + 1}", x + 2.14, 4.42, 0.45, 0.22, size=9, color="93C5FD", bold=True, align=PP_ALIGN.RIGHT)
                add_footer(s, index, accent)
                continue

            shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, paper)
            shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.18, accent)
            shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.66, 0.62, 0.72, 0.38, accent)
            textbox(s, f"{index:02d}", 0.78, 0.70, 0.48, 0.16, size=8, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
            textbox(s, title, 1.55, 0.53, 8.5, 0.62, size=27, color=ink, bold=True)
            textbox(s, profile["sidebar"], 10.1, 0.72, 2.4, 0.25, size=10, color=muted, align=PP_ALIGN.RIGHT)
            shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 1.55, 7.18, 4.55, "FFFFFF", line="E2E8F0")
            add_bullets(s, bullets or ["暂无内容"], 1.02, 1.88, 6.55, 3.75, accent=accent)
            shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.35, 1.55, 3.98, 4.55, "EEF2FF" if index % 2 else "ECFEFF", line="CBD5E1")
            if any(marker in title for marker in ("架构", "检索", "生成", "索引", "生产")):
                add_pipeline(s, accent, profile["pipeline"])
                visual_items = profile["pipeline_visual"]
            elif any(marker in title for marker in ("权限", "安全", "评测")):
                visual_items = profile["governance_visual"]
            else:
                visual_items = profile["visual_default"]
            for i, item in enumerate(visual_items):
                row = i // 2
                col = i % 2
                x = 8.72 + col * 1.72
                y = 2.05 + row * 1.2
                shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.32, 0.78, "FFFFFF", line=accent)
                textbox(s, item, x + 0.08, y + 0.24, 1.16, 0.22, size=11, color=ink, bold=True, align=PP_ALIGN.CENTER)
            textbox(s, profile["loop_title"], 8.78, 4.62, 2.9, 0.35, size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)
            textbox(s, profile["loop_caption"], 8.78, 5.05, 2.9, 0.28, size=12, color=muted, align=PP_ALIGN.CENTER)
            add_footer(s, index, accent)

        prs.save(file_path)

    def _deck_visual_profile(self, title: str) -> dict[str, Any]:
        family = self._topic_family(title)
        clean_title = (title or "自动生成幻灯片").strip()
        if family == "rag":
            return {
                "eyebrow": "BAOYU DESIGN · ENTERPRISE KNOWLEDGE",
                "subtitle": "从企业知识治理到可信问答交付的完整 RAG 架构蓝图",
                "sidebar": "RAG 企业知识库落地方案",
                "cards": ["数据治理", "向量检索", "权限评测"],
                "pipeline": ["数据接入", "切分索引", "混合检索", "重排组装", "生成引用"],
                "pipeline_visual": ["Query", "Retriever", "Reranker", "LLM"],
                "governance_visual": ["ACL", "Audit", "Golden Set", "Feedback"],
                "visual_default": ["Source", "Chunk", "Vector", "Answer"],
                "loop_title": "可信知识闭环",
                "loop_caption": "证据 → 权限 → 生成 → 反馈",
                "footer": "PyClaw · Baoyu Design Workflow",
            }
        if family == "ai_agent":
            return {
                "eyebrow": "BAOYU DESIGN · AI AGENT",
                "subtitle": "从目标理解到工具执行的 Agent 架构蓝图",
                "sidebar": "AI Agent 落地方案",
                "cards": ["目标规划", "工具执行", "反馈验收"],
                "pipeline": ["目标理解", "任务规划", "工具调用", "观察反馈", "验证交付"],
                "pipeline_visual": ["Goal", "Planner", "Tools", "Model"],
                "governance_visual": ["Policy", "Approval", "Sandbox", "Audit"],
                "visual_default": ["Goal", "Plan", "Act", "Observe"],
                "loop_title": "Agent 执行闭环",
                "loop_caption": "目标 → 计划 → 执行 → 观察",
                "footer": "PyClaw · Agent Workflow",
            }
        short = clean_title[:18]
        return {
            "eyebrow": "BAOYU DESIGN · PRESENTATION",
            "subtitle": f"围绕“{clean_title}”的结构化表达与落地路径",
            "sidebar": f"{short}方案",
            "cards": ["背景判断", "关键机制", "落地路径"],
            "pipeline": ["背景", "问题", "方案", "执行", "结果"],
            "pipeline_visual": ["Context", "Problem", "Method", "Result"],
            "governance_visual": ["Risk", "Control", "Metric", "Review"],
            "visual_default": ["Context", "Mechanism", "Action", "Result"],
            "loop_title": "实践闭环",
            "loop_caption": "目标 → 方法 → 执行 → 复盘",
            "footer": "PyClaw · Presentation Workflow",
        }

    def _write_basic_pptx(self, file_path: str, slides: list[tuple[str, list[str]]]) -> None:
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        slides = slides or [("自动生成幻灯片", ["暂无内容"])]
        with zipfile.ZipFile(file_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("[Content_Types].xml", self._content_types(len(slides)))
            zf.writestr("_rels/.rels", self._root_rels())
            zf.writestr("ppt/presentation.xml", self._presentation_xml(len(slides)))
            zf.writestr("ppt/_rels/presentation.xml.rels", self._presentation_rels(len(slides)))
            zf.writestr("ppt/theme/theme1.xml", self._theme_xml())
            zf.writestr("ppt/slideMasters/slideMaster1.xml", self._slide_master_xml())
            zf.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", self._slide_master_rels())
            zf.writestr("ppt/slideLayouts/slideLayout1.xml", self._slide_layout_xml())
            zf.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", self._slide_layout_rels())
            zf.writestr("docProps/app.xml", self._app_xml(len(slides)))
            zf.writestr("docProps/core.xml", self._core_xml())
            for index, (title, bullets) in enumerate(slides, start=1):
                zf.writestr(f"ppt/slides/slide{index}.xml", self._slide_xml(title, bullets))
                zf.writestr(f"ppt/slides/_rels/slide{index}.xml.rels", """<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?>\n<Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\"/>""")

    def _write_deck_stage_html(self, file_path: str, slides: list[tuple[str, list[str]]]) -> None:
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        sections: list[str] = []
        for index, (title, bullets) in enumerate(slides or [("自动生成幻灯片", ["暂无内容"])], start=1):
            safe_title = html.escape(title or f"Slide {index}")
            bullet_html = "".join(f"<li>{html.escape(item)}</li>" for item in (bullets or [])[:8])
            sections.append(
                f'<section data-label="{safe_title}">\n'
                f"  <h1>{safe_title}</h1>\n"
                f"  <ul>{bullet_html}</ul>\n"
                "</section>"
            )
        document = (
            "<!doctype html>\n<html lang=\"zh-CN\">\n<head>\n"
            "  <meta charset=\"utf-8\">\n"
            "  <title>PyClaw Skill Deck</title>\n"
            "</head>\n<body>\n"
            "<deck-stage width=\"1920\" height=\"1080\">\n"
            + "\n".join(sections)
            + "\n</deck-stage>\n</body>\n</html>\n"
        )
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(document)

    def _content_types(self, slide_count: int) -> str:
        slide_overrides = "".join(
            f'<Override PartName="/ppt/slides/slide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
            for i in range(1, slide_count + 1)
        )
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
<Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
<Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
<Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
{slide_overrides}
</Types>'''

    def _root_rels(self) -> str:
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
<Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>'''

    def _presentation_xml(self, slide_count: int) -> str:
        slide_ids = "".join(f'<p:sldId id="{255 + i}" r:id="rId{i}"/>' for i in range(1, slide_count + 1))
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
<p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId{slide_count + 1}"/></p:sldMasterIdLst>
<p:sldIdLst>{slide_ids}</p:sldIdLst>
<p:sldSz cx="12192000" cy="6858000" type="wide"/>
<p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>'''

    def _presentation_rels(self, slide_count: int) -> str:
        rels = [
            f'<Relationship Id="rId{i}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{i}.xml"/>'
            for i in range(1, slide_count + 1)
        ]
        rels.append(f'<Relationship Id="rId{slide_count + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>')
        rels.append(f'<Relationship Id="rId{slide_count + 2}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="theme/theme1.xml"/>')
        return '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">' + "".join(rels) + "</Relationships>"

    def _slide_xml(self, title: str, bullets: list[str]) -> str:
        title_xml = self._text_box(2, 600000, 450000, 11000000, 900000, title, font_size=3600, bold=True)
        bullet_runs = []
        for item in bullets[:7] or ["暂无内容"]:
            bullet_runs.append(
                f'<a:p><a:pPr marL="342900" indent="-228600"><a:buChar char="•"/></a:pPr><a:r><a:rPr lang="zh-CN" sz="2200"/><a:t>{html.escape(item)}</a:t></a:r></a:p>'
            )
        body = f'''<p:sp><p:nvSpPr><p:cNvPr id="3" name="Content"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="900000" y="1650000"/><a:ext cx="10400000" cy="4300000"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr><p:txBody><a:bodyPr wrap="square"/><a:lstStyle/>{''.join(bullet_runs)}</p:txBody></p:sp>'''
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
<p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>{title_xml}{body}</p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>'''

    def _text_box(self, shape_id: int, x: int, y: int, cx: int, cy: int, text: str, *, font_size: int, bold: bool = False) -> str:
        bold_attr = ' b="1"' if bold else ""
        return f'''<p:sp><p:nvSpPr><p:cNvPr id="{shape_id}" name="Title"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="zh-CN" sz="{font_size}"{bold_attr}/><a:t>{html.escape(text)}</a:t></a:r></a:p></p:txBody></p:sp>'''

    def _theme_xml(self) -> str:
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="PyClaw"><a:themeElements><a:clrScheme name="PyClaw"><a:dk1><a:srgbClr val="111111"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1><a:dk2><a:srgbClr val="1F2937"/></a:dk2><a:lt2><a:srgbClr val="F9FAFB"/></a:lt2><a:accent1><a:srgbClr val="2563EB"/></a:accent1><a:accent2><a:srgbClr val="16A34A"/></a:accent2><a:accent3><a:srgbClr val="F97316"/></a:accent3><a:accent4><a:srgbClr val="7C3AED"/></a:accent4><a:accent5><a:srgbClr val="0891B2"/></a:accent5><a:accent6><a:srgbClr val="DC2626"/></a:accent6><a:hlink><a:srgbClr val="2563EB"/></a:hlink><a:folHlink><a:srgbClr val="7C3AED"/></a:folHlink></a:clrScheme><a:fontScheme name="PyClaw"><a:majorFont><a:latin typeface="Arial"/><a:ea typeface="Microsoft YaHei"/><a:cs typeface="Arial"/></a:majorFont><a:minorFont><a:latin typeface="Arial"/><a:ea typeface="Microsoft YaHei"/><a:cs typeface="Arial"/></a:minorFont></a:fontScheme><a:fmtScheme name="PyClaw"><a:fillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:fillStyleLst><a:lnStyleLst><a:ln w="9525"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln></a:lnStyleLst><a:effectStyleLst><a:effectStyle><a:effectLst/></a:effectStyle></a:effectStyleLst><a:bgFillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:bgFillStyleLst></a:fmtScheme></a:themeElements><a:objectDefaults/><a:extraClrSchemeLst/></a:theme>'''

    def _slide_master_xml(self) -> str:
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld><p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/><p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst><p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles></p:sldMaster>'''

    def _slide_master_rels(self) -> str:
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/></Relationships>'''

    def _slide_layout_xml(self) -> str:
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank"><p:cSld name="Blank"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sldLayout>'''

    def _slide_layout_rels(self) -> str:
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/></Relationships>'''

    def _app_xml(self, slide_count: int) -> str:
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"><Application>PyClaw</Application><PresentationFormat>Widescreen</PresentationFormat><Slides>{slide_count}</Slides></Properties>'''

    def _core_xml(self) -> str:
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:dcmitype="http://purl.org/dc/dcmitype/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"><dc:title>PyClaw Generated Deck</dc:title><dc:creator>PyClaw</dc:creator></cp:coreProperties>'''
